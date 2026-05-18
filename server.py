"""
ATNE — Adaptador de Textos a Necessitats Educatives
Servidor FastAPI · Jesuïtes Educació

Executa:  python server.py
Obre:     http://localhost:8000
"""

import asyncio
import concurrent.futures
from collections import defaultdict, deque
import hashlib
import hmac
import json
import os
import random
import re
import secrets
import sys
import time
from pathlib import Path
from urllib.parse import quote as _urlquote

# ── BUG FIX Sprint B (2026-04-16) ───────────────────────────────────────────
# Quan arrenquem amb `python server.py`, aquest mòdul s'executa com a
# `__main__`, no com a `server`. Quan `generador_lliure/orquestrador.py` fa
# `from server import _model_for`, Python no troba `server` a sys.modules i
# re-carrega el fitxer des de zero, creant una **segona instància del mòdul**
# amb `_MODEL_CONFIG` reiniciat als defaults (sense cridar _load_system_config).
# Conseqüència: el mode rotate/fixed de /admin no s'aplicava a les crides que
# passessin per imports lazy del mòdul `server` (ex: generar_stream).
#
# Fix: alias-em el mòdul `__main__` com a `server` a sys.modules perquè els
# imports posteriors retornin la MATEIXA instància (amb el _MODEL_CONFIG ja
# carregat per l'event startup). Una sola línia, zero refactor.
if __name__ == "__main__" or __name__ == "__mp_main__":
    sys.modules["server"] = sys.modules[__name__]
# ────────────────────────────────────────────────────────────────────────────

import corpus_reader
import instruction_catalog
import instruction_filter
import requests
import uvicorn
from dotenv import load_dotenv
from fastapi import FastAPI, Body, Depends, HTTPException, Request, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, StreamingResponse, JSONResponse, FileResponse, RedirectResponse

# ── Configuració ────────────────────────────────────────────────────────────

load_dotenv()

# Versió del prompt — git hash curt + data, identifica la versió del codi
# que ha generat cada adaptació del pilot. Imprescindible per a comparatives
# cegues entre setmanes (cada refactor del prompt invalida comparacions).
# Pot sobreescriure's amb env var ATNE_PROMPT_VERSION (útil a Cloud Run on
# git no està disponible — el Cloud Build el resol i l'injecta).
def _resolve_prompt_version() -> str:
    env_v = os.getenv("ATNE_PROMPT_VERSION", "").strip()
    if env_v:
        return env_v
    try:
        import subprocess
        h = subprocess.check_output(
            ["git", "rev-parse", "--short", "HEAD"],
            cwd=Path(__file__).parent,
            stderr=subprocess.DEVNULL,
            timeout=2,
        ).decode().strip()
        date = time.strftime("%Y-%m-%d")
        return f"{date}_{h}" if h else f"{date}_unknown"
    except Exception:
        return f"{time.strftime('%Y-%m-%d')}_unknown"

ATNE_PROMPT_VERSION = _resolve_prompt_version()
print(f"[ATNE] prompt_version: {ATNE_PROMPT_VERSION}", flush=True)

# Claus API, àlies de models i resolució estan a `adaptation/llm_clients.py`.
# Aquí les re-exposem al namespace de `server` per mantenir el contracte amb
# callers externs (snapshot_contract, generador_lliure, tests).
from adaptation.llm_clients import (
    ATNE_MODEL,
    GEMINI_API_KEY,
    GEMINI_API_KEYS,
    GEMMA4_API_KEY,
    GEMMA4_API_KEYS,
    MISTRAL_API_KEY,
    OPENROUTER_API_KEY,
    OPENROUTER_API_KEYS,
    _MODEL_ALIASES,
    _resolve_model,
)

# ── Sprint 1B — selector de model per fase ─────────────────────────────────
#
# _MODEL_CONFIG mapeja cada fase del pipeline (generate/adapt/refine/
# complements/auditor) al model a usar. Es carrega al startup des de la
# taula system_config de Supabase i es pot sobreescriure en calent des de
# /api/admin/config (PUT). Si la DB no respon al startup, cau a ATNE_MODEL
# com a fallback segur. El dict és read-only per als workers: només el
# startup i el PUT d'admin l'escriuen, i el PUT només muta claus atòmiques
# → segur per accés concurrent des dels SSE workers sense lock.
_MODEL_CONFIG: dict[str, str] = {
    "generate": ATNE_MODEL,
    "adapt": ATNE_MODEL,
    "adapt_flash": ATNE_MODEL,   # mode Flash (prompt MVP)
    "refine": ATNE_MODEL,
    "complements": ATNE_MODEL,
    "auditor": "gpt-4o-mini",
    # Fase il·lustracions: tradueix concept catala → query EN + brief FLUX.
    # Default Gemma 4 31B (free tier) — Gemma 3 discontinuat per Google el 2026-05.
    # Configurable des d'/admin com la resta de fases.
    "illustration_translate": "gemma-4-31b-it",
}


def _model_for(phase: str, override: str = "") -> str:
    """Retorna el model_id a usar per a una fase donada.

    Prioritat:
    1. override explícit (ex: payload.model del frontend)
    2. _MODEL_CONFIG[phase], que pot ser:
       - str: mode fix (comportament històric, retrocompatible)
       - dict amb {"mode": "fixed", "model": "..."}: mode fix explícit
       - dict amb {"mode": "rotate", "models": [...], "strategy": "random"}:
         rotació silenciosa, cada crida tria aleatòriament un dels models
    3. ATNE_MODEL (fallback final)

    La rotació és per a validació cega en el pilot (2026-04-20..05-08):
    el docent no veu quin model ha disparat, però el model_id retornat
    es pot registrar a logs per a anàlisi estadística posterior.
    """
    if override:
        return override.strip()
    config = _MODEL_CONFIG.get(phase, ATNE_MODEL)
    if isinstance(config, str):
        return config
    if isinstance(config, dict):
        mode = config.get("mode", "fixed")
        if mode == "fixed":
            return config.get("model") or ATNE_MODEL
        if mode == "rotate":
            models = config.get("models") or []
            if not models:
                return ATNE_MODEL
            if len(models) == 1:
                return models[0]
            # Estratègia random: cada crida tria un model independent.
            # No fem round-robin per evitar persistència d'estat; random
            # sobre un nombre suficient de crides és estadísticament equivalent.
            return random.choice(models)
    return ATNE_MODEL


# Estimació de cost per crida (mitjana per a text educatiu ~1k-2k paraules
# input + ~600-1200 output). Serveix per alimentar el budget tracker del
# pilot i el dashboard /admin. No és facturació real — és una aproximació
# suficient per detectar drifts d'ús i avisar el docent abans de saturar
# el budget_eur_max. Fonts: preus públics proveïdors 2026-04, memòria
# project_llicons_costos_api.md i project_estrategia_escalat.md.
_MODEL_COST_EUR_PER_CALL: dict[str, float] = {
    "gemma-4-31b-it":       0.0,      # Free tier Gemma 4 dense (claus Google)
    "gemma-4-26b-a4b-it":   0.0,      # Free tier Gemma 4 MoE (claus Google) — ~3x més ràpid
    "gemini-2.5-flash":     0.0,      # Free tier Gemini (claus Google)
    "gemini-2.5-flash-lite": 0.0,     # Free tier Gemini (~0.7s, més ràpid)
    "gpt-4o-mini":          0.0036,   # ~2k in + 1k out
    "gpt-4o":               0.045,    # idem ~12× més car
    "gpt-4.1-mini":         0.006,    # high-tier OpenAI preu contingut
    "mistral-small-latest": 0.012,
    "mistral-large-latest": 0.048,
    "qwen/qwen3.5-27b":              0.0003,   # OpenRouter pay-per-use ultra-baix
    "qwen/qwen3.5-9b":               0.0001,   # idem, més petit i barat
    "qwen/qwen3-235b-a22b:free":     0.0,      # Qwen 3 235B MoE — gratuït OpenRouter
    "qwen/qwen3-30b-a3b:free":       0.0,      # Qwen 3 30B MoE — gratuït OpenRouter
    "deepseek/deepseek-chat-v3-0324:free": 0.0,# DeepSeek V3 — gratuït OpenRouter
}


def _estimate_cost_eur(models_per_phase: dict) -> float:
    """Suma estimada de cost per a una adaptació multi-fase.

    `models_per_phase` admet valors string (una crida) o llista (múltiples
    crides en la mateixa fase, ex: retries d'adapt o refines successius).
    Models desconeguts s'aproximen a 0.002€/crida (conservador).
    """
    total = 0.0
    for _phase, entry in models_per_phase.items():
        if entry is None:
            continue
        items = entry if isinstance(entry, list) else [entry]
        for m in items:
            if not isinstance(m, str) or not m:
                continue
            _prov, spec = _resolve_model(m)
            total += _MODEL_COST_EUR_PER_CALL.get(spec, 0.002)
    return round(total, 6)


# Flags runtime carregades també de system_config al startup. Els defaults
# són els valors inicials del SQL 1B; si la DB falla, aquests valen.
_AUDITOR_ENABLED_RUNTIME: bool = False
_ADMIN_BUDGET_EUR_MAX: float = 30.0
_PILOT_ACTIVE: bool = True


SUPABASE_URL = os.getenv("SUPABASE_URL", "")
SUPABASE_ANON_KEY = os.getenv("SUPABASE_ANON_KEY", "")
# Service-role (opcional). Si està present, l'usem per inserts crítics
# (telemetria pilot) perquè salta RLS. Mai exposar al frontend.
SUPABASE_SERVICE_KEY = (
    os.getenv("SUPABASE_SERVICE_KEY", "")
    or os.getenv("SUPABASE_SERVICE_ROLE_KEY", "")
)

SUPABASE_HEADERS = {
    "apikey": SUPABASE_ANON_KEY,
    "Authorization": f"Bearer {SUPABASE_ANON_KEY}",
    "Content-Type": "application/json",
}

# Headers per escriptures de telemetria del pilot. Usen service-role si està
# configurada (salta RLS); en cas contrari cauen a l'anon (que depèn de la
# política RLS de la taula). Mai exposar SUPABASE_SERVICE_KEY al client.
_TELEMETRY_WRITE_KEY = SUPABASE_SERVICE_KEY or SUPABASE_ANON_KEY
SUPABASE_WRITE_HEADERS = {
    "apikey": _TELEMETRY_WRITE_KEY,
    "Authorization": f"Bearer {_TELEMETRY_WRITE_KEY}",
    "Content-Type": "application/json",
}
_USING_SERVICE_KEY = bool(SUPABASE_SERVICE_KEY)

# Fitxer de fallback per a events de telemetria que no han pogut entrar a
# Supabase (xarxa caiguda, RLS, 5xx…). NDJSON append-only, pensat per ser
# rejugat manualment via /api/admin/pilot/replay-fallback. A Cloud Run el
# directori /tmp és tmpfs (efímer entre restarts) però els logs stdout sí
# són durables — per això SEMPRE imprimim el payload a stdout també.
_PILOT_FALLBACK_DIR = Path(os.getenv("ATNE_PILOT_FALLBACK_DIR", "/tmp")) \
    if os.name != "nt" else Path(os.getenv("ATNE_PILOT_FALLBACK_DIR", os.getenv("TEMP", "."))).resolve()
_PILOT_FALLBACK_FILE = _PILOT_FALLBACK_DIR / "atne_pilot_fallback.ndjson"


def _append_telemetry_fallback(kind: str, payload: dict, error: str) -> None:
    """Escriu l'event fallit a stdout (sempre, durable a Cloud Logging) i a
    un fitxer NDJSON local (recuperable si el procés segueix viu). Mai
    llança excepció — això mateix és la xarxa de seguretat de la xarxa de
    seguretat.
    """
    import datetime as _dt
    record = {
        "ts": _dt.datetime.utcnow().isoformat() + "Z",
        "kind": kind,
        "error": error,
        "payload": payload,
    }
    # 1. stdout amb tag indexable per a Cloud Logging
    try:
        print(f"[TELEMETRY_FAILED] {json.dumps(record, ensure_ascii=False, default=str)}", flush=True)
    except Exception:
        pass
    # 2. fitxer local NDJSON (millor effort)
    try:
        _PILOT_FALLBACK_DIR.mkdir(parents=True, exist_ok=True)
        with _PILOT_FALLBACK_FILE.open("a", encoding="utf-8") as f:
            f.write(json.dumps(record, ensure_ascii=False, default=str) + "\n")
    except Exception:
        pass


def _load_system_config() -> dict:
    """Llegeix la taula system_config de Supabase i actualitza l'estat runtime.

    Crida al startup FastAPI i a cada PUT /api/admin/config. Tolerant a
    errors: si la DB no respon, deixa els valors inicials de _MODEL_CONFIG
    i dels flags globals (fallback segur per a l'arrencada offline).
    Retorna un dict amb l'estat actual per a debugging/GET d'admin.
    """
    global _AUDITOR_ENABLED_RUNTIME, _ADMIN_BUDGET_EUR_MAX, _PILOT_ACTIVE
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        print("[ATNE] Sense SUPABASE_URL/KEY — _MODEL_CONFIG manté defaults")
        return {"model_config": dict(_MODEL_CONFIG), "source": "defaults"}
    try:
        r = requests.get(
            f"{SUPABASE_URL}/rest/v1/system_config?select=key,value",
            headers=SUPABASE_HEADERS,
            timeout=5,
        )
        r.raise_for_status()
        rows = r.json()
    except Exception as e:
        print(f"[ATNE] ERROR carregant system_config: {e} — manté defaults")
        return {"model_config": dict(_MODEL_CONFIG), "source": "error", "error": str(e)}

    # Mapa: clau de DB → clau de _MODEL_CONFIG
    model_key_map = {
        "atne_model_generate":               "generate",
        "atne_model_adapt":                  "adapt",
        "atne_model_adapt_flash":            "adapt_flash",
        "atne_model_refine":                 "refine",
        "atne_model_complements":            "complements",
        "atne_model_auditor":                "auditor",
        "atne_model_illustration_translate": "illustration_translate",
    }
    for row in rows:
        key = row.get("key")
        val = row.get("value")
        if key in model_key_map and isinstance(val, dict):
            # Format nou amb rotació silenciosa:
            # {"mode": "fixed", "model": "..."} o
            # {"mode": "rotate", "models": [...], "strategy": "random"}
            if "mode" in val:
                _MODEL_CONFIG[model_key_map[key]] = val
            else:
                # Format legacy: {"model_id": "...", ...}
                model_id = val.get("model_id")
                if model_id:
                    _MODEL_CONFIG[model_key_map[key]] = model_id
        elif key == "atne_auditor_enabled":
            _AUDITOR_ENABLED_RUNTIME = bool(val)
        elif key == "admin_budget_eur_max":
            try:
                _ADMIN_BUDGET_EUR_MAX = float(val) if val is not None else 30.0
            except (TypeError, ValueError):
                _ADMIN_BUDGET_EUR_MAX = 30.0
        elif key == "pilot_active":
            _PILOT_ACTIVE = bool(val)

    print(f"[ATNE] _MODEL_CONFIG carregat de system_config: {_MODEL_CONFIG}")
    print(f"[ATNE] auditor_enabled={_AUDITOR_ENABLED_RUNTIME} budget={_ADMIN_BUDGET_EUR_MAX}€ pilot_active={_PILOT_ACTIVE}")
    return {
        "model_config": dict(_MODEL_CONFIG),
        "auditor_enabled": _AUDITOR_ENABLED_RUNTIME,
        "admin_budget_eur_max": _ADMIN_BUDGET_EUR_MAX,
        "pilot_active": _PILOT_ACTIVE,
        "source": "supabase",
    }


PROFILES_DIR = Path(__file__).parent / "profiles"
PROFILES_DIR.mkdir(exist_ok=True)

UI_DIR = Path(__file__).parent / "ui"


# ── Sprint 1B — auth admin (simple cookie signada) ─────────────────────────
#
# El /admin està protegit per password. Un sol rol "admin" (el Miquel).
# Flux: POST /api/admin/login amb {password} → si OK, seta cookie signada
# amb HMAC(secret, "admin:<ts>"). _require_admin() verifica signatura i
# expiració (8h). Per a producció (Cloud Run) cal fixar ADMIN_PASSWORD i
# ADMIN_SESSION_SECRET via env vars; si no, les cookies moriran a cada
# restart del pod (que és OK per al MVP del pilot).
ADMIN_PASSWORD = os.getenv("ADMIN_PASSWORD", "")
ADMIN_SESSION_SECRET = os.getenv("ADMIN_SESSION_SECRET") or secrets.token_hex(32)
ADMIN_SESSION_TTL_SEC = 8 * 3600  # 8h
# Logins lanet separats per comes que sempre tindran rol admin (sense necessitat de BD).
# Ex: ATNE_ADMIN_LOGINS=claudeai.je,miquel.amor
_ADMIN_LOGINS_RAW = os.getenv("ATNE_ADMIN_LOGINS", "")
_ADMIN_DOCENT_IDS: set[str] = set()


def _admin_sign(payload: str) -> str:
    sig = hmac.new(
        ADMIN_SESSION_SECRET.encode(),
        payload.encode(),
        hashlib.sha256,
    ).hexdigest()
    return f"{payload}.{sig}"


def _admin_verify(token: str) -> bool:
    """True si el token és signat correctament i no ha expirat."""
    if not token or "." not in token:
        return False
    try:
        payload, sig = token.rsplit(".", 1)
        expected = hmac.new(
            ADMIN_SESSION_SECRET.encode(),
            payload.encode(),
            hashlib.sha256,
        ).hexdigest()
        if not hmac.compare_digest(sig, expected):
            return False
        parts = payload.split(":")
        if len(parts) != 2 or parts[0] != "admin":
            return False
        ts = int(parts[1])
        if time.time() - ts > ADMIN_SESSION_TTL_SEC:
            return False
        return True
    except Exception:
        return False


def _verify_session(token: str) -> str | None:
    """Retorna el login si la cookie de sessió és vàlida, o None."""
    if not token or "." not in token:
        return None
    try:
        payload, sig = token.rsplit(".", 1)
        expected = hmac.new(
            ADMIN_SESSION_SECRET.encode(),
            payload.encode(),
            hashlib.sha256,
        ).hexdigest()
        if not hmac.compare_digest(sig, expected):
            return None
        parts = payload.split(":", 2)  # session:login:ts
        if len(parts) != 3 or parts[0] != "session":
            return None
        ts = int(parts[2])
        if time.time() - ts > 8 * 3600:
            return None
        return parts[1]  # login
    except Exception:
        return None


def _safe_error(e: Exception, prefix: str = "Error intern") -> str:
    """Missatge genèric per al client; detall complet al log del servidor."""
    print(f"[ATNE:error] {type(e).__name__}: {e}", flush=True)
    return prefix


def _is_admin_login(login: str) -> bool:
    """True si el login és a la llista ATNE_ADMIN_LOGINS."""
    if not login or not _ADMIN_LOGINS_RAW:
        return False
    return login in {l.strip() for l in _ADMIN_LOGINS_RAW.split(",") if l.strip()}


def _require_admin(request: Request) -> bool:
    """Dependency FastAPI. Accepta: (1) cookie admin HMAC, (2) sessió de login admin."""
    if _admin_verify(request.cookies.get("atne_admin", "")):
        return True
    session_login = _verify_session(request.cookies.get("atne_session", ""))
    if session_login and _is_admin_login(session_login):
        return True
    raise HTTPException(status_code=401, detail="Admin auth required")


# Rate limiter en memòria — clau: string, valor: deque de timestamps
_rate_limits: dict[str, deque] = defaultdict(deque)


def _rate_check(key: str, limit: int, window_sec: int) -> None:
    """Llança 429 si key ha superat limit peticions en window_sec segons."""
    now = time.monotonic()
    q = _rate_limits[key]
    while q and now - q[0] > window_sec:
        q.popleft()
    if len(q) >= limit:
        raise HTTPException(status_code=429, detail="Massa peticions, espera un moment")
    q.append(now)

# Client Gemini compartit per al health-check — definit a adaptation/llm_clients.
from adaptation.llm_clients import gemini_client

print(f"[ATNE] Model actiu: {ATNE_MODEL}")

# ── FastAPI app ─────────────────────────────────────────────────────────────

app = FastAPI(title="ATNE", version="0.1.0")

# ── Auth (lanet — tokenNet via bridge PHP) ─────────────────────────────────
#
# Middleware ASGI pur (no BaseHTTPMiddleware) per compatibilitat amb les
# respostes streaming (SSE del pipeline LLM). Valida el token lanet contra
# el bridge PHP de FJE (LANET_BRIDGE_URL). Els endpoints públics i els
# /api/admin/* (auth pròpia via HMAC cookie) queden exempts.
#
# Flux:
#   Frontend → redirigeix a lanet_bridge.php (FJE) → torna amb token a URL
#   → JS guarda token a localStorage → envia com Authorization: Bearer
#   → aquest middleware valida el token via POST al bridge PHP

LANET_BRIDGE_URL = os.getenv("LANET_BRIDGE_URL", "")

# Cache simple: Key = token; value = (expires_at_monotonic, login).
# TTL 2 min per limitar crides al bridge sense deixar tokens revocats actius.
_LANET_AUTH_CACHE: dict[str, tuple[float, str]] = {}
_LANET_AUTH_CACHE_TTL = 120.0

ATNE_PUBLIC_API_PATHS = {
    "/api/health",
    "/api/runtime-config",
    "/api/auth/exchange",   # intercanvi token LaNet → cookie httpOnly (sense auth prèvia)
}

def _atne_is_public_path(path: str) -> bool:
    if not path.startswith("/api/"):
        return True
    if path in ATNE_PUBLIC_API_PATHS:
        return True
    if path.startswith("/api/admin/") or path.startswith("/api/audit/"):
        return True
    return False


def _lanet_validate_token(token: str) -> tuple[int, str]:
    """Valida el token lanet via POST al bridge PHP.

    Retorna (status_code, login). status_code:
      200 → token vàlid, login retornat
      401 → token rebutjat pel bridge
      503 → bridge no configurat o no accessible
    """
    if not LANET_BRIDGE_URL:
        return 503, ""
    import time as _time
    now = _time.monotonic()
    cached = _LANET_AUTH_CACHE.get(token)
    if cached and cached[0] > now:
        return 200, cached[1]
    try:
        r = requests.post(
            LANET_BRIDGE_URL,
            json={"token": token},
            timeout=5,
        )
        if r.status_code != 200:
            return 401, ""
        data = r.json()
        login = (data.get("login") or "").strip()
        if not login:
            return 401, ""
        _LANET_AUTH_CACHE[token] = (now + _LANET_AUTH_CACHE_TTL, login)
        if len(_LANET_AUTH_CACHE) > 500:
            expired = [k for k, (exp, _) in _LANET_AUTH_CACHE.items() if exp <= now]
            for k in expired:
                _LANET_AUTH_CACHE.pop(k, None)
        return 200, login
    except Exception as e:
        print(f"[ATNE:auth] Error contactant bridge lanet: {e}")
        return 503, ""


class _AtneAuthMiddleware:
    def __init__(self, app):
        self.app = app

    async def __call__(self, scope, receive, send):
        if scope["type"] != "http":
            await self.app(scope, receive, send)
            return
        path = scope.get("path", "")
        method = scope.get("method", "GET")
        if method == "OPTIONS" or _atne_is_public_path(path):
            await self.app(scope, receive, send)
            return
        # Bypass per a tests locals automatitzats (ATNE_TEST_BYPASS_AUTH=1).
        if os.getenv("ATNE_TEST_BYPASS_AUTH") == "1":
            client_ip = (scope.get("client") or ("", 0))[0]
            if client_ip in ("127.0.0.1", "::1"):
                scope.setdefault("state", {})
                scope["state"]["user_login"] = "test_docent"
                await self.app(scope, receive, send)
                return
        # 1. Cookie de sessió httpOnly (intercanviada via /api/auth/exchange)
        cookie_login = None
        for k, v in scope.get("headers", []):
            if k == b"cookie":
                for part in v.decode("latin-1", errors="replace").split(";"):
                    part = part.strip()
                    if part.startswith("atne_session="):
                        cookie_login = _verify_session(part[len("atne_session="):].strip())
                break
        if cookie_login:
            scope.setdefault("state", {})
            scope["state"]["user_login"] = cookie_login
            scope["state"]["user_email"] = cookie_login
            print(f"[ATNE:auth] {cookie_login} {method} {path} (cookie)", flush=True)
            await self.app(scope, receive, send)
            return

        # 2. Fallback: Authorization: Bearer (token LaNet directe)
        if not LANET_BRIDGE_URL:
            await self._send_error(send, 503, "Auth no configurada al servidor")
            return
        auth_header = ""
        for k, v in scope.get("headers", []):
            if k == b"authorization":
                auth_header = v.decode("latin-1", errors="replace")
                break
        if not auth_header.lower().startswith("bearer "):
            await self._send_error(send, 401, "No autenticat")
            return
        token = auth_header[7:].strip()
        status, login = _lanet_validate_token(token)
        if status == 503:
            await self._send_error(send, 503, "No es pot validar sessió")
            return
        if status != 200:
            await self._send_error(send, 401, "Token invàlid o expirat")
            return
        scope.setdefault("state", {})
        scope["state"]["user_login"] = login
        # Mantenim user_email com a àlies per compatibilitat amb codi existent
        scope["state"]["user_email"] = login
        print(f"[ATNE:auth] {login} {method} {path}", flush=True)
        await self.app(scope, receive, send)

    async def _send_error(self, send, status: int, msg: str):
        body = json.dumps({"ok": False, "error": msg}).encode("utf-8")
        await send({
            "type": "http.response.start",
            "status": status,
            "headers": [
                (b"content-type", b"application/json; charset=utf-8"),
                (b"content-length", str(len(body)).encode("ascii")),
            ],
        })
        await send({"type": "http.response.body", "body": body})


class _AtneSecurityHeadersMiddleware:
    """Afegeix headers de seguretat bàsics a totes les respostes."""
    def __init__(self, app):
        self.app = app

    async def __call__(self, scope, receive, send):
        if scope["type"] != "http":
            await self.app(scope, receive, send)
            return

        async def _send_with_headers(message):
            if message["type"] == "http.response.start":
                headers = list(message.get("headers") or [])
                def _set(h, v):
                    for i, (k, _) in enumerate(headers):
                        if k.lower() == h:
                            headers[i] = (h, v)
                            return
                    headers.append((h, v))
                _set(b"x-content-type-options", b"nosniff")
                _set(b"x-frame-options", b"DENY")
                _set(b"referrer-policy", b"strict-origin-when-cross-origin")
                _set(b"permissions-policy", b"camera=(), microphone=(), geolocation=()")
                _set(b"strict-transport-security", b"max-age=31536000; includeSubDomains")
                _set(b"content-security-policy", (
                    b"default-src 'self'; "
                    b"script-src 'self' 'unsafe-inline' https://cdn.jsdelivr.net; "
                    b"style-src 'self' 'unsafe-inline' https://fonts.googleapis.com; "
                    b"font-src 'self' data: https://fonts.gstatic.com https://cdn.jsdelivr.net; "
                    b"img-src 'self' data: blob: https:; "
                    b"connect-src 'self' https:; "
                    b"object-src 'none'; "
                    b"frame-src 'none'; "
                    b"frame-ancestors 'none'; "
                    b"base-uri 'self';"
                ))
                message["headers"] = headers
            await send(message)

        await self.app(scope, receive, _send_with_headers)


# Ordre de registre (primer = més extern en la request):
#  1. Security headers (outermost, aplica a totes les respostes inclòs 401)
#  2. CORS (gestiona preflight OPTIONS sense passar per auth)
#  3. Auth (innermost, verifica JWT per /api/* protegits)
app.add_middleware(_AtneSecurityHeadersMiddleware)

ATNE_ALLOWED_ORIGINS = [
    "http://localhost:8000",
    "http://127.0.0.1:8000",
    "http://localhost:8001",
    "http://127.0.0.1:8001",
    "http://localhost:8080",
    "http://127.0.0.1:8080",
    "https://atne-1050342211642.europe-west1.run.app",
]
app.add_middleware(
    CORSMiddleware,
    allow_origins=ATNE_ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "PATCH", "DELETE", "OPTIONS"],
    allow_headers=["Authorization", "Content-Type", "X-Requested-With"],
)

app.add_middleware(_AtneAuthMiddleware)


@app.on_event("startup")
async def _atne_startup():
    """Càrrega de la configuració runtime des de Supabase.

    Sprint 1B: el selector de model per fase viu a system_config. Carreguem
    els valors un cop al boot i els refrescarem via PUT /api/admin/config.
    Si la càrrega falla, els defaults de _MODEL_CONFIG fan de fallback.

    Probe telemetria pilot (afegit 2026-05-15): valida que el backend pot
    escriure a atne_pilot_events. Si RLS bloqueja l'insert (bug que va fer
    perdre suggeriments als homòlegs DOP/InfPri), ho diem A LA CARA al log
    de boot perquè no torni a passar en silenci.
    """
    _load_system_config()

    # ── Asserts de seguretat de la service-role key ──────────────────────
    # La service-role bypassa RLS — un leak al frontend permetria modificar
    # qualsevol taula. Aquí verifiquem el bàsic abans d'usar-la.
    if SUPABASE_SERVICE_KEY:
        if SUPABASE_SERVICE_KEY == SUPABASE_ANON_KEY:
            raise RuntimeError(
                "[ATNE startup] SUPABASE_SERVICE_KEY == SUPABASE_ANON_KEY. "
                "Has copiat la clau equivocada — la service-role es troba a "
                "Project Settings → API → 'service_role secret' (sb_secret_… o JWT)."
            )
        if not (
            SUPABASE_SERVICE_KEY.startswith("sb_secret_")
            or SUPABASE_SERVICE_KEY.startswith("eyJ")
        ):
            print(
                "[ATNE startup] ⚠️  SUPABASE_SERVICE_KEY té format inesperat "
                "(no comença per 'sb_secret_' ni 'eyJ'). Continuem, però "
                "revisa que sigui la clau correcta.",
                flush=True,
            )
        # Mai imprimim la clau sencera — només una empremta per identificar-la.
        _masked = f"{SUPABASE_SERVICE_KEY[:12]}…{SUPABASE_SERVICE_KEY[-4:]}"
        print(
            f"[ATNE startup] SUPABASE_SERVICE_KEY configurada ({_masked}) — "
            "telemetria pilot bypassa RLS.",
            flush=True,
        )
    else:
        print(
            "[ATNE startup] SUPABASE_SERVICE_KEY no configurada — la "
            "telemetria depèn de polítiques RLS d'INSERT per anon (vegeu "
            "docs/sql/fix_rls_pilot_events_20260515.sql).",
            flush=True,
        )

    if SUPABASE_URL and _TELEMETRY_WRITE_KEY:
        import datetime as _dt
        try:
            probe = {
                "event_type": "client_error",
                "data": {"_boot_probe": True, "ts": _dt.datetime.utcnow().isoformat() + "Z"},
            }
            r = requests.post(
                f"{SUPABASE_URL}/rest/v1/atne_pilot_events",
                headers={**SUPABASE_WRITE_HEADERS, "Prefer": "return=minimal"},
                json=probe,
                timeout=5,
            )
            if r.status_code in (200, 201, 204):
                print(
                    f"[ATNE startup] OK telemetria pilot: insert a atne_pilot_events funciona "
                    f"(service_key={_USING_SERVICE_KEY}).",
                    flush=True,
                )
            else:
                print(
                    "\n" + "=" * 72 + "\n"
                    "[ATNE startup] ⚠️  TELEMETRIA PILOT TRENCADA — els events "
                    "(suggeriments, refines, etc.) NO entraran a Supabase.\n"
                    f"  status={r.status_code}  body={r.text[:300]}\n"
                    f"  service_key_configured={_USING_SERVICE_KEY}\n"
                    "  Fix: aplica docs/sql/fix_rls_pilot_events_20260515.sql al SQL editor "
                    "de Supabase, o configura SUPABASE_SERVICE_KEY al Cloud Run.\n"
                    + "=" * 72,
                    flush=True,
                )
        except Exception as e:
            print(f"[ATNE startup] probe telemetria error: {e!r}", flush=True)


# ── Esborranys (drafts) del Pas 2 ──────────────────────────────────────────
# Extret a `routes/drafts.py` (refactor 2026-04-21). 5 endpoints sota /api/drafts.
# El registrem via include_router just despres de configurar el middleware.
from routes.drafts import router as _drafts_router  # noqa: E402
app.include_router(_drafts_router)

# Biblioteca d'adaptacions desades (Pas 3) a `atne_adaptations`.
# Extret a `routes/adaptations.py` (2026-04-22). 4 endpoints sota /api/adaptations.
from routes.adaptations import router as _adaptations_router  # noqa: E402
app.include_router(_adaptations_router)


# ── Pàgines HTML ────────────────────────────────────────────────────────────

@app.get("/favicon.ico", include_in_schema=False)
async def favicon():
    from fastapi.responses import Response
    return Response(status_code=204)


@app.get("/")
async def index():
    return RedirectResponse(url="/ui/atne/home.html")


@app.get("/legacy", response_class=HTMLResponse)
async def index_legacy():
    """UI antic (preservat per a debug / admin). Veure bug 4 (2026-04-19)."""
    return HTMLResponse(
        (UI_DIR / "index.html").read_text(encoding="utf-8"),
        headers={"Cache-Control": "no-store"},
    )


@app.get("/logo-fje.jpg")
async def serve_logo_fje():
    """Serveix el logo FJE des del directori arrel."""
    logo = Path(__file__).parent / "logo fje.jpg"
    if logo.exists():
        return FileResponse(str(logo), media_type="image/jpeg")
    raise HTTPException(404, "Logo no trobat")


@app.get("/ui/{path:path}")
async def serve_static(path: str):
    file = UI_DIR / path
    if not file.exists() or not file.is_file():
        raise HTTPException(404, "Fitxer no trobat")
    content_types = {
        ".html": "text/html",
        ".css": "text/css",
        ".js": "application/javascript",
        ".png": "image/png",
        ".svg": "image/svg+xml",
        ".ico": "image/x-icon",
        ".woff2": "font/woff2",
        ".woff": "font/woff",
        ".ttf": "font/ttf",
    }
    ct = content_types.get(file.suffix, "application/octet-stream")
    return FileResponse(file, media_type=ct)



# ── System prompt per Gemini ───────────────────────────────────────────────

# ── Prompt v2 RAG: instruccions carregades del corpus (no hardcoded) ──────
# Inicialitzar cache del corpus al arrencar
corpus_reader.load_corpus()
_corpus_stats = corpus_reader.get_all_loaded_stats()
print(f"[corpus_reader] Carregat: {len(_corpus_stats['profiles'])} perfils, "
      f"{len(_corpus_stats['mecr'])} MECR, {len(_corpus_stats['dua'])} DUA, "
      f"{len(_corpus_stats['genres'])} gèneres, {_corpus_stats['crossings']} creuaments")


# Les funcions _get_active_profiles / _str_to_bool / build_persona_audience /
# build_system_prompt s han mogut a adaptation/prompt_builder.py. Les re-exposem
# aqui per mantenir el contracte amb generador_lliure, tests i snapshot.
from adaptation.prompt_builder import (
    _get_active_profiles,
    _str_to_bool,
    build_persona_audience,
    build_system_prompt,
)



# ── Post-processament Python (verificació post-LLM) ──────────────────────
#
# Extret a `adaptation/post_process.py` (refactor 2026-04-21).
# Re-importem aqui tots els simbols perque codi extern que fa
# `from server import clean_gemini_output` (o qualsevol altre) continui
# funcionant sense canvis. Els tests externs (tests/.tmp/*.py) tambe fan
# servir `server.MECR_MAX_WORDS` i `server.FORBIDDEN_WORDS`.
from adaptation.post_process import (  # noqa: E402,F401
    MECR_MAX_WORDS,
    FORBIDDEN_WORDS,
    post_process_adaptation,
    clean_gemini_output,
    _strip_latex_artifacts,
    _fix_english_words,
    _fix_typos,
    _fix_word_concatenations,
    _post_process_llm_output,
    # Constants internes re-exportades per compatibilitat (les usa algun test?)
    _LATEX_PATTERNS,
    _ENGLISH_REPLACEMENTS,
    _TYPO_FIXES,
    _CONCAT_WORD_RE,
)



# Els tres wrappers _call_llm / _call_llm_raw / _call_llm_stream s hi mouen
# a adaptation/llm_clients.py. Els re-exposem al namespace de server per
# compatibilitat amb generador_lliure, tests i el snapshot contract.
from adaptation.llm_clients import (
    _call_llm,
    _call_llm_raw,
    _call_llm_stream,
)






def _log_session(session: dict) -> None:
    """Insereix una fila a atne_sessions (Supabase). Fire-and-forget.

    Si l'insert falla (xarxa, RLS, 5xx…) el payload sencer queda a stdout
    (tag [TELEMETRY_FAILED], durable a Cloud Logging) i al fitxer fallback.
    Mai es perd silenciosament.
    """
    if not SUPABASE_URL or not _TELEMETRY_WRITE_KEY:
        _append_telemetry_fallback("atne_sessions", session, "supabase_no_configurat")
        return
    try:
        r = requests.post(
            f"{SUPABASE_URL}/rest/v1/atne_sessions",
            headers={**SUPABASE_WRITE_HEADERS, "Prefer": "return=minimal"},
            json=session,
            timeout=5,
        )
        if r.status_code not in (200, 201, 204):
            _append_telemetry_fallback(
                "atne_sessions",
                session,
                f"http_{r.status_code}: {r.text[:400]}",
            )
    except Exception as e:
        _append_telemetry_fallback("atne_sessions", session, f"exception: {e!r}")


def _log_pilot_event(event: dict) -> None:
    """Insereix una fila a atne_pilot_events (Supabase). Fire-and-forget.

    Si l'insert falla per RLS, 5xx o xarxa, el payload queda durable a
    stdout (tag [TELEMETRY_FAILED]) i al fitxer fallback NDJSON. Es pot
    rejugar amb POST /api/admin/pilot/replay-fallback.
    """
    if not SUPABASE_URL or not _TELEMETRY_WRITE_KEY:
        _append_telemetry_fallback("atne_pilot_events", event, "supabase_no_configurat")
        return
    try:
        r = requests.post(
            f"{SUPABASE_URL}/rest/v1/atne_pilot_events",
            headers={**SUPABASE_WRITE_HEADERS, "Prefer": "return=minimal"},
            json=event,
            timeout=5,
        )
        if r.status_code not in (200, 201, 204):
            _append_telemetry_fallback(
                "atne_pilot_events",
                event,
                f"http_{r.status_code}: {r.text[:400]}",
            )
    except Exception as e:
        _append_telemetry_fallback("atne_pilot_events", event, f"exception: {e!r}")


def _docent_hash_from_id(docent_id: str) -> str:
    """SHA256 truncat del docent_id (email). Estable, no reversible."""
    if not docent_id:
        return ""
    salt = os.getenv("ATNE_DOCENT_SALT", "atne-pilot-2026")
    h = hashlib.sha256(f"{docent_id.strip().lower()}:{salt}".encode()).hexdigest()
    return h[:16]


# VERIFY_SYSTEM, _verify_adaptation, run_adaptation + estat d'audit
# (_ATNE_LAST_ADAPTATION, _ATNE_ADAPTATIONS_LOG, _ATNE_ADAPTATIONS_MAX) s'han
# mogut a adaptation/orchestrator.py. Re-exposem run_adaptation aqui per
# mantenir el contracte. El buffer d'adaptacions es llegeix via `orchestrator._ATNE_*`
# des dels endpoints d'audit (lectura dinamica, no stale).
from adaptation import orchestrator as _orchestrator
from adaptation.orchestrator import run_adaptation


# ── API endpoints ───────────────────────────────────────────────────────────

@app.get("/api/health")
async def health():
    """
    Verifica l'estat complet del sistema:
    - Supabase (vector store)
    - Claus API de tots els models LLM suportats (booleans, sense valor)
    - Connectivitat amb LanguageTool
    - Configuració del pipeline de qualitat (auditor LLM opt-in)
    Retorna 200 si Supabase + model default funcionen, 503 en cas contrari.
    """
    checks = {
        "supabase": False,
        "llm": False,
        "model_default": ATNE_MODEL,
        "keys": {
            "gemma4": {
                "configured": bool(GEMMA4_API_KEYS and any(GEMMA4_API_KEYS)),
                "count": len([k for k in GEMMA4_API_KEYS if k]),
            },
            "mistral": {
                "configured": bool(MISTRAL_API_KEY),
            },
            "openai": {
                "configured": bool(os.getenv("OPENAI_API_KEY")),
            },
            "gemini": {
                "configured": bool(GEMINI_API_KEYS and any(GEMINI_API_KEYS)),
                "count": len([k for k in GEMINI_API_KEYS if k]),
            },
        },
        "languagetool": {
            "url": LANGUAGETOOL_URL,
            "reachable": False,
        },
        "auditor": {
            "enabled": _AUDITOR_ENABLED_RUNTIME,
            "model": ATNE_AUDITOR_MODEL,
            "can_run": bool(os.getenv("OPENAI_API_KEY")),
        },
        "telemetry_write": {
            "using_service_key": _USING_SERVICE_KEY,
            "can_write_pilot_events": None,
            "error": None,
        },
    }

    # Supabase check
    try:
        resp = requests.get(
            f"{SUPABASE_URL}/rest/v1/rag_fje",
            headers=SUPABASE_HEADERS,
            params={"select": "id", "limit": "1"},
            timeout=5,
        )
        checks["supabase"] = resp.status_code == 200
    except Exception:
        pass

    # Model default check
    try:
        if ATNE_MODEL == "mistral":
            r = requests.get(
                "https://api.mistral.ai/v1/models",
                headers={"Authorization": f"Bearer {MISTRAL_API_KEY}"},
                timeout=5,
            )
            checks["llm"] = r.status_code == 200
        elif ATNE_MODEL == "gemma4":
            gemini_client.models.get(model="gemma-4-31b-it")
            checks["llm"] = True
        elif ATNE_MODEL == "gpt":
            # Només comprovar que la clau hi és (evita cost de crida)
            checks["llm"] = bool(os.getenv("OPENAI_API_KEY"))
        else:
            gemini_client.models.get(model="gemini-2.5-flash")
            checks["llm"] = True
    except Exception:
        pass

    # LanguageTool connectivity check (quick probe a /languages endpoint)
    try:
        base = LANGUAGETOOL_URL.rsplit("/check", 1)[0]
        r = requests.get(f"{base}/languages", timeout=4)
        checks["languagetool"]["reachable"] = r.status_code == 200
    except Exception:
        pass

    # Telemetry write probe (CRÍTIC) — sense això es perden dades del pilot.
    # Detecta el bug del 2026-05-15 (RLS bloquejant inserts a atne_pilot_events).
    try:
        import datetime as _dt
        probe_row = {
            "event_type": "client_error",  # whitelist-safe
            "step": None,
            "data": {"_health_probe": True, "ts": _dt.datetime.utcnow().isoformat() + "Z"},
        }
        r = requests.post(
            f"{SUPABASE_URL}/rest/v1/atne_pilot_events",
            headers={**SUPABASE_WRITE_HEADERS, "Prefer": "return=minimal"},
            json=probe_row,
            timeout=5,
        )
        if r.status_code in (200, 201, 204):
            checks["telemetry_write"]["can_write_pilot_events"] = True
        else:
            checks["telemetry_write"]["can_write_pilot_events"] = False
            checks["telemetry_write"]["error"] = f"http_{r.status_code}: {r.text[:200]}"
    except Exception as e:
        checks["telemetry_write"]["can_write_pilot_events"] = False
        checks["telemetry_write"]["error"] = str(e)

    ok = checks["supabase"] and checks["llm"]
    return JSONResponse({"ok": ok, **checks}, status_code=200 if ok else 503)


# ── Admin API (Sprint 1B) ───────────────────────────────────────────────────
#
# Endpoints per gestionar la configuració runtime del pilot des del
# dashboard /admin. Protegits per password via cookie signada.

_ALLOWED_MODEL_KEYS = {
    "atne_model_generate":               "generate",
    "atne_model_adapt":                  "adapt",
    "atne_model_adapt_flash":            "adapt_flash",
    "atne_model_refine":                 "refine",
    "atne_model_complements":            "complements",
    "atne_model_illustration_translate": "illustration_translate",
    "atne_model_auditor":                "auditor",
}

_ALLOWED_MODELS = [
    "gemma-4-31b-it",
    "gemma-4-26b-a4b-it",
    "gemini-2.5-flash",
    "gemini-2.5-flash-lite",
    "gpt-4o-mini",
    "gpt-4o",
    "gpt-4.1-mini",
    "mistral-small-latest",
    "mistral-large-latest",
    "qwen/qwen3.5-27b",
    "qwen/qwen3.5-9b",
    "qwen/qwen3-235b-a22b:free",
    "qwen/qwen3-30b-a3b:free",
    "deepseek/deepseek-chat-v3-0324:free",
]


@app.post("/api/auth/exchange")
async def auth_exchange(request: Request, payload: dict = Body(...)):
    """Intercanvia token LaNet per cookie httpOnly ATNE (endpoint públic).

    El frontend crida aquest endpoint just després del redirect del bridge.
    A partir d'aquí les peticions autenticen via cookie (no Bearer header).
    """
    _rate_check(f"exchange:{request.client.host}", 20, 60)
    token = (payload.get("token") or "").strip()
    if not token:
        raise HTTPException(400, "Token manquant")
    status, login = _lanet_validate_token(token)
    if status == 503:
        raise HTTPException(503, "Bridge no disponible")
    if status != 200:
        raise HTTPException(401, "Token invàlid")
    session_val = _admin_sign(f"session:{login}:{int(time.time())}")
    is_prod = bool(os.getenv("K_SERVICE"))  # Cloud Run seta K_SERVICE automàticament
    resp = JSONResponse({"ok": True, "login": login})
    resp.set_cookie(
        "atne_session", session_val,
        httponly=True,
        secure=is_prod,
        samesite="strict",
        max_age=8 * 3600,
        path="/",
    )
    return resp


@app.post("/api/admin/login")
async def admin_login(request: Request, payload: dict = Body(...)):
    _rate_check(f"adminlogin:{request.client.host}", 5, 300)
    password = (payload.get("password") or "").strip()
    if not ADMIN_PASSWORD:
        raise HTTPException(500, "ADMIN_PASSWORD no configurat al servidor")
    if not secrets.compare_digest(password, ADMIN_PASSWORD):
        raise HTTPException(401, "Password incorrecte")
    token = _admin_sign(f"admin:{int(time.time())}")
    resp = JSONResponse({"ok": True, "ttl_seconds": ADMIN_SESSION_TTL_SEC})
    resp.set_cookie(
        "atne_admin",
        token,
        httponly=True,
        samesite="strict",
        max_age=ADMIN_SESSION_TTL_SEC,
        # secure=False a local, True a Cloud Run darrere HTTPS. El flag
        # s'activa quan FORCE_HTTPS_COOKIE és "1" (env var de producció).
        secure=os.getenv("FORCE_HTTPS_COOKIE", "").strip() == "1",
    )
    return resp


@app.post("/api/admin/logout")
async def admin_logout():
    resp = JSONResponse({"ok": True})
    resp.delete_cookie("atne_admin")
    return resp


@app.get("/api/admin/whoami")
async def admin_whoami(request: Request):
    """Check lleuger per al frontend: retorna si hi ha sessió activa (200) o no (401).
    Útil per saber si cal mostrar pantalla de login.
    Accepta tant la cookie admin HMAC com la sessió lanet si el login és admin.
    """
    if _admin_verify(request.cookies.get("atne_admin", "")):
        return {"ok": True}
    session_login = _verify_session(request.cookies.get("atne_session", ""))
    if session_login and _is_admin_login(session_login):
        return {"ok": True}
    raise HTTPException(401, "No autenticat")


@app.get("/api/runtime-config")
async def runtime_config():
    """Endpoint lleuger (no requereix auth) per al frontend docent.

    Exposa només els model_id actuals per fase i els costos aproximats.
    Serveix perquè el frontend pugui saber quin model s'està fent servir a
    cada fase i desar-ho a history.models_per_phase. NO retorna passwords
    ni claus API ni flags administratius.
    """
    return {
        "model_config": dict(_MODEL_CONFIG),
        "model_costs_eur_per_call": dict(_MODEL_COST_EUR_PER_CALL),
        "pilot_active": _PILOT_ACTIVE,
    }


@app.get("/api/admin/config")
async def admin_get_config(_: bool = Depends(_require_admin)):
    return {
        "model_config": dict(_MODEL_CONFIG),
        "auditor_enabled": _AUDITOR_ENABLED_RUNTIME,
        "admin_budget_eur_max": _ADMIN_BUDGET_EUR_MAX,
        "pilot_active": _PILOT_ACTIVE,
        "models_available": list(_ALLOWED_MODELS),
        "model_key_map": dict(_ALLOWED_MODEL_KEYS),
        "model_costs_eur_per_call": dict(_MODEL_COST_EUR_PER_CALL),
        "atne_model_default": ATNE_MODEL,
    }


@app.put("/api/admin/config")
async def admin_put_config(payload: dict = Body(...), _: bool = Depends(_require_admin)):
    """Actualitza una o més claus de system_config a Supabase i refresca
    _MODEL_CONFIG in-memory.

    Payload exemple:
    {
      "atne_model_adapt": "gpt-4o-mini",
      "atne_model_refine": "gemma-4-31b-it",
      "atne_auditor_enabled": false,
      "admin_budget_eur_max": 50,
      "pilot_active": true
    }
    """
    if not payload or not isinstance(payload, dict):
        raise HTTPException(400, "Cal proporcionar claus a actualitzar")
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        raise HTTPException(500, "Supabase no configurat al servidor")

    updates: list[dict] = []
    for key, value in payload.items():
        if key in _ALLOWED_MODEL_KEYS:
            # Admetem dues formes:
            # 1) str: "gemma-4-31b-it" → mode fix (retrocompat)
            # 2) dict: {"mode":"fixed","model":"..."} o
            #         {"mode":"rotate","models":[...],"strategy":"random"}
            if isinstance(value, str):
                if not value.strip():
                    raise HTTPException(400, f"{key} no pot ser una string buida")
                if value not in _ALLOWED_MODELS:
                    raise HTTPException(
                        400,
                        f"{key}={value} no és a _ALLOWED_MODELS: {_ALLOWED_MODELS}",
                    )
                stored_value = {
                    "mode": "fixed",
                    "model": value,
                    "set_by": "admin",
                }
            elif isinstance(value, dict):
                mode = value.get("mode", "fixed")
                if mode == "fixed":
                    m = value.get("model", "").strip()
                    if not m:
                        raise HTTPException(400, f"{key}: mode fixed requereix 'model'")
                    if m not in _ALLOWED_MODELS:
                        raise HTTPException(
                            400,
                            f"{key}={m} no és a _ALLOWED_MODELS: {_ALLOWED_MODELS}",
                        )
                    stored_value = {"mode": "fixed", "model": m, "set_by": "admin"}
                elif mode == "rotate":
                    models_list = value.get("models") or []
                    if not isinstance(models_list, list) or not models_list:
                        raise HTTPException(
                            400,
                            f"{key}: mode rotate requereix 'models' com a llista no buida",
                        )
                    for m in models_list:
                        if m not in _ALLOWED_MODELS:
                            raise HTTPException(
                                400,
                                f"{key}: model {m} no és a _ALLOWED_MODELS",
                            )
                    # Si només hi ha 1 model, reduim a mode fix (simplificació)
                    if len(models_list) == 1:
                        stored_value = {
                            "mode": "fixed",
                            "model": models_list[0],
                            "set_by": "admin",
                        }
                    else:
                        strategy = value.get("strategy", "random")
                        if strategy not in ("random",):
                            raise HTTPException(
                                400,
                                f"{key}: estrategia '{strategy}' no suportada (només 'random')",
                            )
                        stored_value = {
                            "mode": "rotate",
                            "models": list(models_list),
                            "strategy": strategy,
                            "set_by": "admin",
                        }
                else:
                    raise HTTPException(
                        400,
                        f"{key}: mode '{mode}' desconegut (esperava 'fixed' o 'rotate')",
                    )
            else:
                raise HTTPException(
                    400,
                    f"{key} ha de ser str (model_id) o dict (config rotate)",
                )
            updates.append({
                "key": key,
                "value": stored_value,
                "updated_by": "admin",
            })
        elif key == "atne_auditor_enabled":
            updates.append({"key": key, "value": bool(value), "updated_by": "admin"})
        elif key == "pilot_active":
            updates.append({"key": key, "value": bool(value), "updated_by": "admin"})
        elif key == "admin_budget_eur_max":
            try:
                v = float(value)
            except (TypeError, ValueError):
                raise HTTPException(400, f"{key} ha de ser numèric")
            if v < 0 or v > 10000:
                raise HTTPException(400, f"{key}={v} fora de rang raonable")
            updates.append({"key": key, "value": v, "updated_by": "admin"})
        else:
            raise HTTPException(400, f"Clau desconeguda: {key}")

    if not updates:
        raise HTTPException(400, "Cap clau vàlida per actualitzar")

    # UPSERT batch al Supabase. La PK és `key`, per tant on_conflict=key
    # fusiona els rows existents. Prefer: resolution=merge-duplicates és
    # la sintaxi de PostgREST per al comportament UPSERT.
    try:
        r = requests.post(
            f"{SUPABASE_URL}/rest/v1/system_config?on_conflict=key",
            headers={
                **SUPABASE_HEADERS,
                "Prefer": "resolution=merge-duplicates,return=minimal",
            },
            json=updates,
            timeout=10,
        )
        if r.status_code not in (200, 201, 204):
            raise RuntimeError(f"HTTP {r.status_code}: {r.text[:300]}")
    except Exception as e:
        raise HTTPException(500, f"Error escrivint a Supabase system_config: {e}")

    # Refrescar estat in-memory (re-llegeix TOT, no només les claus actualitzades)
    state = _load_system_config()
    return {"ok": True, "updated": len(updates), "state": state}


@app.delete("/api/admin/history")
async def admin_wipe_history(payload: dict = Body(default={}), _: bool = Depends(_require_admin)):
    """Esborra registres de l'historial. **Acció destructiva.**

    Payload opcional:
        confirm (str, obligatori): ha de ser exactament "FULL_WIPE" per esborrar
            TOTS els registres. Requisit de seguretat per evitar wipes accidentals.
        docent_hash (str, opcional): si s'especifica, només esborra els registres
            amb aquest docent_hash. Ignora confirm en aquest cas.

    Afegit al Sprint B (2026-04-16) per netejar la taula abans del pilot
    20/04-08/05 i començar amb mem\u00f2ria neta + 2-3 exemples de mostra.

    Retorna {ok, deleted_count}.
    """
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        raise HTTPException(500, "Supabase no configurat al servidor")

    docent_hash = (payload.get("docent_hash") or "").strip()
    confirm = (payload.get("confirm") or "").strip()

    if docent_hash:
        # Wipe per docent — no requereix confirm
        filter_clause = f"docent_hash=eq.{docent_hash}"
    else:
        # Full wipe — requereix confirm explícit
        if confirm != "FULL_WIPE":
            raise HTTPException(
                400,
                "Full wipe requereix payload {\"confirm\":\"FULL_WIPE\"}. "
                "Si vols netejar només els teus registres, passa docent_hash."
            )
        # PostgREST requereix un filtre per a DELETE; usem id=gt.0 com a
        # "tots els registres amb id > 0" (tots, perqu\u00e8 id \u00e9s serial).
        filter_clause = "id=gt.0"

    try:
        resp = requests.delete(
            f"{SUPABASE_URL}/rest/v1/history?{filter_clause}",
            headers={**SUPABASE_HEADERS, "Prefer": "return=representation"},
            timeout=15,
        )
        if resp.status_code not in (200, 204):
            raise RuntimeError(f"HTTP {resp.status_code}: {resp.text[:300]}")
        # Amb return=representation retorna l'array de rows esborrats
        deleted = 0
        try:
            deleted = len(resp.json()) if resp.text else 0
        except Exception:
            deleted = 0
        return {"ok": True, "deleted_count": deleted, "filter": filter_clause}
    except Exception as e:
        raise HTTPException(500, f"Error esborrant history: {e}")


# ── Perfils CRUD ────────────────────────────────────────────────────────────

@app.get("/api/profiles")
async def list_profiles():
    profiles = []
    for f in sorted(PROFILES_DIR.glob("*.json")):
        try:
            data = json.loads(f.read_text(encoding="utf-8"))
            profiles.append({"nom": data.get("nom", f.stem), "fitxer": f.stem})
        except Exception:
            pass
    return profiles


@app.get("/api/admin/analytics")
async def admin_analytics(_: bool = Depends(_require_admin)):
    """Retorna mètriques agregades de atne_sessions per al dashboard del pilot."""
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        return {"ok": False, "error": "Supabase no configurat"}

    def _get(url):
        r = requests.get(url, headers=SUPABASE_HEADERS, timeout=10)
        return r.json() if r.status_code == 200 else []

    base = f"{SUPABASE_URL}/rest/v1/atne_sessions"

    # Totes les sessions
    rows = _get(f"{base}?select=ts,model,profile_type,conditions,etapa,mecr_sortida,latency_ms,n_instructions,verify_score,docent_id&order=ts.desc&limit=500")

    if not rows:
        return {"ok": True, "total": 0, "by_model": {}, "by_etapa": {}, "by_profile": {},
                "docents_actius": 0, "latency_avg_ms": None, "sessions_avui": 0,
                "recent": []}

    from collections import Counter
    import datetime

    total = len(rows)
    avui = datetime.date.today().isoformat()
    sessions_avui = sum(1 for r in rows if (r.get("ts") or "").startswith(avui))

    by_model   = dict(Counter(r.get("model", "?") for r in rows).most_common())
    by_etapa   = dict(Counter(r.get("etapa", "?") for r in rows if r.get("etapa")).most_common())
    by_profile = dict(Counter(r.get("profile_type", "?") for r in rows).most_common())

    # Condicions (pot ser array)
    cond_counter: Counter = Counter()
    for r in rows:
        conds = r.get("conditions") or []
        for c in conds:
            if c and c != "desconegut":
                cond_counter[c] += 1
    by_conditions = dict(cond_counter.most_common(10))

    latencies = [r["latency_ms"] for r in rows if r.get("latency_ms")]
    latency_avg = int(sum(latencies) / len(latencies)) if latencies else None
    latency_p90 = int(sorted(latencies)[int(len(latencies) * 0.9)]) if len(latencies) >= 5 else None

    scores = [r["verify_score"] for r in rows if r.get("verify_score") is not None]
    verify_avg = round(sum(scores) / len(scores), 2) if scores else None

    docents = {r["docent_id"] for r in rows if r.get("docent_id")}

    recent = []
    for r in rows[:20]:
        ts = (r.get("ts") or "")[:16].replace("T", " ")
        recent.append({
            "ts": ts,
            "model": r.get("model", "?"),
            "profile_type": r.get("profile_type", "?"),
            "etapa": r.get("etapa", ""),
            "mecr": r.get("mecr_sortida", ""),
            "latency_ms": r.get("latency_ms"),
            "n_instr": r.get("n_instructions"),
            "score": r.get("verify_score"),
        })

    return {
        "ok": True,
        "total": total,
        "sessions_avui": sessions_avui,
        "docents_actius": len(docents),
        "by_model": by_model,
        "by_etapa": by_etapa,
        "by_profile": by_profile,
        "by_conditions": by_conditions,
        "latency_avg_ms": latency_avg,
        "latency_p90_ms": latency_p90,
        "verify_avg": verify_avg,
        "recent": recent,
    }


# ── Pilot — events UX granulars + consentiment ─────────────────────────────
#
# Sprint 1C (2026-04-22): no podem avaluar el pilot només amb la dada final
# (rating + adaptat). Cal el procés: refines, edits, complements
# generats/esborrats, exports, biblioteca, canvis de model. Aquests endpoints
# són *públics* (no admin) — els frontends els criden a cada acció rellevant.
#
# Whitelist d'event_types acceptats. Tot el que no estigui aquí es rebutja
# (evita pollution de la taula amb events oportunistics o tests).
_PILOT_EVENT_TYPES = {
    # Adaptació
    "adapt_started", "adapt_done", "adapt_error",
    # Refinaments
    "refine_started", "refined", "refine_submitted", "redo", "redo_rubric",
    # Complements
    "complement_generated", "complement_deleted", "complement_edited",
    # Edició / exports
    "manual_edit", "exported", "copied", "saved",
    # Navegació / biblioteca
    "biblioteca_opened", "draft_loaded", "pas_change",
    # Model / config
    "model_switch", "preset_applied",
    # Interaccions UI (pilot)
    "font_changed", "complement_toggled", "help_opened",
    # Navegació i comportament de lectura
    "page_view", "page_leave", "scroll_depth", "form_abandoned",
    # Edició post-generació
    "manual_edit_delta",
    # Feedback / rúbrica / suggeriments
    "rubric_submitted", "feedback_skipped", "feedback_submitted", "suggestion_submitted",
    # Errors client
    "client_error",
    # Consent flow
    "consent_shown",
}


@app.post("/api/pilot/event")
async def pilot_event(payload: dict = Body(...)):
    """Registra un event UX granular del pilot. Fire-and-forget al backend.

    Payload mínim: { event_type: str }
    Payload complet: {
        event_type: str,        # whitelist _PILOT_EVENT_TYPES
        session_id: str?,       # adapt_id o UUID client
        history_id: int?,       # FK lògica a history.id
        step: str?,             # 'pas1' | 'pas2' | 'pas3' | 'pas4'
        docent_id: str?,        # email (si conegut, sinó backend hash anònim)
        data: dict?,            # payload específic (preset, target, format…)
    }
    """
    event_type = (payload.get("event_type") or "").strip()
    if event_type not in _PILOT_EVENT_TYPES:
        return JSONResponse(
            {"ok": False, "error": f"event_type desconegut: {event_type!r}"},
            status_code=400,
        )

    docent_id = (payload.get("docent_id") or "").strip()
    docent_hash = _docent_hash_from_id(docent_id) if docent_id else \
                  (payload.get("docent_hash") or _get_current_docent_hash())

    event = {
        "event_type": event_type,
        "step": (payload.get("step") or None),
        "docent_id": docent_id or None,
        "docent_hash": docent_hash or None,
        "session_id": payload.get("session_id") or None,
        "history_id": payload.get("history_id") or None,
        "data": payload.get("data") or {},
        "prompt_version": payload.get("prompt_version") or ATNE_PROMPT_VERSION,
    }

    # Events crítics (dades irrecuperables si es perden) → insert SÍNCRON
    # amb retorn d'estat real al client. Així el frontend pot avisar el
    # docent si el seu suggeriment NO ha entrat i pot copiar-lo manualment.
    # La resta queden fire-and-forget per no bloquejar la UI.
    _CRITICAL_EVENTS = {"suggestion_submitted", "client_error"}
    if event_type in _CRITICAL_EVENTS:
        if not SUPABASE_URL or not _TELEMETRY_WRITE_KEY:
            _append_telemetry_fallback("atne_pilot_events", event, "supabase_no_configurat")
            return JSONResponse(
                {"ok": False, "error": "Supabase no configurat — guardat localment al servidor"},
                status_code=503,
            )
        try:
            r = requests.post(
                f"{SUPABASE_URL}/rest/v1/atne_pilot_events",
                headers={**SUPABASE_WRITE_HEADERS, "Prefer": "return=minimal"},
                json=event,
                timeout=8,
            )
            if r.status_code in (200, 201, 204):
                return {"ok": True}
            err = f"http_{r.status_code}: {r.text[:200]}"
            _append_telemetry_fallback("atne_pilot_events", event, err)
            return JSONResponse(
                {"ok": False, "error": err, "saved_to_fallback": True},
                status_code=502,
            )
        except Exception as e:
            _append_telemetry_fallback("atne_pilot_events", event, f"exception: {e!r}")
            return JSONResponse(
                {"ok": False, "error": str(e), "saved_to_fallback": True},
                status_code=502,
            )

    # Esdeveniments no crítics: fire-and-forget (el fallback intern del
    # _log_pilot_event ja captura errors a stdout/disc).
    import threading as _th
    _th.Thread(target=_log_pilot_event, args=(event,), daemon=True).start()
    return {"ok": True}


@app.post("/api/pilot/consent")
async def pilot_consent(request: Request, payload: dict = Body(...)):
    """Registra el consentiment informat del docent (RGPD art. 7 + AI Act).

    Payload: {
        docent_id: str (email),
        decision: 'accepted' | 'declined' | 'revoked',
        consent_text_version: str?  (default ve de system_config)
    }
    """
    docent_id = (payload.get("docent_id") or "").strip().lower()
    decision = (payload.get("decision") or "").strip().lower()
    if not docent_id or decision not in ("accepted", "declined", "revoked"):
        return JSONResponse(
            {"ok": False, "error": "docent_id obligatori i decision in {accepted,declined,revoked}"},
            status_code=400,
        )
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        return {"ok": False, "error": "Supabase no configurat"}

    docent_hash = _docent_hash_from_id(docent_id)
    ip = request.client.host if request.client else ""
    ip_hash = hashlib.sha256(f"{ip}:{os.getenv('ATNE_DOCENT_SALT', 'atne-pilot-2026')}".encode()).hexdigest()[:16]
    user_agent = (request.headers.get("user-agent") or "")[:300]

    row = {
        "docent_id": docent_id,
        "docent_hash": docent_hash,
        "decision": decision,
        "dpia_version": payload.get("dpia_version") or "2026-04-22",
        "consent_text_version": payload.get("consent_text_version") or "v1.0",
        "user_agent": user_agent,
        "ip_hash": ip_hash,
    }
    try:
        r = requests.post(
            f"{SUPABASE_URL}/rest/v1/atne_pilot_consent",
            headers={**SUPABASE_WRITE_HEADERS, "Prefer": "return=minimal"},
            json=row,
            timeout=5,
        )
        if r.status_code in (200, 201, 204):
            return {"ok": True}
        _append_telemetry_fallback(
            "atne_pilot_consent", row, f"http_{r.status_code}: {r.text[:400]}"
        )
        return {"ok": False, "error": r.text}
    except Exception as e:
        _append_telemetry_fallback("atne_pilot_consent", row, f"exception: {e!r}")
        return {"ok": False, "error": str(e)}


@app.get("/api/pilot/pending-feedback")
async def pilot_pending_feedback(docent_id: str = "", limit: int = 3):
    """Retorna adaptacions del docent que ha consumit (exported/copied) però
    no ha valorat (sense rating ni review_items). Sprint 1D follow-up: la
    home mostra una targeta convidant a valorar-les en una sessió posterior,
    quan ja les ha pogut usar a classe.

    Filtre:
      - docent_id (email) o docent_hash
      - exported=true OR copied=true
      - rating IS NULL AND review_items IS NULL
      - created_at > now - 7 days
      - límit configurable (default 3, cap 10)

    NO es protegeix amb auth admin: el docent ha de poder consultar les
    seves adaptacions pendents. El filtre per docent_id assegura que només
    veu les seves.
    """
    if not docent_id:
        return {"ok": True, "items": []}
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        return {"ok": False, "items": [], "error": "Supabase no configurat"}

    docent_id = docent_id.strip().lower()
    docent_hash = _docent_hash_from_id(docent_id)
    n = max(1, min(int(limit or 3), 10))

    import datetime
    cutoff = (datetime.datetime.utcnow() - datetime.timedelta(days=7)).isoformat() + "Z"

    # Filtrem per docent_hash (anonimització) i fem dues queries i unim:
    # una per exported=true, una per copied=true (PostgREST no facilita un OR
    # complex amb diversos camps amb facilitat). Després unim+deduplim per id.
    base = (
        f"{SUPABASE_URL}/rest/v1/history"
        f"?select=id,profile_name,created_at,exported,copied,original_text,model_used,prompt_version"
        f"&docent_hash=eq.{docent_hash}"
        f"&rating=is.null"
        f"&review_items=is.null"
        f"&created_at=gte.{cutoff}"
        f"&order=created_at.desc"
        f"&limit={n * 2}"  # marge per dedup
    )
    try:
        r1 = requests.get(base + "&exported=eq.true",
                          headers=SUPABASE_HEADERS, timeout=10)
        rows1 = r1.json() if r1.status_code == 200 else []
        r2 = requests.get(base + "&copied=eq.true",
                          headers=SUPABASE_HEADERS, timeout=10)
        rows2 = r2.json() if r2.status_code == 200 else []
    except Exception as e:
        return {"ok": False, "items": [], "error": str(e)}

    # Unim i deduplim per id, conservant ordre per created_at desc
    by_id: dict = {}
    for row in (rows1 + rows2):
        rid = row.get("id")
        if rid and rid not in by_id:
            by_id[rid] = row
    rows = sorted(
        by_id.values(),
        key=lambda r: r.get("created_at") or "",
        reverse=True,
    )[:n]

    items = []
    now = datetime.datetime.utcnow()
    for row in rows:
        # Calcular antiguitat per al display ("fa 2h", "fa 1 dia"…)
        age_text = ""
        try:
            ts = (row.get("created_at") or "").replace("Z", "")
            dt = datetime.datetime.fromisoformat(ts)
            diff = (now - dt).total_seconds()
            if diff < 3600:
                age_text = f"fa {max(1, int(diff / 60))} min"
            elif diff < 86400:
                age_text = f"fa {int(diff / 3600)} h"
            else:
                d = int(diff / 86400)
                age_text = "fa 1 dia" if d == 1 else f"fa {d} dies"
        except Exception:
            pass

        # Excerpt curt (primers 80 chars del text original) per identificar
        excerpt = (row.get("original_text") or "")[:80].strip()
        if len(row.get("original_text") or "") > 80:
            excerpt += "…"

        items.append({
            "id": row.get("id"),
            "profile_name": row.get("profile_name") or "Adaptació",
            "created_at": row.get("created_at"),
            "age_text": age_text,
            "excerpt": excerpt,
            "exported": bool(row.get("exported")),
            "copied": bool(row.get("copied")),
            "model_used": row.get("model_used"),
            "prompt_version": row.get("prompt_version"),
        })

    return {"ok": True, "items": items}


@app.get("/api/pilot/consent/{docent_id}")
async def pilot_consent_status(docent_id: str):
    """Retorna l'estat del consentiment d'un docent. Permet al frontend saber
    si cal redirigir a la pantalla de consent.html o no."""
    docent_id = docent_id.strip().lower()
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        return {"ok": False, "decision": None}
    try:
        r = requests.get(
            f"{SUPABASE_URL}/rest/v1/atne_pilot_consent"
            f"?docent_id=eq.{docent_id}&select=decision,ts,dpia_version&order=ts.desc&limit=1",
            headers=SUPABASE_HEADERS,
            timeout=5,
        )
        rows = r.json() if r.status_code == 200 else []
        if rows:
            return {"ok": True, "decision": rows[0].get("decision"), "ts": rows[0].get("ts")}
        return {"ok": True, "decision": None}
    except Exception as e:
        return {"ok": False, "error": str(e)}


@app.get("/api/admin/pilot-metrics")
async def admin_pilot_metrics(_: bool = Depends(_require_admin)):
    """Mètriques per al dashboard /admin/pilot.

    Combina 4 fonts:
      - atne_sessions  → sessions, models usats, latency, verify, cost
      - atne_pilot_events  → funnel UX (refines, complements, exports)
      - history  → feedback (rating, review_items), edit_manual, copied
      - atne_pilot_consent  → docents amb consentiment vàlid

    Tot per al pilot 2026-04-20 → 2026-05-08. Limita a 1000 rows per font.
    """
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        return {"ok": False, "error": "Supabase no configurat"}

    # IMPORTANT: usem SUPABASE_WRITE_HEADERS (service-key) per LLEGIR. Causa:
    # les escriptures de telemetria del pilot usen service-key per bypassar
    # RLS (atne_pilot_events típicament té RLS permissive INSERT-only per
    # anon). Si llegim amb anon-key, RLS bloqueja i veiem 0 resultats encara
    # que les dades existeixin. L'endpoint ja està protegit per _require_admin
    # → és segur usar service-key aquí (mai s'exposa al client).
    _read_headers = SUPABASE_WRITE_HEADERS if _USING_SERVICE_KEY else SUPABASE_HEADERS

    def _get(url):
        try:
            r = requests.get(url, headers=_read_headers, timeout=10)
            return r.json() if r.status_code == 200 else []
        except Exception:
            return []

    base = f"{SUPABASE_URL}/rest/v1"

    # Mapping docent_id → email local (hash → miquel.amor sense domini)
    docents_rows = _get(f"{base}/atne_docents?select=id,alias,email&limit=200")
    _id_to_alias: dict[str, str] = {}
    for d in docents_rows:
        did = d.get("id")
        if did:
            email = (d.get("email") or "").split("@")[0]
            _id_to_alias[did] = email or d.get("alias") or did[:8]

    def _alias(docent_id: str) -> str:
        if not docent_id:
            return "–"
        return _id_to_alias.get(docent_id, docent_id[:8])

    sessions = _get(
        f"{base}/atne_sessions"
        f"?select=ts,model,profile_type,conditions,etapa,mecr_sortida,latency_ms,"
        f"n_instructions,verify_score,docent_id,prompt_version,cost_estimat_eur"
        f"&order=ts.desc&limit=1000"
    )
    events_raw = _get(
        f"{base}/atne_pilot_events"
        f"?select=ts,event_type,docent_hash,docent_id,session_id,history_id,step,data,prompt_version"
        f"&order=ts.desc&limit=2000"
    )
    # Filtra events de probe (health + boot) que es generen automàticament
    # per detectar regressions del bug RLS 2026-05-15. No són events reals.
    def _is_probe(e):
        d = e.get("data") or {}
        return bool(d.get("_health_probe") or d.get("_boot_probe"))
    events = [e for e in events_raw if not _is_probe(e)]

    # Suggestions i client_errors: queries dedicades. El loop principal sobre
    # `events` només veu els 2000 events més recents — amb volum normal del
    # pilot (5-20 events/sessió), suggestions de fa pocs dies poden quedar
    # fora del window. Aquests queries garanteixen que TOTS els suggeriments
    # es vegin al dashboard, independentment del volum d'altres events.
    suggestions_dedicated = _get(
        f"{base}/atne_pilot_events"
        f"?select=ts,docent_id,data"
        f"&event_type=eq.suggestion_submitted"
        f"&order=ts.desc&limit=100"
    )
    client_errors_dedicated = _get(
        f"{base}/atne_pilot_events"
        f"?select=ts,docent_id,data"
        f"&event_type=eq.client_error"
        f"&order=ts.desc&limit=100"
    )
    # Refines: dos tipus d'events ens interessen amb el seu context complet.
    #  - refine_submitted: refinament del text (steppers + instrucció lliure)
    #  - redo_rubric: regeneració amb rúbrica (problemes marcats + observació)
    refines_dedicated = _get(
        f"{base}/atne_pilot_events"
        f"?select=ts,event_type,docent_id,data"
        f"&event_type=in.(refine_submitted,redo_rubric)"
        f"&order=ts.desc&limit=100"
    )
    feedback = _get(
        f"{base}/history"
        f"?select=created_at,rating,review_items,comment,refine_count,exported,copied,"
        f"edit_manual,time_on_step4_ms,model_used,prompt_version,docent_hash,cost_estimat_eur"
        f"&order=created_at.desc&limit=1000"
    )
    consents = _get(
        f"{base}/atne_pilot_consent"
        f"?select=ts,docent_id,decision&order=ts.desc&limit=500"
    )

    from collections import Counter, defaultdict
    import datetime

    today = datetime.date.today().isoformat()

    # ── Sessions agregats ───────────────────────────────────────────────
    total = len(sessions)
    sessions_avui = sum(1 for s in sessions if (s.get("ts") or "").startswith(today))
    docents_actius = len({s["docent_id"] for s in sessions if s.get("docent_id")})

    by_model = dict(Counter(s.get("model", "?") for s in sessions).most_common())
    by_etapa = dict(Counter(s.get("etapa", "?") for s in sessions if s.get("etapa")).most_common())
    by_profile = dict(Counter(s.get("profile_type", "?") for s in sessions).most_common())
    by_prompt_version = dict(Counter(s.get("prompt_version", "?") for s in sessions if s.get("prompt_version")).most_common())

    cond_counter: Counter = Counter()
    for s in sessions:
        for c in (s.get("conditions") or []):
            if c and c != "desconegut":
                cond_counter[c] += 1
    by_conditions = dict(cond_counter.most_common(10))

    latencies = [s["latency_ms"] for s in sessions if s.get("latency_ms")]
    latency_avg = int(sum(latencies) / len(latencies)) if latencies else None
    latency_p90 = int(sorted(latencies)[int(len(latencies) * 0.9)]) if len(latencies) >= 5 else None

    scores = [s["verify_score"] for s in sessions if s.get("verify_score") is not None]
    verify_avg = round(sum(scores) / len(scores), 2) if scores else None

    cost_total = round(sum(float(s.get("cost_estimat_eur") or 0) for s in sessions), 4)
    cost_per_model = defaultdict(float)
    for s in sessions:
        cost_per_model[s.get("model", "?")] += float(s.get("cost_estimat_eur") or 0)
    cost_per_model = {k: round(v, 4) for k, v in cost_per_model.items()}

    # ── Events UX agregats ──────────────────────────────────────────────
    by_event_type = dict(Counter(e.get("event_type", "?") for e in events).most_common())

    refines_per_session: dict[str, int] = defaultdict(int)
    complement_actions: Counter = Counter()
    exports_format: Counter = Counter()
    presets_used: Counter = Counter()
    for e in events:
        et = e.get("event_type")
        sid = e.get("session_id") or "no_sess"
        data = e.get("data") or {}
        if et == "refined":
            refines_per_session[sid] += 1
            if data.get("preset"):
                presets_used[data["preset"]] += 1
        elif et in ("complement_generated", "complement_deleted", "complement_edited"):
            ctype = data.get("type") or "?"
            complement_actions[f"{ctype}:{et.replace('complement_', '')}"] += 1
        elif et == "exported":
            exports_format[data.get("format") or "?"] += 1

    refine_dist = Counter(refines_per_session.values())
    refine_dist_summary = {
        "0": sum(1 for s in sessions if (s.get("session_id") or "no_sess") not in refines_per_session),
        "1": refine_dist.get(1, 0),
        "2": refine_dist.get(2, 0),
        "3": refine_dist.get(3, 0),
        "4+": sum(c for n, c in refine_dist.items() if n >= 4),
    }

    # ── Feedback agregat ────────────────────────────────────────────────
    rated = [h for h in feedback if h.get("rating") is not None]
    rating_dist = dict(Counter(h["rating"] for h in rated))
    rating_avg = round(sum(h["rating"] for h in rated) / len(rated), 2) if rated else None
    feedback_rate = round(len(rated) / len(feedback), 3) if feedback else None

    # Review items: agregar checkboxes (usar_classe / retocar / no_usar i altres)
    # provinents del pill «Valora aquesta adaptació» del Pas 3 + modal antic.
    review_items_counter: Counter = Counter()
    altres_texts: list[str] = []
    feedback_comments: list[str] = []
    for h in feedback:
        ri = h.get("review_items") or {}
        if isinstance(ri, dict):
            for k, v in ri.items():
                if k == "altres_text":
                    if v: altres_texts.append(str(v)[:200])
                elif v:
                    review_items_counter[k] += 1
        # comment del pill: text lliure que escriu el docent quan valora
        cmt = (h.get("comment") or "").strip()
        if cmt:
            ts = (h.get("created_at") or "")[:19].replace("T", " ")
            feedback_comments.append(f"[{ts}] [Valora] {cmt[:300]}")

    edit_rate = (sum(1 for h in feedback if h.get("edit_manual"))
                 / len(feedback)) if feedback else None
    copied_rate = (sum(1 for h in feedback if h.get("copied"))
                   / len(feedback)) if feedback else None
    exported_rate = (sum(1 for h in feedback if h.get("exported"))
                     / len(feedback)) if feedback else None

    times_step4 = [h["time_on_step4_ms"] for h in feedback if h.get("time_on_step4_ms")]
    median_time_step4 = sorted(times_step4)[len(times_step4) // 2] if times_step4 else None

    # ── Funnel d'adopció Stanford SCALE ─────────────────────────────────
    docents_seen = {c["docent_id"] for c in consents if c.get("docent_id")}
    docents_consented = {c["docent_id"] for c in consents if c.get("decision") == "accepted"}
    docents_adapted = {s["docent_id"] for s in sessions if s.get("docent_id")}
    docents_refined = {e["docent_id"] for e in events
                       if e.get("event_type") == "refined" and e.get("docent_id")}
    docents_exported = {e["docent_id"] for e in events
                        if e.get("event_type") == "exported" and e.get("docent_id")}

    funnel = {
        "consent_seen": len(docents_seen),
        "consent_accepted": len(docents_consented),
        "adapt_at_least_one": len(docents_adapted),
        "refine_at_least_one": len(docents_refined),
        "export_at_least_one": len(docents_exported),
    }

    # ── Recents (últimes 20 sessions) ───────────────────────────────────
    recent = []
    for s in sessions[:20]:
        ts = (s.get("ts") or "")[:16].replace("T", " ")
        recent.append({
            "ts": ts,
            "docent": _alias(s.get("docent_id", "")),
            "model": s.get("model", "?"),
            "etapa": s.get("etapa", ""),
            "mecr": s.get("mecr_sortida", ""),
            "latency_ms": s.get("latency_ms"),
            "score": s.get("verify_score"),
            "prompt_version": s.get("prompt_version"),
            "cost_eur": s.get("cost_estimat_eur"),
        })

    # ── Llista docents actius amb noms ──────────────────────────────────
    docents_actius_ids = {s["docent_id"] for s in sessions if s.get("docent_id")}
    docents_list = sorted([_alias(did) for did in docents_actius_ids])

    # ── Suggeriments i errors client (queries dedicades, NO depenen del
    #    window de 2000 events més recents).
    suggestions = []
    for e in suggestions_dedicated:
        data = e.get("data") or {}
        ts = (e.get("ts") or "")[:19].replace("T", " ")
        text = str(data.get("text") or "")[:300]
        if text and len(suggestions) < 100:
            suggestions.append({"ts": ts, "docent": _alias(e.get("docent_id") or ""), "text": text})

    client_errors = []
    for e in client_errors_dedicated:
        data = e.get("data") or {}
        ts = (e.get("ts") or "")[:19].replace("T", " ")
        if len(client_errors) < 100:
            client_errors.append({
                "ts": ts,
                "docent": _alias(e.get("docent_id") or ""),
                "page": str(data.get("page") or "")[:40],
                "msg": str(data.get("msg") or "")[:120],
                "type": data.get("type") or "js",
            })

    # Refines amb context (text/rúbrica): per a cada esdeveniment, exposem
    # QUÈ ha demanat el docent (problemes marcats + text lliure) perquè
    # el dashboard pugui mostrar-ho amb un format llegible.
    # També extreu els textos lliures cap a `comentaris_lliures` perquè
    # la secció de "Comentaris lliures" del dashboard tingui dades reals
    # (abans depenia del modal antic de /home, sempre buit).
    refines_detall = []
    comentaris_lliures: list[str] = []
    for e in refines_dedicated:
        data = e.get("data") or {}
        ts = (e.get("ts") or "")[:19].replace("T", " ")
        et = e.get("event_type")
        if et == "refine_submitted":
            user_text = (data.get("user_instruction") or "").strip()
            refines_detall.append({
                "ts": ts,
                "tipus": "refine_text",
                "docent": _alias(e.get("docent_id") or ""),
                "preset": data.get("preset"),
                "len": data.get("len"),
                "simp": data.get("simp"),
                "tone": data.get("tone"),
                "revisa_catala": bool(data.get("revisa_catala")),
                "instruccio": user_text or data.get("instruction_full"),
                "problems": [],
            })
            if user_text and len(comentaris_lliures) < 40:
                comentaris_lliures.append(f"[{ts}] [Refer] {user_text}")
        elif et == "redo_rubric":
            user_text = (data.get("user_observation") or "").strip()
            refines_detall.append({
                "ts": ts,
                "tipus": "rubric_redo",
                "docent": _alias(e.get("docent_id") or ""),
                "problems": data.get("problems") or [],
                "instruccio": user_text or None,
                "preserve_text": bool(data.get("preserve_text")),
                "preset": None, "len": None, "simp": None, "tone": None,
                "revisa_catala": False,
            })
            if user_text and len(comentaris_lliures) < 40:
                comentaris_lliures.append(f"[{ts}] [Rúbrica] {user_text}")

    body = {
        "ok": True,
        "prompt_version_actual": ATNE_PROMPT_VERSION,
        "totals": {
            "sessions": total,
            "sessions_avui": sessions_avui,
            "docents_actius": docents_actius,
            "events": len(events),
            "feedback_rebut": len(rated),
            "consents": len(consents),
            "cost_eur_total": cost_total,
        },
        "docents_list": docents_list,
        "by_model": by_model,
        "by_etapa": by_etapa,
        "by_profile": by_profile,
        "by_conditions": by_conditions,
        "by_prompt_version": by_prompt_version,
        "by_event_type": by_event_type,
        "complement_actions": dict(complement_actions.most_common(20)),
        "presets_used": dict(presets_used),
        "exports_format": dict(exports_format),
        "refine_distribution": refine_dist_summary,
        "latency_avg_ms": latency_avg,
        "latency_p90_ms": latency_p90,
        "verify_avg": verify_avg,
        "feedback": {
            "rate": feedback_rate,
            "rating_avg": rating_avg,
            "rating_dist": rating_dist,
            "review_items": dict(review_items_counter),
            # Combinem comentaris lliures de 3 fonts:
            # - pill «Valora aquesta adaptació» (history.comment): font principal
            # - popovers Refer/Rúbrica del Pas 3 (events refine_submitted/redo_rubric)
            # - modal antic de /home (history.review_items.altres_text): residuals
            "altres_texts": (feedback_comments + comentaris_lliures + altres_texts)[:50],
            # Checkboxes d'ús del pill «Valora»: comptadors d'intenció
            "usage_intent": {
                "usar_classe": review_items_counter.get("usar_classe", 0),
                "retocar":     review_items_counter.get("retocar", 0),
                "no_usar":     review_items_counter.get("no_usar", 0),
            },
            "edit_rate": edit_rate,
            "copied_rate": copied_rate,
            "exported_rate": exported_rate,
            "median_time_step4_ms": median_time_step4,
        },
        "cost_per_model_eur": cost_per_model,
        "funnel": funnel,
        "recent": recent,
        "suggestions": suggestions,
        "client_errors": client_errors,
        "refines_detall": refines_detall,
    }
    # Cap cache: les dades de pilot canvien constantment. Sense aquest header,
    # alguns navegadors/proxies cachejaven la resposta i el docent veia dades
    # antigues després d'una adaptació recent.
    return JSONResponse(content=body, headers={
        "Cache-Control": "no-store, no-cache, must-revalidate",
        "Pragma": "no-cache",
    })


@app.get("/api/admin/pilot/fallback")
async def admin_pilot_fallback_list(_: bool = Depends(_require_admin)):
    """Llista els events de telemetria que no van poder entrar a Supabase
    (RLS, 5xx, xarxa…) i van quedar al fitxer fallback NDJSON local.

    NOTA: a Cloud Run el directori /tmp és tmpfs (efímer). Els events
    queden també a stdout amb el tag `[TELEMETRY_FAILED]` (durable a
    Cloud Logging). Per recuperació real cal mirar Cloud Logging.
    """
    if not _PILOT_FALLBACK_FILE.exists():
        return {"ok": True, "items": [], "path": str(_PILOT_FALLBACK_FILE), "msg": "Sense fallback registrat."}
    items = []
    try:
        with _PILOT_FALLBACK_FILE.open("r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    items.append(json.loads(line))
                except Exception:
                    items.append({"_raw": line[:500]})
    except Exception as e:
        return {"ok": False, "error": str(e)}
    return {
        "ok": True,
        "path": str(_PILOT_FALLBACK_FILE),
        "count": len(items),
        "items": items[-200:],  # darreres 200
        "using_service_key": _USING_SERVICE_KEY,
    }


@app.post("/api/admin/pilot/replay-fallback")
async def admin_pilot_fallback_replay(_: bool = Depends(_require_admin)):
    """Reintenta inserir a Supabase tots els events del fallback NDJSON.
    Si l'insert té èxit l'event es marca com a `_replayed` i s'arxiva
    a un fitxer rotat (`.ndjson.replayed-<ts>`). No esborra dades.
    """
    if not _PILOT_FALLBACK_FILE.exists():
        return {"ok": True, "replayed": 0, "failed": 0, "msg": "Sense fallback."}
    if not SUPABASE_URL or not _TELEMETRY_WRITE_KEY:
        return {"ok": False, "error": "Supabase no configurat"}

    items = []
    try:
        with _PILOT_FALLBACK_FILE.open("r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    items.append(json.loads(line))
                except Exception:
                    continue
    except Exception as e:
        return {"ok": False, "error": str(e)}

    replayed = 0
    failed = 0
    still_failing: list[dict] = []
    for rec in items:
        kind = rec.get("kind")
        payload = rec.get("payload") or {}
        if not kind or not payload:
            continue
        table = {
            "atne_pilot_events": "atne_pilot_events",
            "atne_pilot_consent": "atne_pilot_consent",
            "atne_sessions": "atne_sessions",
        }.get(kind)
        if not table:
            still_failing.append(rec)
            continue
        try:
            r = requests.post(
                f"{SUPABASE_URL}/rest/v1/{table}",
                headers={**SUPABASE_WRITE_HEADERS, "Prefer": "return=minimal"},
                json=payload,
                timeout=8,
            )
            if r.status_code in (200, 201, 204):
                replayed += 1
            else:
                failed += 1
                rec["replay_error"] = f"http_{r.status_code}: {r.text[:200]}"
                still_failing.append(rec)
        except Exception as e:
            failed += 1
            rec["replay_error"] = f"exception: {e!r}"
            still_failing.append(rec)

    # Arxiva el fitxer original i reescriu només amb els que segueixen fallant
    import datetime as _dt
    try:
        archive = _PILOT_FALLBACK_FILE.with_suffix(
            f".ndjson.replayed-{_dt.datetime.utcnow().strftime('%Y%m%dT%H%M%SZ')}"
        )
        _PILOT_FALLBACK_FILE.rename(archive)
    except Exception:
        archive = None
    if still_failing:
        try:
            with _PILOT_FALLBACK_FILE.open("w", encoding="utf-8") as f:
                for rec in still_failing:
                    f.write(json.dumps(rec, ensure_ascii=False, default=str) + "\n")
        except Exception:
            pass

    return {
        "ok": True,
        "replayed": replayed,
        "failed": failed,
        "archived_to": str(archive) if archive else None,
    }


@app.get("/api/audit/last-adaptation")
async def audit_last_adaptation(_: bool = Depends(_require_admin)):
    """Retorna la darrera adaptació (retrocompatibilitat amb UI antiga)."""
    if not _orchestrator._ATNE_LAST_ADAPTATION:
        return {"ok": True, "empty": True, "msg": "No hi ha cap adaptació registrada des del reinici del servidor."}
    return {"ok": True, "empty": False, "data": _orchestrator._ATNE_LAST_ADAPTATION}


def _persist_adaptation_to_supabase(entry: dict) -> None:
    """Persisteix l'entrada del buffer a la taula atne_prompt_debug de Supabase.
    Fire-and-forget (errors només es loguegen, no aturen l'adaptació)."""
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        return
    try:
        requests.post(
            f"{SUPABASE_URL}/rest/v1/atne_prompt_debug",
            headers={**SUPABASE_HEADERS, "Prefer": "return=minimal"},
            json={
                "adapt_id": entry.get("id"),
                "ts_ms": int(entry.get("ts", 0) * 1000),
                "docent_id": entry.get("docent_id") or None,
                "model": entry.get("model"),
                "data": entry,  # jsonb amb tot
            },
            timeout=5,
        )
    except Exception as e:
        print(f"[audit] error persist atne_prompt_debug: {e}", flush=True)


def _list_adaptations_from_supabase(limit: int = 20) -> list:
    """Fallback quan el buffer en memòria està buit (ex: Cloud Run amb múltiples
    instàncies). Llegeix les últimes N de Supabase."""
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        return []
    try:
        r = requests.get(
            f"{SUPABASE_URL}/rest/v1/atne_prompt_debug"
            f"?select=data&order=ts_ms.desc&limit={limit}",
            headers=SUPABASE_HEADERS, timeout=10,
        )
        if r.status_code == 200:
            return [row["data"] for row in r.json() if row.get("data")]
    except Exception as e:
        print(f"[audit] error list atne_prompt_debug: {e}", flush=True)
    return []


def _get_adaptation_from_supabase(adapt_id: str) -> dict:
    """Recupera una adaptació per id si no és al buffer en memòria."""
    if not SUPABASE_URL or not SUPABASE_ANON_KEY:
        return {}
    try:
        r = requests.get(
            f"{SUPABASE_URL}/rest/v1/atne_prompt_debug"
            f"?adapt_id=eq.{_urlquote(str(adapt_id), safe='')}&select=data&limit=1",
            headers=SUPABASE_HEADERS, timeout=5,
        )
        if r.status_code == 200:
            rows = r.json()
            if rows:
                return rows[0].get("data", {})
    except Exception as e:
        print(f"[audit] error get atne_prompt_debug: {e}", flush=True)
    return {}


@app.get("/api/audit/adaptations")
async def audit_adaptations_list(_: bool = Depends(_require_admin)):
    """Retorna un resum de les últimes N adaptacions (N més recents primer).

    Primer mira el buffer en memòria. Si està buit (Cloud Run amb múltiples
    instàncies o instància acabada de reiniciar), cau a Supabase.
    """
    _log = _orchestrator._ATNE_ADAPTATIONS_LOG
    source_entries = list(reversed(_log)) if _log \
                     else _list_adaptations_from_supabase(limit=_orchestrator._ATNE_ADAPTATIONS_MAX)
    source = "memory" if _log else "supabase"
    summary = []
    for e in source_entries:
        chars = (e.get("profile") or {}).get("caracteristiques") or {}
        n_cond = sum(1 for v in chars.values() if isinstance(v, dict) and v.get("actiu"))
        summary.append({
            "id": e.get("id"),
            "ts": e.get("ts"),
            "iso": e.get("iso"),
            "docent_id": e.get("docent_id"),
            "model": e.get("model"),
            "mecr_sortida": (e.get("params") or {}).get("mecr_sortida"),
            "n_instructions": e.get("n_instructions"),
            "n_conditions": n_cond,
            "input_chars": e.get("text_input_len"),
            "output_chars": e.get("adapted_output_len"),
            "profile_name": (e.get("profile") or {}).get("nom", ""),
        })
    return {"ok": True, "count": len(summary), "adaptations": summary, "source": source}


@app.get("/api/audit/adaptations/{adapt_id}")
async def audit_adaptation_detail(adapt_id: str, _: bool = Depends(_require_admin)):
    """Retorna el detall d'una adaptació concreta: buffer primer, Supabase fallback."""
    for e in _orchestrator._ATNE_ADAPTATIONS_LOG:
        if e.get("id") == adapt_id:
            return {"ok": True, "data": e, "source": "memory"}
    sb = _get_adaptation_from_supabase(adapt_id)
    if sb:
        return {"ok": True, "data": sb, "source": "supabase"}
    return {"ok": False, "error": "No trobada ni al buffer ni a Supabase"}


@app.post("/api/audit/instruction-map")
async def audit_instruction_map(
    payload: dict = Body(...),
    _: bool = Depends(_require_admin),
):
    """Auditoria en viu: quines instruccions s'activen per a una combinació.

    Input JSON:
      {
        "profile": {caracteristiques: {tdah: {actiu, grau, ...}, ...}},
        "params": {mecr_sortida, dua, complements, ...}
      }

    Retorna per a cada instrucció del catàleg:
      - id, text, macro, activation
      - status: 'ACTIVE' | 'SUPPRESSED' | 'NOT_MATCHED'
      - reason: perquè queda fora si NOT_MATCHED o SUPPRESSED
    Més estadístiques i el resultat filtrat agrupat per macro.
    """
    from instruction_catalog import CATALOG, MACRODIRECTIVES

    profile = payload.get("profile") or {"caracteristiques": {}}
    params = payload.get("params") or {}

    # Filtra via la funció real que fa servir el pipeline
    filtered = instruction_filter.get_instructions(profile, params)

    active_ids = set()
    for macro_id, macro in (filtered.get("macrodirectives") or {}).items():
        for instr in macro.get("instruccions", []):
            active_ids.add(instr["id"])
    suppressed_ids = set(filtered.get("suppressed") or [])
    audit_list = filtered.get("audit") or []
    audit_by_id = {a["id"]: a for a in audit_list}

    # Estat de TOTES les instruccions del catàleg
    all_rows = []
    for iid, instr in CATALOG.items():
        activation = instr.get("activation", "?")
        macro_id = instr.get("macro", "ALTRES")
        macro_label = MACRODIRECTIVES.get(macro_id, {}).get("label", macro_id)
        if iid in active_ids:
            status = "ACTIVE"
            reason = "inclosa al prompt"
        elif iid in suppressed_ids:
            status = "SUPPRESSED"
            reason = audit_by_id.get(iid, {}).get("motiu", "suprimida")
        else:
            # No activada per les condicions actuals (p.ex. PERFIL sense perfil actiu)
            status = "NOT_MATCHED"
            if activation == "PERFIL":
                reason = f"perfils objectiu {instr.get('profiles', [])} no actius"
            elif activation == "COMPLEMENT":
                reason = f"complement '{instr.get('complement', '?')}' no triat"
            elif activation == "NIVELL":
                reason = f"MECR '{params.get('mecr_sortida', '?')}' no dispara"
            else:
                reason = "no inclosa per condicions"
        all_rows.append({
            "id": iid,
            "text": instr.get("text", "")[:300],
            "macro_id": macro_id,
            "macro_label": macro_label,
            "activation": activation,
            "status": status,
            "reason": reason,
        })

    # Resum per activació i status
    counts = {
        "total": len(all_rows),
        "active": len(active_ids),
        "suppressed": len(suppressed_ids),
        "not_matched": len(all_rows) - len(active_ids) - len(suppressed_ids),
        "by_activation": {},
    }
    for row in all_rows:
        act = row["activation"]
        counts["by_activation"].setdefault(act, {"active": 0, "suppressed": 0, "not_matched": 0})
        counts["by_activation"][act][row["status"].lower()] += 1

    return {
        "ok": True,
        "params": params,
        "profile_summary": {
            "caracteristiques_actives": [
                k for k, v in (profile.get("caracteristiques") or {}).items()
                if isinstance(v, dict) and v.get("actiu")
            ],
        },
        "counts": counts,
        "macrodirectives": filtered.get("macrodirectives", {}),
        "instructions": all_rows,
    }


@app.post("/api/profiles")
async def save_profile(payload: dict = Body(...)):
    nom = payload.get("nom", "sense_nom")
    # Sanititzar nom per a filename
    safe = re.sub(r'[^\w\s-]', '', nom).strip().replace(" ", "_")[:50]
    if not safe:
        safe = "perfil"
    path = PROFILES_DIR / f"{safe}.json"
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"ok": True, "fitxer": safe}


@app.get("/api/profiles/{nom}")
async def load_profile(nom: str):
    path = PROFILES_DIR / f"{nom}.json"
    if not path.exists():
        raise HTTPException(404, "Perfil no trobat")
    return json.loads(path.read_text(encoding="utf-8"))


@app.delete("/api/profiles/{nom}")
async def delete_profile(nom: str):
    path = PROFILES_DIR / f"{nom}.json"
    if path.exists():
        path.unlink()
    return {"ok": True}


# ── Historial i feedback ──────────────────────────────────────────────────

@app.get("/api/history")
async def list_history(limit: int = 30):
    """Llista les últimes adaptacions de l'historial.

    Camps retornats (Sprint B 2026-04-16):
    - id, created_at, profile_name, original_text, adapted_text (preview)
    - profile_json, context_json, params_json (per "carregar text + perfil")
    - rating (feedback del docent)
    - source (paste|upload|generated) — pilot anònima, Sprint B
    - model_used (quin LLM va adaptar) — Sprint 1A

    Ordenat per created_at desc, limitat a `limit` (default 30).
    """
    # Bug 5 (2026-04-19): abans retornàvem sempre HTTP 200 quan Supabase fallava
    # (el frontend ho llegia com "cap document recent" i amagava l'error real).
    # Ara els errors de xarxa/Supabase retornen HTTP 503 perquè el frontend,
    # que ja gestiona status >= 400, els mostri al docent.
    try:
        resp = requests.get(
            f"{SUPABASE_URL}/rest/v1/history"
            f"?select=id,created_at,profile_name,original_text,adapted_text,"
            f"profile_json,context_json,params_json,rating,source,model_used,etapa,curs"
            f"&order=created_at.desc&limit={limit}",
            headers=SUPABASE_HEADERS,
            timeout=10,
        )
        if resp.status_code == 200:
            return {"ok": True, "items": resp.json()}
        return JSONResponse(
            {"ok": False, "error": resp.text or f"Supabase HTTP {resp.status_code}"},
            status_code=503,
        )
    except Exception as e:
        return JSONResponse(
            {"ok": False, "error": str(e) or type(e).__name__},
            status_code=503,
        )


# Sprint 1A+1B: columnes ampliades de history acceptades per POST i PATCH
# (veure docs/sql/sprint1a_alter_history.sql i sprint1b_admin_config.sql)
_HISTORY_INSERTABLE_FIELDS = {
    # Legacy
    "profile_name", "profile_json", "context_json", "params_json",
    "original_text", "adapted_text",
    # Sprint 1A — instrumentació
    "model_used", "endpoint", "duration_ms", "refine_count", "edit_manual",
    "exported", "etapa", "curs", "perfil_kind", "via",
    "n_words_in", "n_words_out", "docent_hash", "quality_summary",
    "auditor_used",
    # Sprint 1B — selector model per fase + captures pilot
    "models_per_phase", "cost_estimat_eur", "copied",
    "time_on_step4_ms", "review_items",
    # Sprint B (2026-04-16) — memòria pilot anònima
    "source",  # 'paste' | 'upload' | 'generated'
    # Sprint 1C (2026-04-22) — versionat del prompt
    "prompt_version",
}

_HISTORY_UPDATABLE_FIELDS = {
    # Feedback Sprint 0
    "rating", "comment",
    # Sprint 1A — feedback ampliat
    "fb_used_in_class", "fb_needs_redo", "fb_level_ok",
    # Sprint 1A — instrumentació post-generació
    "refine_count", "edit_manual", "exported", "duration_ms",
    "adapted_text", "quality_summary",
    # Sprint 1B — captures pilot
    "copied", "time_on_step4_ms", "review_items", "cost_estimat_eur",
    "models_per_phase",
    # Sprint 1C — versionat del prompt
    "prompt_version",
}


def _get_current_docent_hash() -> str:
    """Retorna un hash del docent actual per distingir-lo al dashboard sense
    revelar la identitat. Al pilot inicial ve d'env var; a Cloud Run amb
    IAP vindrà del header X-Forwarded-User-Email (TODO pilot 2)."""
    email = os.getenv("ATNE_DOCENT_EMAIL", "anonim@fje.edu").strip().lower()
    salt = os.getenv("ATNE_DOCENT_SALT", "atne-pilot-2026")
    h = hashlib.sha256(f"{email}:{salt}".encode()).hexdigest()
    return h[:16]


@app.get("/api/history/recent")
async def history_recent(docent_id: str = "", limit: int = 5):
    """Últimes adaptacions del docent per a la home (Flash + Taller)."""
    if not docent_id:
        return {"ok": True, "items": []}
    try:
        resp = requests.get(
            f"{SUPABASE_URL}/rest/v1/history"
            f"?docent_hash=eq.{docent_id}"
            f"&select=id,created_at,profile_name,mode"
            f"&order=created_at.desc&limit={limit}",
            headers=SUPABASE_HEADERS, timeout=5,
        )
        if resp.status_code == 200:
            return {"ok": True, "items": resp.json()}
        return {"ok": True, "items": []}
    except Exception:
        return {"ok": True, "items": []}


@app.post("/api/history")
async def save_history(payload: dict = Body(...)):
    """Desa una adaptació a l'historial de Supabase.

    Accepta camps legacy (profile, context, params, original, adapted) i els
    camps nous Sprint 1A+1B (model_used, models_per_phase, duration_ms,
    cost_estimat_eur, copied, time_on_step4_ms, review_items, etc.). Els
    camps no reconeguts es descarten silenciosament (additive compat).
    """
    row: dict = {
        # Compat: els camps del frontend actual porten alies diferents
        "profile_name": payload.get("profile_name", ""),
        "profile_json": payload.get("profile", payload.get("profile_json", {})),
        "context_json": payload.get("context", payload.get("context_json", {})),
        "params_json": payload.get("params", payload.get("params_json", {})),
        "original_text": payload.get("original", payload.get("original_text", "")),
        "adapted_text": payload.get("adapted", payload.get("adapted_text", "")),
    }
    # Camps nous Sprint 1A+1B: només els que arriben explícitament
    for field in _HISTORY_INSERTABLE_FIELDS:
        if field in payload and field not in row:
            row[field] = payload[field]
    # Omplir docent_hash si el frontend no l'ha enviat (backend-side)
    if "docent_hash" not in row:
        row["docent_hash"] = _get_current_docent_hash()
    # Default endpoint
    if "endpoint" not in row:
        row["endpoint"] = payload.get("endpoint") or "/api/adapt"
    # Sprint 1C: estampar prompt_version automàticament al backend si el client
    # no l'ha enviat. Permet correlacionar adaptacions amb la versió del codi.
    if "prompt_version" not in row:
        row["prompt_version"] = ATNE_PROMPT_VERSION

    # Sprint 1C: estimar cost_estimat_eur al backend (per-token) si el client
    # no l'ha enviat o ha enviat 0 (estimació per-call coarse). El backend té
    # informació més precisa: model_used + n_words_in + n_words_out.
    try:
        from adaptation.pricing import estimate_cost_eur as _est_cost
        existing_cost = row.get("cost_estimat_eur")
        if existing_cost is None or existing_cost == 0:
            n_in = row.get("n_words_in") or 0
            n_out = row.get("n_words_out") or 0
            # Aprox 6 chars/paraula per a català (incloent espais)
            chars_in = int(n_in) * 6
            chars_out = int(n_out) * 6
            model = row.get("model_used")
            cost = _est_cost(model, chars_in, chars_out)
            if cost is not None:
                row["cost_estimat_eur"] = cost
    except Exception as _e:
        print(f"[history] error estimant cost: {_e}", flush=True)

    try:
        resp = requests.post(
            f"{SUPABASE_URL}/rest/v1/history",
            headers={**SUPABASE_HEADERS, "Prefer": "return=representation"},
            json=row,
            timeout=10,
        )
        if resp.status_code in (200, 201):
            data = resp.json()
            return {"ok": True, "id": data[0]["id"] if data else None}
        return {"ok": False, "error": resp.text}
    except Exception as e:
        return {"ok": False, "error": str(e)}


@app.post("/api/history/{history_id}/beacon")
async def history_beacon(history_id: int, payload: dict = Body(...)):
    """Endpoint POST equivalent al PATCH per a navigator.sendBeacon().

    sendBeacon() només fa POST i serveix per enviar dades en moments de
    teardown de la pàgina (beforeunload). Al pilot, el fem servir per
    capturar el `time_on_step4_ms` quan el docent tanca la pestanya sense
    passar pel botó "Nova adaptació".
    """
    # Reutilitza la whitelist de l'update tradicional
    update: dict = {}
    for field in _HISTORY_UPDATABLE_FIELDS:
        if field in payload:
            update[field] = payload[field]
    if not update:
        return {"ok": False, "error": "Cap camp conegut al payload"}
    try:
        resp = requests.patch(
            f"{SUPABASE_URL}/rest/v1/history?id=eq.{history_id}",
            headers={**SUPABASE_HEADERS, "Prefer": "return=minimal"},
            json=update,
            timeout=5,
        )
        return {"ok": resp.status_code in (200, 204)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


@app.patch("/api/history/{history_id}")
async def update_history_feedback(history_id: int, payload: dict = Body(...)):
    """Actualitza camps d'una entrada de l'historial.

    Whitelist: només camps a _HISTORY_UPDATABLE_FIELDS + rating+comment
    legacy. Si el payload porta `rating` o `comment`, també seta `rated_at`
    a now() perquè el dashboard pugui filtrar per feedback rebut.
    """
    update: dict = {}
    has_feedback = False
    for field in _HISTORY_UPDATABLE_FIELDS:
        if field in payload:
            update[field] = payload[field]
            if field in ("rating", "comment"):
                has_feedback = True
    if not update:
        return {"ok": False, "error": "Cap camp conegut al payload"}
    if has_feedback:
        update["rated_at"] = "now()"
    try:
        resp = requests.patch(
            f"{SUPABASE_URL}/rest/v1/history?id=eq.{history_id}",
            headers={**SUPABASE_HEADERS, "Prefer": "return=representation"},
            json=update,
            timeout=10,
        )
        if resp.status_code in (200, 204):
            return {"ok": True}
        return {"ok": False, "error": resp.text}
    except Exception as e:
        return {"ok": False, "error": str(e)}


# ── Extracció de text des de fitxer (PDF/DOCX/MD/TXT) ─────────────────────

MAX_UPLOAD_BYTES = 5 * 1024 * 1024  # 5 MB

@app.post("/api/extract-text")
async def extract_text_from_file(file: UploadFile = File(...)):
    """
    Rep un fitxer PDF/DOCX/MD/TXT i retorna el text pla extret.
    Límit: 5 MB. Format detectat per l'extensió.
    """
    filename = (file.filename or "").lower()
    ext = filename.rsplit(".", 1)[-1] if "." in filename else ""

    if ext not in ("pdf", "docx", "md", "txt"):
        return JSONResponse(
            {"error": f"Format no suportat: .{ext}. Accepta: PDF, DOCX, MD, TXT."},
            status_code=400,
        )

    raw = await file.read()
    if len(raw) > MAX_UPLOAD_BYTES:
        return JSONResponse(
            {"error": f"Fitxer massa gran ({len(raw)//1024} KB). Màxim: 5 MB."},
            status_code=400,
        )
    if not raw:
        return JSONResponse({"error": "Fitxer buit."}, status_code=400)

    text = ""
    try:
        if ext == "txt" or ext == "md":
            # Provar UTF-8 primer, fallback a latin-1
            try:
                text = raw.decode("utf-8")
            except UnicodeDecodeError:
                text = raw.decode("latin-1", errors="replace")
            # Per MD, mantenim el text tal qual (l'LLM entén markdown)

        elif ext == "pdf":
            import io
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw))
            pages = []
            for page in reader.pages:
                try:
                    pages.append(page.extract_text() or "")
                except Exception:
                    continue
            text = "\n\n".join(p.strip() for p in pages if p.strip())

        elif ext == "docx":
            import io
            from docx import Document  # python-docx
            doc = Document(io.BytesIO(raw))
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(parts)

    except Exception as e:
        return JSONResponse(
            {"error": _safe_error(e, "No s'ha pogut extreure el text")},
            status_code=500,
        )

    text = text.strip()
    if not text:
        return JSONResponse(
            {"error": "No s'ha pogut extreure text llegible del fitxer. "
                      "Si és un PDF escanejat, caldria OCR (no suportat ara)."},
            status_code=422,
        )

    paraules = len(text.split())
    return {
        "text": text,
        "paraules": paraules,
        "format_detectat": ext,
        "filename": file.filename,
    }


# ── Adaptació de documents amb preservació de format ─────────────────────

from fastapi import Form as _Form  # noqa: E402

@app.post("/api/adapt-pdf")
async def adapt_pdf_document(
    file: UploadFile = File(...),
    profile: str = _Form("{}"),
    context: str = _Form("{}"),
    params: str = _Form("{}"),
    model: str = _Form(""),
    output_format: str = _Form("pdf"),
):
    """
    Rep un PDF, adapta el text preservant el layout i retorna el PDF adaptat.
    Els camps profile/context/params són JSON strings (mateixos camps que /api/adapt).
    """
    import io
    from adaptation.document_adapter import (
        is_scanned_pdf, extract_pdf_text_map, inject_pdf_adapted, batch_adapt_text_map,
    )
    from adaptation.prompt_builder import build_system_prompt

    raw = await file.read()
    if len(raw) > MAX_UPLOAD_BYTES:
        return JSONResponse({"error": f"Fitxer massa gran ({len(raw)//1024} KB). Màxim: 5 MB."}, status_code=400)
    if not raw:
        return JSONResponse({"error": "Fitxer buit."}, status_code=400)

    try:
        profile_d = json.loads(profile)
        context_d = json.loads(context)
        params_d = json.loads(params)
    except json.JSONDecodeError as e:
        return JSONResponse({"error": f"JSON invàlid als paràmetres: {e}"}, status_code=400)

    if is_scanned_pdf(raw):
        return JSONResponse(
            {"error": "PDF escanejat (sense text seleccionable). "
                      "Exporta el document com a PDF des del programa original i torna-ho a provar."},
            status_code=422,
        )

    text_map = extract_pdf_text_map(raw)
    if not text_map:
        return JSONResponse({"error": "No s'ha trobat text adaptable al PDF."}, status_code=422)

    try:
        params_d["document_mode"] = True
        system_prompt = build_system_prompt(profile_d, context_d, params_d)
        active_model = _model_for("adapt", override=model)
        adapted = batch_adapt_text_map(text_map, active_model, system_prompt)
    except Exception as e:
        return JSONResponse({"error": _safe_error(e, "Error durant l'adaptació del PDF")}, status_code=500)

    base_name = (file.filename or "document").rsplit(".", 1)[0]

    if output_format == "docx":
        from adaptation.document_adapter import build_docx_from_adapted
        try:
            docx_out = build_docx_from_adapted(text_map, adapted)
        except Exception as e:
            return JSONResponse({"error": _safe_error(e, "Error generant el Word")}, status_code=500)
        filename_out = base_name + "_adaptat.docx"
        return StreamingResponse(
            io.BytesIO(docx_out),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f'attachment; filename="{filename_out}"'},
        )

    try:
        pdf_out = inject_pdf_adapted(raw, adapted)
    except Exception as e:
        return JSONResponse({"error": _safe_error(e, "Error reinjectant el PDF")}, status_code=500)

    filename_out = base_name + "_adaptat.pdf"
    return StreamingResponse(
        io.BytesIO(pdf_out),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename_out}"'},
    )


@app.post("/api/adapt-pptx")
async def adapt_pptx_document(
    file: UploadFile = File(...),
    profile: str = _Form("{}"),
    context: str = _Form("{}"),
    params: str = _Form("{}"),
    model: str = _Form(""),
):
    """
    Rep un PPTX, adapta el text preservant el layout i retorna el PPTX adaptat.
    Els camps profile/context/params són JSON strings (mateixos camps que /api/adapt).
    """
    import io
    from adaptation.document_adapter import (
        extract_pptx_text_map, inject_pptx_adapted, batch_adapt_text_map,
    )
    from adaptation.prompt_builder import build_system_prompt

    raw = await file.read()
    if len(raw) > MAX_UPLOAD_BYTES:
        return JSONResponse({"error": f"Fitxer massa gran ({len(raw)//1024} KB). Màxim: 5 MB."}, status_code=400)
    if not raw:
        return JSONResponse({"error": "Fitxer buit."}, status_code=400)

    filename = (file.filename or "").lower()
    if not filename.endswith(".pptx"):
        return JSONResponse({"error": "Format no suportat. Accepta: .pptx"}, status_code=400)

    try:
        profile_d = json.loads(profile)
        context_d = json.loads(context)
        params_d = json.loads(params)
    except json.JSONDecodeError as e:
        return JSONResponse({"error": f"JSON invàlid als paràmetres: {e}"}, status_code=400)

    text_map = extract_pptx_text_map(raw)
    if not text_map:
        return JSONResponse({"error": "No s'ha trobat text adaptable al PPTX."}, status_code=422)

    # Avisar de formes no accessibles (SmartArt, text dins d'imatges)
    warnings = []
    try:
        from pptx import Presentation
        import io as _io
        prs = Presentation(_io.BytesIO(raw))
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                from pptx.util import Pt
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    continue  # imatges — no es toquen
                if not shape.has_text_frame and hasattr(shape, "name"):
                    if "SmartArt" in shape.name or shape.shape_type == 14:  # 14=FREEFORM
                        warnings.append(f"Diapositiva {slide_idx+1}: «{shape.name}» (SmartArt/forma) no s'ha adaptat.")
    except Exception:
        pass

    try:
        params_d["document_mode"] = True
        system_prompt = build_system_prompt(profile_d, context_d, params_d)
        active_model = _model_for("adapt", override=model)
        adapted = batch_adapt_text_map(text_map, active_model, system_prompt)
        pptx_out = inject_pptx_adapted(raw, adapted)
    except Exception as e:
        return JSONResponse({"error": _safe_error(e, "Error durant l'adaptació del PPTX")}, status_code=500)

    filename_out = (file.filename or "presentacio").rsplit(".", 1)[0] + "_adaptat.pptx"
    headers = {"Content-Disposition": f'attachment; filename="{filename_out}"'}
    if warnings:
        headers["X-Adapt-Warnings"] = json.dumps(warnings[:10], ensure_ascii=False)
    return StreamingResponse(
        io.BytesIO(pptx_out),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=headers,
    )


# ── Generació de text base (per a docents sense text propi) ───────────────

GENERATE_EXTENSIONS = {
    "micro": (
        "Micro (50-100 paraules). Prioritza la SÍNTESI extrema. "
        "Una sola idea principal, sense exemples llargs. Pensa en un tuit o "
        "una entrada d'enciclopèdia molt breu."
    ),
    "curt": (
        "Curt (~200 paraules, marge 180-240). Concís, 2-3 idees principals "
        "lligades. Sense subapartats. Una entrada de blog breu o un paràgraf "
        "introductori d'un manual."
    ),
    "estandard": (
        "Estàndard (~400 paraules, marge 350-450). Amb desenvolupament: "
        "introducció breu, cos amb 3-4 idees connectades amb exemples concrets, "
        "i tancament. Sense subtítols obligatoris però admet 1-2 si ajuden."
    ),
    "extens": (
        "Extens (+600 paraules, marge 600-900). Desenvolupament complet amb "
        "subapartats clars (usa títols H2 o H3), detalls tècnics, exemples "
        "concrets, contextualització i tancament. Estructura editorial."
    ),
}

GENERATE_TONS = {
    "neutre": (
        "neutre i acadèmic. Vocabulari precís, frases ben construïdes, "
        "distancia objectiva. Cap referència personal del tipus 'tu' o 'jo'."
    ),
    "informal": (
        "informal i col·loquial. Tutejant directament al lector ('tu', 'sabies "
        "que...?'), amb expressions properes i exemples del dia a dia. Estàndard "
        "oral català però respectant la normativa de l'IEC."
    ),
    "creatiu": (
        "creatiu i literari. Evocador, amb imatges, metàfores, ritme i música "
        "del llenguatge. Permet recursos retòrics: anàfora, comparació, "
        "personificació. Sense perdre rigor sobre el tema."
    ),
    "motivador": (
        "motivador i engrescador. Comença amb un ganxo (pregunta retòrica, "
        "anècdota intrigant o dada sorprenent) i manté l'interès amb verbs "
        "actius i frases dinàmiques."
    ),
    "reflexiu": (
        "reflexiu. Convida a pensar amb preguntes obertes intercalades, "
        "contrastos, exploració de matisos. No dona respostes tancades."
    ),
    "empatic": (
        "empàtic i cuidadós, especialment sensible amb temes delicats "
        "(trauma, vulnerabilitat, identitat). Vocabulari respectuós, evita "
        "judicis de valor, dona espai al lector."
    ),
    "humoristic": (
        "humorístic i divertit, amb tocs lleugers, jocs de paraules, "
        "comparacions enginyoses. Sense caure en la frivolitat: el contingut "
        "ha de ser rigorós, el to amable."
    ),
    "solemne": (
        "solemne, formal i ple de respecte. Propi de textos històrics, "
        "commemoratius o de memòria. Vocabulari elevat, frases ben mesurades, "
        "absència d'humor."
    ),
}


# Bug 2 (2026-04-19): el camp `extensio` arribava al generador com a string
# descriptiva ("200-300 paraules") i no es traduïa a cap target numèric, de
# manera que el generador sempre acabava retornant ~400 paraules (default).
# Aquest mapping fa el pont entre l'etiqueta humana i el `target_words` que
# espera `generador_lliure`. Si el valor no matcheja cap patró conegut, el
# mòdul downstream usa el seu default actual (400).
_EXTENSIO_TO_TARGET_WORDS = {
    "50-100 paraules": 75,
    "100-200 paraules": 150,
    "200-300 paraules": 250,
    "300-500 paraules": 400,
    "500-800 paraules": 650,
    "800-1200 paraules": 1000,
}


def _resolve_target_words(payload: dict) -> dict:
    """Retorna una còpia del payload amb `target_words` injectat quan
    l'`extensio` coincideix amb una etiqueta coneguda. No muta l'entrada.

    Si el client ja ha passat `target_words` explícit, el respectem.
    """
    if not isinstance(payload, dict):
        return payload
    if payload.get("target_words"):
        return payload
    ext = (payload.get("extensio") or "").strip().lower()
    target = _EXTENSIO_TO_TARGET_WORDS.get(ext)
    if target is None:
        return payload  # default downstream (400) queda intacte
    return {**payload, "target_words": target}


@app.post("/api/generate-text")
async def generate_text(payload: dict = Body(...)):
    """
    Genera un text base segons context i paràmetres.
    Per a docents que no disposen del text que volen adaptar.

    Des del sprint 2026-04-15 delega al mòdul `generador_lliure` que usa
    un prompt mínim (~110 paraules) i `_call_llm_raw` sense el prefix
    "TEXT ORIGINAL A ADAPTAR". Zero pipeline de qualitat: el text del
    model va directe a la resposta. Motivat pel cas del castell medieval
    on el pipeline antic contaminava el registre. Veure pla a
    `.claude/plans/sorted-juggling-locket.md`.

    Payload:
        tema: str (required)
        genere: str (gènere discursiu, ex: "Article divulgatiu")
        tipologia: str (expositiva | narrativa | descriptiva | argumentativa | instructiva | dialogada)
        to: str (neutre | proper | formal | divulgatiu | ...)
        extensio: str (curt | estandard | extens)
        notes: str (instruccions addicionals, opcional)
        context: dict amb etapa, curs, ambit, materia (opcional)
        saber_curricular: str (opcional, Sprint C — stub curriculum KG)
        model: str (opcional, override del model admin)
    """
    from generador_lliure import generar as generar_text_lliure

    payload = _resolve_target_words(payload)
    payload["model"] = _model_for("generate", override=(payload.get("model") or "").strip())

    try:
        result = generar_text_lliure(payload)
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:
        return JSONResponse(
            {"error": _safe_error(e, "Error generant el text")},
            status_code=500,
        )

    return result


@app.post("/api/generate-text-stream")
async def generate_text_stream(request: Request, payload: dict = Body(...)):
    """Variant streaming (SSE) de `/api/generate-text`.

    Emet events Server-Sent Events a mesura que el LLM produeix tokens.
    Afegit 2026-04-16 per a la UX del pilot: amb Gemma/Qwen generant 60-90s,
    veure pantalla buida és insostenible. El cost i els tokens totals són
    idèntics a la versió no streaming; només canvia el transport.

    Events:
        data: {"type":"start","model":"...","target_words":N}
        data: {"type":"chunk","text":"..."}   (N repeticions)
        data: {"type":"done","text":"...","paraules":N,"duration_ms":M,"model":"..."}
        data: {"type":"error","message":"..."}

    Payload idèntic a /api/generate-text.
    """
    _rate_check(f"gen:{request.client.host}", 15, 60)
    from generador_lliure import generar_stream as generar_text_stream_lliure

    # Bug 2 (2026-04-19): mateix mapping que /api/generate-text perquè el
    # streaming també respecti l'extensió demanada.
    payload = _resolve_target_words(payload)
    # Desacoblament: el server resol el model i l'injecta al payload perquè
    # generador_lliure no hagi d'importar res de server.py.
    payload["model"] = _model_for("generate", override=(payload.get("model") or "").strip())

    async def gen():
        # Els chunks del LLM arriben sync; els emetem com events SSE.
        # Fem servir un ThreadPoolExecutor perquè l'iterador s\u00edncron
        # del SDK (google-genai / openai) no bloquegi l'event loop.
        loop = asyncio.get_event_loop()
        queue: asyncio.Queue = asyncio.Queue()

        def worker():
            try:
                for event in generar_text_stream_lliure(payload):
                    loop.call_soon_threadsafe(queue.put_nowait, event)
            except Exception as e:
                loop.call_soon_threadsafe(
                    queue.put_nowait,
                    {"type": "error", "message": _safe_error(e, "Error generant el text")},
                )
            finally:
                loop.call_soon_threadsafe(queue.put_nowait, None)  # sentinella fi

        loop.run_in_executor(None, worker)

        while True:
            event = await queue.get()
            if event is None:
                break
            yield f"data: {json.dumps(event, ensure_ascii=False)}\n\n"
            if event.get("type") in ("done", "error"):
                break

    return StreamingResponse(
        gen(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",  # desactiva buffering de nginx si hi \u00e9s
        },
    )


# ── Refinament de text (sense regenerar) ──────────────────────────────────

REFINE_PRESETS = {
    "catala": (
        "REGLA SUPREMA: actua com un corrector ortogràfic CONSERVADOR. "
        "NOMÉS corregeix errors que estiguis 100% segur que són errors. "
        "Si tens CAP dubte sobre si una paraula o forma és correcta, DEIXA-LA TAL COM ESTÀ. "
        "És PITJOR introduir un error nou que deixar un error existent. "
        "NO reformulis frases. NO canviïs sinònims. NO reorganitzis el text.\n\n"
        "\n"
        "CORRECCIONS CERTES A APLICAR:\n"
        "1. Paraules en altres llengües: tradueix NOMÉS les paraules clarament no catalanes "
        "(castellà, francès, anglès). Ex: 'ensemble' → 'conjunt', 'approach' → 'enfocament'. "
        "Si és un terme tècnic acceptat en català, NO el canviïs.\n"
        "2. Accents diacrítics IEC 2017: només els 15 vigents "
        "(bé, déu, és, mà, més, món, pèl, què, sé, sí, sóc, són, té, ús, vós). "
        "La resta NO porten diacrític: 'dona' (no 'dóna'), 'os' (no 'ós'), 'net' (no 'nét').\n"
        "3. Apostrofacions: l'home, l'escola, l'una, d'ahir. "
        "NO apostrofis davant h aspirada o consonant inicial pronunciada.\n"
        "4. Aglutinacions barbarismes: 'tal i com' → 'tal com', 'donar-se compte' → 'adonar-se'.\n"
        "5. Formes verbals: 'hauries' (no 'hauríes'), 'seria' (no 'séria').\n"
        "6. Concordances de gènere/nombre evidents.\n"
        "\n"
        "COSES QUE NO HAS DE TOCAR SI NO ESTÀS 100% SEGUR:\n"
        "- No canviïs preposicions regim ('pensar en' vs 'pensar a') si l'original funciona.\n"
        "- No afegeixis ni treguis 'ha de', 'cal', 'és menester'.\n"
        "- No transformis estructures passives/actives.\n"
        "- No inventis paraules ni les substitueixis per sinònims.\n"
        "- No canviïs temps verbals ni modes.\n"
        "\n"
        "VERIFICACIÓ FINAL OBLIGATÒRIA: abans de retornar el text, rellegeix-lo "
        "paraula per paraula. Si hi ha CAP paraula que no existeixi al diccionari "
        "normatiu (DIEC2), ELIMINA-LA o substitueix-la per la forma correcta. "
        "Si has introduït canvis que no tens 100% de certesa, REVERTEIX-LOS a "
        "l'original. Retorna NOMÉS text 100% normatiu."
    ),
    "simplificar": (
        "Simplifica el llenguatge sense canviar el contingut ni la longitud. "
        "Substitueix paraules complexes per sinònims més freqüents. Trenca "
        "frases llargues en frases més curtes. Mantén el to general. "
        "OBJECTIU: mateix missatge, vocabulari més senzill."
    ),
    "enriquir": (
        "Enriqueix el text UN GRAÓ respecte al seu nivell actual. "
        "PRIMER observa el text original: el seu to, la seva complexitat sintàctica "
        "i el seu destinatari implícit. DESPRÉS eleva'l lleugerament sense canviar "
        "el gènere, el to ni la longitud: "
        "substitueix paraules massa freqüents per sinònims més precisos; "
        "afegeix detalls descriptius o matisos on el text és massa esquemàtic. "
        "REGLA CLAU: respecta el nivell de partida. "
        "Si el text és per a lectors joves o emergents → enriqueix el vocabulari "
        "però mantén les frases curtes i simples, sense subordinació. "
        "Si el text és de nivell intermedi o superior → pots afegir alguna subordinada "
        "i connectors argumentatius amb mesura. "
        "OBJECTIU: una versió lleugerament millor del mateix text, no un canvi de registre."
    ),
    "ampliar": (
        "AMPLIA EL TEXT. El text resultant HA DE SER MÉS LLARG que l'original, "
        "no més curt. Afegeix exemples concrets, contextualització, detalls "
        "explicatius i matisos que enriqueixin el contingut. Mantén l'estructura "
        "general i el to. OBJECTIU OBLIGATORI: el text ampliat ha de tenir "
        "aproximadament entre un 30% i un 50% més de paraules que l'original. "
        "NO escurcis. NO retornis un text més breu. Si l'original té 200 paraules, "
        "el resultat ha de tenir com a mínim 260 paraules."
    ),
    "escurcar": (
        "ESCURÇA EL TEXT. El text resultant HA DE SER MÉS CURT que l'original, "
        "no més llarg. Elimina paraules redundants, repeticions, frases "
        "secundàries, exemples superflus i digressions. Mantén TOTES les idees "
        "principals intactes. OBJECTIU OBLIGATORI: el text escurçat ha de tenir "
        "aproximadament entre un 25% i un 40% menys de paraules que l'original. "
        "NO ampliïs. NO retornis un text més llarg. Si l'original té 200 paraules, "
        "el resultat ha de tenir entre 120 i 150 paraules."
    ),
    "to_mes_proper": (
        "Reescriu el text amb un to més proper i informal, parlant "
        "directament al lector ('tu', 'sabies que...?'). Mantén el contingut "
        "intacte."
    ),
    "to_mes_formal": (
        "Reescriu el text amb un to més formal i acadèmic. Elimina "
        "expressions col·loquials. Mantén el contingut intacte."
    ),
}


LANGUAGETOOL_URL = os.getenv("LANGUAGETOOL_URL", "https://api.languagetool.org/v2/check")

# Model LLM utilitzat per a l'auditor pedagògic (Layer 3 del pipeline de qualitat).
# Per defecte gpt-4o-mini (barat + ràpid). Canviable via env var si cal,
# però ATENCIÓ: usem OPENAI_API_KEY institucional → NO canviar a models premium
# sense aprovació explícita (gpt-4o, gpt-4.1 són 10-30x més cars).
ATNE_AUDITOR_MODEL = os.getenv("ATNE_AUDITOR_MODEL", "gpt-4o-mini")

# Auditor DESACTIVAT per defecte des del 14/04/26: proves empíriques amb
# 3 models (Gemma/GPT/Mistral) van mostrar ~15% precisió i ~85% falsos positius.
# El docent el pot reactivar via UI (toggle al Pas 2) per a casos específics.
ATNE_AUDITOR_ENABLED = os.getenv("ATNE_AUDITOR_ENABLED", "false").lower() == "true"


def _languagetool_correct(text: str, lang: str = "ca") -> tuple[str, int, list[dict]]:
    """
    Corregeix un text via LanguageTool API pública (determinista, no LLM).
    Retorna: (text_corregit, n_canvis, llista_canvis).
    Si LanguageTool falla, retorna (text, 0, []).
    """
    from adaptation.lang_config import get_lt_code
    try:
        resp = requests.post(
            LANGUAGETOOL_URL,
            data={
                "text": text,
                "language": get_lt_code(lang),
                "level": "picky",
                "enabledOnly": "false",
            },
            headers={"Accept": "application/json"},
            timeout=30,
        )
        if resp.status_code != 200:
            return text, 0, []
        data = resp.json()
    except Exception as e:
        print(f"[LanguageTool] Error: {type(e).__name__}: {e}")
        return text, 0, []

    matches = data.get("matches", [])
    if not matches:
        return text, 0, []

    # Aplicar correccions d'esquerra a dreta inverses (per no desplaçar offsets)
    sorted_matches = sorted(matches, key=lambda m: m.get("offset", 0), reverse=True)

    corrected = text
    changes: list[dict] = []
    for m in sorted_matches:
        offset = m.get("offset", 0)
        length = m.get("length", 0)
        replacements = m.get("replacements", [])
        if not replacements or length <= 0:
            continue
        new_value = replacements[0].get("value", "")
        if new_value is None:
            continue
        # Offsets LT en UTF-16 code units: cal convertir per extreure/substituir
        _b_cur = corrected.encode('utf-16-le')
        _b_start = offset * 2
        _b_end = (offset + length) * 2
        old_value = _b_cur[_b_start:_b_end].decode('utf-16-le', errors='replace')
        # Saltar si la substitució és idèntica
        if new_value == old_value:
            continue
        # Protecció markdown: no tocar backticks (trenca blocs de codi)
        if '`' in old_value:
            continue
        corrected = _lt_splice(corrected, offset, length, new_value)
        changes.append({
            "original": old_value,
            "corregit": new_value,
            "missatge": m.get("shortMessage") or m.get("message", ""),
            "regla": m.get("rule", {}).get("id", ""),
        })

    # Retornar canvis en ordre d'aparició al text
    changes.reverse()
    return corrected, len(changes), changes


# ═══ Pipeline de qualitat català (post-processament) ══════════════════════════

# Taula etapa+curs → MECR aproximat per al target de llegibilitat
# Usada quan el client no envia target_mecr explícit.
def _mecr_from_etapa_curs(etapa: str, curs: str = "") -> str:
    _MAP = {
        "infantil":    {"P3": "pre-A1", "P4": "pre-A1", "P5": "pre-A1"},
        "primaria":    {"1r": "A1", "2n": "A1", "3r": "A1", "4t": "A2", "5e": "A2", "6e": "B1"},
        "ESO":         {"1r": "B1", "2n": "B1", "3r": "B2", "4t": "B2"},
        "batxillerat": {"1r": "B2", "2n": "C1"},
        "FP":          {"1r_CFGB": "A2", "2n_CFGB": "A2", "1r_CGM": "B1", "2n_CGM": "B1", "1r_CGS": "B2", "2n_CGS": "B2"},
    }
    _FALLBACK = {"infantil": "pre-A1", "primaria": "A1", "ESO": "B1", "batxillerat": "B2", "FP": "B1"}
    return _MAP.get(etapa, {}).get(curs) or _FALLBACK.get(etapa, "B1")


# ═══ Filtre de caràcters exòtics (CJK, ciríl·lic, àrab, etc.) ══════════════
# Rangs Unicode permesos per a text català. Tota la resta genera alerta.
_ALLOWED_UNICODE_RANGES = [
    (0x0009, 0x000A),   # Tab, LF
    (0x000D, 0x000D),   # CR
    (0x0020, 0x007E),   # ASCII printable
    (0x00A0, 0x024F),   # Latin-1 Supplement + Latin Extended-A/B (à, é, ç, ñ, etc.)
    (0x2000, 0x206F),   # General Punctuation (em dash, smart quotes, bullet, ellipsis)
    (0x20A0, 0x20CF),   # Currency symbols (€, £, ¥)
    (0x2100, 0x214F),   # Letterlike (™, ℃, ℉)
    (0x2190, 0x21FF),   # Arrows (↑ ↓ → ←) — per esquemes
    (0x2200, 0x22FF),   # Mathematical Operators (+, ±, ×)
    (0x2500, 0x257F),   # Box drawing (├ ─ └) — per esquemes
    (0x25A0, 0x25FF),   # Geometric shapes (● ■ ▲)
    (0x2600, 0x27BF),   # Miscellaneous Symbols (☀ ♪ ⚡)
    (0x1F300, 0x1F9FF), # Emojis (alguns complements demanen pictogrames)
    (0x1FA70, 0x1FAFF), # Emojis v13+
]


def _is_char_allowed(ch: str) -> bool:
    cp = ord(ch)
    for start, end in _ALLOWED_UNICODE_RANGES:
        if start <= cp <= end:
            return True
    return False


def _detect_script(cp: int) -> str:
    """Identifica aproximadament l'script d'un codepoint."""
    if 0x0400 <= cp <= 0x04FF or 0x0500 <= cp <= 0x052F:
        return "ciríl·lic"
    if 0x0590 <= cp <= 0x05FF:
        return "hebreu"
    if 0x0600 <= cp <= 0x06FF or 0x0750 <= cp <= 0x077F or 0x08A0 <= cp <= 0x08FF:
        return "àrab"
    if 0x0900 <= cp <= 0x097F:
        return "devanagari (hindi)"
    if 0x0E00 <= cp <= 0x0E7F:
        return "tailandès"
    if 0x3000 <= cp <= 0x303F:
        return "CJK puntuació"
    if 0x3040 <= cp <= 0x309F:
        return "hiragana (japonès)"
    if 0x30A0 <= cp <= 0x30FF:
        return "katakana (japonès)"
    if 0x4E00 <= cp <= 0x9FFF:
        return "xinès"
    if 0xAC00 <= cp <= 0xD7AF:
        return "coreà (hangul)"
    if 0xE000 <= cp <= 0xF8FF:
        return "ús privat"
    if 0xFB00 <= cp <= 0xFB4F:
        return "presentacions alfabètiques"
    return "desconegut"


def _exotic_char_scan(text: str) -> list[dict]:
    """
    Escaneja el text en cerca de caràcters no esperats per a text català
    (CJK, ciríl·lic, àrab, hangul, ús privat, etc.). Retorna una llista
    amb un màxim de 10 entrades úniques amb {caracter, codepoint, script,
    context, ocurrencies}.

    Motivació: els LLMs (vam veure Gemma amb '홈olatge') ocasionalment
    insereixen caràcters d'altres scripts per artefactes de tokenització.
    Aquest filtre ho detecta de manera determinista abans que cap altre
    pas del pipeline.
    """
    if not text:
        return []

    grouped: dict[str, dict] = {}
    for i, ch in enumerate(text):
        if ch in (" ", "\n", "\t", "\r"):
            continue
        if _is_char_allowed(ch):
            continue
        cp = ord(ch)
        key = ch
        if key not in grouped:
            context_start = max(0, i - 18)
            context_end = min(len(text), i + 18)
            context_raw = text[context_start:context_end]
            context = context_raw.replace(ch, f"«{ch}»", 1)
            grouped[key] = {
                "caracter": ch,
                "codepoint": f"U+{cp:04X}",
                "script": _detect_script(cp),
                "context": context,
                "ocurrencies": 1,
            }
        else:
            grouped[key]["ocurrencies"] += 1
        if len(grouped) >= 10:
            break

    return list(grouped.values())


# Regles de LanguageTool que indiquen paraula desconeguda al diccionari català
_LT_UNKNOWN_WORD_RULES = {
    "MORFOLOGIK_RULE_CA_ES",  # Catalan ortografia general
    "MORFOLOGIK_RULE_CA_ES_V",
    "MORFOLOGIK_RULE_CA_ES_VALENCIA",
}

# Mots catalans legítims que el diccionari de LanguageTool marca erròniament.
# Comparació case-insensitive (lower) sobre el fragment exacte del match.
# Afegir només mots verificats; cada entrada evita un fals positiu al Quality Report.
_LOCAL_WORD_WHITELIST = frozenset({
    "rails",   # via fèrria, mot català legítim (glossari FJE 15/04)
    "raïls",   # variant amb dièresi, també correcte
})

# Llindars de llegibilitat per MECR (mitjana paraules/frase i % paraules llargues >7 caràcters)
_READABILITY_TARGETS = {
    "pre-A1": {"max_wps": 8, "max_long_pct": 12},
    "A1": {"max_wps": 10, "max_long_pct": 15},
    "A2": {"max_wps": 14, "max_long_pct": 20},
    "B1": {"max_wps": 18, "max_long_pct": 28},
    "B2": {"max_wps": 25, "max_long_pct": 35},
    "C1": {"max_wps": 30, "max_long_pct": 42},
}


def _lt_splice(text: str, offset_u16: int, length_u16: int, new_value: str) -> str:
    """Splice a LT match respectant UTF-16 code units.

    LanguageTool retorna offsets i lengths en codeunits UTF-16 (com JS).
    Python str opera en codepoints. Els emojis astral (🌍, 💧, 🐰…) ocupen
    2 codeunits UTF-16 però 1 codepoint Python — sense aquesta conversió,
    les correccions cauen desfasades dins d'altres paraules.
    """
    b = text.encode('utf-16-le')
    start = offset_u16 * 2
    end = (offset_u16 + length_u16) * 2
    prefix = b[:start].decode('utf-16-le', errors='replace')
    suffix = b[end:].decode('utf-16-le', errors='replace')
    return prefix + new_value + suffix


def _languagetool_full_analysis(text: str, lang: str = "ca") -> dict:
    """
    Crida LanguageTool i separa matches en 3 categories:
    - correccions automàtiques (ortografia, gramàtica amb suggeriment clar)
    - paraules desconegudes (warnings, no auto-corregides)
    - avisos d'estil (info, no crítics)
    Retorna un dict amb text corregit i les 3 llistes.
    """
    from adaptation.lang_config import get_lt_code
    result = {
        "text_original": text,
        "text_corregit": text,
        "correccions": [],
        "paraules_desconegudes": [],
        "avisos_estil": [],
        "lt_disponible": False,
    }

    try:
        resp = requests.post(
            LANGUAGETOOL_URL,
            data={
                "text": text,
                "language": get_lt_code(lang),
                "level": "picky",
                "enabledOnly": "false",
            },
            headers={"Accept": "application/json"},
            timeout=30,
        )
        if resp.status_code != 200:
            return result
        data = resp.json()
        result["lt_disponible"] = True
    except Exception as e:
        print(f"[LanguageTool] Error: {type(e).__name__}: {e}")
        return result

    matches = data.get("matches", [])
    if not matches:
        return result

    sorted_matches = sorted(matches, key=lambda m: m.get("offset", 0), reverse=True)

    corrected = text
    correccions: list[dict] = []
    desconegudes: list[dict] = []
    avisos: list[dict] = []

    for m in sorted_matches:
        offset = m.get("offset", 0)
        length = m.get("length", 0)
        replacements = m.get("replacements", [])
        rule_id = m.get("rule", {}).get("id", "")
        rule_cat = m.get("rule", {}).get("category", {}).get("id", "")
        # Offsets de LT són UTF-16 code units; cal convertir per extreure old_value
        if offset >= 0:
            _b_orig = text.encode('utf-16-le')
            _b_start = offset * 2
            _b_end = (offset + length) * 2
            old_value = _b_orig[_b_start:_b_end].decode('utf-16-le', errors='replace')
        else:
            old_value = ""
        missatge = m.get("shortMessage") or m.get("message", "")

        # Whitelist local: mots catalans legítims que LT marca erròniament.
        # Es comprova abans de qualsevol classificació — si hi és, s'ignora
        # el match sencer (ni correcció, ni avís, ni paraula sospitosa).
        if old_value and old_value.strip(".,;:!?'\"()[]{}").lower() in _LOCAL_WORD_WHITELIST:
            continue

        # Protecció markdown: saltar qualsevol match que abasti backticks.
        # LT 'picky' transforma ` aïllats en apòstrofs tipogràfics, cosa que
        # trenca els blocs de codi (```) i els codis inline (`x`).
        if '`' in old_value:
            continue

        # Estil explícit → warning, MAI auto-aplicar
        if rule_cat in _STYLE_WARNING_CATEGORIES:
            avisos.append({
                "fragment": old_value,
                "suggeriment": replacements[0].get("value", "") if replacements else "",
                "missatge": missatge,
                "regla": rule_id,
            })
            continue

        # Si és paraula desconeguda (MORFOLOGIK) → warning, i mirem si és
        # segur auto-aplicar (només si norm igual: accents, majúscules, etc.)
        if rule_id in _LT_UNKNOWN_WORD_RULES or "MORFOLOGIK" in rule_id:
            desconegudes.append({
                "paraula": old_value,
                "suggeriments": [r.get("value", "") for r in replacements[:3]],
                "missatge": missatge,
                "regla": rule_id,
            })
            if replacements and _suggestion_is_safe(
                old_value, replacements[0].get("value", ""), rule_id, rule_cat
            ):
                new_value = replacements[0].get("value", "")
                if new_value and new_value != old_value:
                    corrected = _lt_splice(corrected, offset, length, new_value)
                    correccions.append({
                        "original": old_value,
                        "corregit": new_value,
                        "missatge": missatge,
                        "regla": rule_id,
                    })
            continue

        # Resta → auto-aplicar NOMÉS si és segur (categoria/prefix o norm)
        if replacements and length > 0:
            new_value = replacements[0].get("value", "")
            if new_value and new_value != old_value and _suggestion_is_safe(
                old_value, new_value, rule_id, rule_cat
            ):
                corrected = _lt_splice(corrected, offset, length, new_value)
                correccions.append({
                    "original": old_value,
                    "corregit": new_value,
                    "missatge": missatge,
                    "regla": rule_id,
                })
            elif replacements:
                # No segur → emetre com a avís (l'usuari decideix)
                avisos.append({
                    "fragment": old_value,
                    "suggeriment": new_value,
                    "missatge": missatge,
                    "regla": rule_id,
                })

    # Revertir ordres per lectura natural
    correccions.reverse()
    desconegudes.reverse()
    avisos.reverse()

    result["text_corregit"] = corrected
    result["correccions"] = correccions
    result["paraules_desconegudes"] = desconegudes
    result["avisos_estil"] = avisos
    return result


## Categories de rules de LanguageTool que considerem SEGURES per auto-aplicació
## (no canvien el lema ni la semàntica — només ortografia / apostrofació / etc.)
_SAFE_RULE_CATEGORIES = {
    "PUNCTUATION",
    "TYPOGRAPHY",
    # CASING tret 2026-04-21 (Miquel): LT baixa "Revolucio Industrial" a
    # minuscules per norma IEC estricta, xoca amb us pedagogic de tractar-ho
    # com a nom propi d'epoca. Ara passa a warning, el docent decideix.
    "HYPHENATION",
    "WHITESPACE",
    "DIACRITICS_CA",
    "DIACRITICS",
}

## Prefixos de rule_id que són segurs (IEC-determinístics)
_SAFE_RULE_PREFIXES = (
    "L_APOSTROF",              # L' apostrofació catalana (el home → l'home)
    "APOSTROPHE",
    "APOSTROFAT",
    "APOSTROF",
    "DIACRITICS",              # accents diacrítics
    "ACCENTUATION",
    "HIAT",                    # hiatus
    "WHITESPACE",
    "DOUBLE_PUNCTUATION",
    # UPPERCASE_SENTENCE_START i UPPERCASE_ trets 2026-04-21: passen a warning
    # (veure nota al _SAFE_RULE_CATEGORIES sobre CASING).
    "NUMBER_SPACE",
    "PUNT_FINAL",              # punt final
    "A_EL",                    # a el → al
    "DE_EL",                   # de el → del
    "PER_A_EL",                # per a el → pel
    "GUIONET",                 # guionet
    "HYPHEN",
    "COMMA_",                  # comes missing o extra
    "COMMA_PARENTHESIS",
)

## Categories que NO auto-apliquem (warnings d'estil, revisió humana recomanada)
_STYLE_WARNING_CATEGORIES = {
    "STYLE",
    "REDUNDANCY",
    "REGISTER",
    "COLLOQUIAL",
}


def _suggestion_is_safe(original: str, suggestion: str,
                        rule_id: str = "", rule_category: str = "") -> bool:
    """
    Un suggeriment es considera "segur" per a auto-aplicació si es compleix
    alguna d'aquestes condicions:
      1. La rule_category o el rule_id coincideix amb les llistes de rules
         segures (apostrofacions, accents, majúscules, puntuació, etc.)
      2. La normalització alfanumèrica (sense accents ni caràcters no alfanums)
         és idèntica entre original i suggeriment — això cobreix variacions
         d'accent, apostrof, majúscules i puntuació sense canviar el lema.

    Exemples de canvis AUTOAPLICABLES:
      - "dona" → "dóna"                (accent)
      - "l home" → "l'home"            (apostrof)
      - "el home" → "l'home"           (apostrof via rule L_APOSTROF)
      - "tambe" → "també"              (accent)
      - "Pompeu fabra" → "Pompeu Fabra" (majúscula)
      - "pel·lícula " → "pel·lícula"   (whitespace)

    Exemples de canvis NO autoaplicables (warnings):
      - "tal i com" → "tal com"        (treu una paraula — canvia lema)
      - "ser vius" → "éssers vius"     (canvia lema)
      - "moltíssim" → "molt"           (sinonimia, revisió humana)
    """
    if not original or not suggestion:
        return False

    rid = (rule_id or "").upper()
    rcat = (rule_category or "").upper()

    # 1. Rule category segura
    if rcat in _SAFE_RULE_CATEGORIES:
        return True

    # 2. Rule ID amb prefix segur
    if any(rid.startswith(p) for p in _SAFE_RULE_PREFIXES):
        return True

    # 3. Categoria explícita de warning d'estil → mai segur
    if rcat in _STYLE_WARNING_CATEGORIES:
        return False

    # 4. Per defecte: no és segur → avís visible, no auto-apply.
    #    Abans hi havia un fallback "normalització alfanumèrica iguals →
    #    segur" que era massa laxista: auto-aplicava qualsevol canvi
    #    d'accent, incloent suggeriments dubtosos de 'picky' com
    #    "diferència" (nom) → "diferencia" (verb conjugat). Ara només
    #    s'auto-aplica si hi ha regla/categoria explícitament segura (#1/#2).
    return False


def _readability_score(text: str, target_mecr: str = "") -> dict:
    """
    Calcula indicadors de llegibilitat del text en català:
    - paraules per frase (mitjana)
    - percentatge de paraules llargues (>7 caràcters)
    - longitud mitjana de paraula
    Compara amb els llindars del MECR objectiu i retorna un estat:
    "ok" / "sobre" (text més complex del que caldria) / "sota" (massa simple).
    """
    import re

    text_clean = text.strip()
    if not text_clean:
        return {"ok": True, "wps": 0, "long_pct": 0, "avg_word_len": 0,
                "target_mecr": target_mecr, "estat": "buit",
                "missatge": "Text buit"}

    # Trencar en frases per punts i signes d'interrogació/exclamació finals
    sentences = [s.strip() for s in re.split(r"[.!?]+", text_clean) if s.strip()]
    n_sentences = max(1, len(sentences))

    words = re.findall(r"\b[\wàèéíòóú·ïü]+\b", text_clean, flags=re.IGNORECASE)
    n_words = max(1, len(words))

    wps = round(n_words / n_sentences, 1)
    long_words = [w for w in words if len(w) > 7]
    long_pct = round(100 * len(long_words) / n_words, 1)
    avg_word_len = round(sum(len(w) for w in words) / n_words, 1)

    result = {
        "wps": wps,
        "long_pct": long_pct,
        "avg_word_len": avg_word_len,
        "n_words": n_words,
        "n_sentences": n_sentences,
        "target_mecr": target_mecr,
    }

    target = _READABILITY_TARGETS.get(target_mecr)
    if not target:
        result["ok"] = True
        result["estat"] = "sense_objectiu"
        result["missatge"] = f"Llegibilitat: {wps} par/frase · {long_pct}% paraules llargues"
        return result

    sobrepassa_wps = wps > target["max_wps"]
    sobrepassa_long = long_pct > target["max_long_pct"]

    if sobrepassa_wps and sobrepassa_long:
        result["ok"] = False
        result["estat"] = "sobre"
        result["missatge"] = (
            f"El text és més complex del que caldria per a {target_mecr}: "
            f"{wps} par/frase (màx {target['max_wps']}) i "
            f"{long_pct}% paraules llargues (màx {target['max_long_pct']}%)."
        )
    elif sobrepassa_wps:
        result["ok"] = False
        result["estat"] = "sobre"
        result["missatge"] = (
            f"Frases massa llargues per a {target_mecr}: "
            f"{wps} par/frase (màx recomanat {target['max_wps']})."
        )
    elif sobrepassa_long:
        result["ok"] = False
        result["estat"] = "sobre"
        result["missatge"] = (
            f"Vocabulari potser massa complex per a {target_mecr}: "
            f"{long_pct}% paraules llargues (màx recomanat {target['max_long_pct']}%)."
        )
    else:
        result["ok"] = True
        result["estat"] = "ok"
        result["missatge"] = (
            f"Llegibilitat adequada per a {target_mecr}: "
            f"{wps} par/frase · {long_pct}% paraules llargues."
        )

    return result


def _llm_audit(text: str, target_mecr: str = "", etapa: str = "") -> dict:
    """
    Layer 3 del pipeline: auditor pedagògic via LLM (GPT-4o-mini per defecte).
    NO modifica el text — només emet avisos qualitatius que LanguageTool no veu:
    frases confuses, salts lògics, vocabulari desajustat, repeticions, connectors
    mal usats, calcs d'altres llengües.

    Retorna: {"avisos": [...], "disponible": bool, "model": str, "error": str?}
    - avisos és una llista de dicts {"tipus", "fragment", "motiu"}
    - disponible indica si s'ha pogut fer la crida (fallback silent a LT si no)
    """
    result = {"avisos": [], "disponible": False, "model": ATNE_AUDITOR_MODEL}

    if not text or not text.strip():
        result["error"] = "text buit"
        return result

    if not os.getenv("OPENAI_API_KEY"):
        result["error"] = "OPENAI_API_KEY no configurada"
        return result

    contexte_etapa = f"alumne de {etapa}" if etapa else "un alumne d'escola catalana"
    contexte_mecr = f"nivell MECR {target_mecr}" if target_mecr else "nivell estàndard"

    prompt = f"""Ets un inspector pedagògic en català. Has d'auditar el text següent pensant en un {contexte_etapa} amb {contexte_mecr}.

NO modifiquis el text. NOMÉS emet fins a 6 avisos qualitatius sobre:
- Frases ambigües o confuses
- Salts lògics o idees que perden el fil
- Vocabulari massa complex o inadequat per al nivell
- Repeticions que no aporten
- Connectors mal usats (connectors causals, consecutius, adversatius usats incorrectament)
- Construccions no naturals en català (calcs del castellà, francès o anglès)

Retorna NOMÉS un objecte JSON vàlid (sense marcadors markdown, sense explicacions extra), amb aquesta estructura exacta:

{{"avisos": [{{"tipus": "confusa|salt|vocabulari|repeticio|connector|calc", "fragment": "fragment literal del text (max 100 car)", "motiu": "explicació breu del problema"}}]}}

Si el text és adequat al nivell i no hi ha problemes rellevants, retorna {{"avisos": []}}.

TEXT A AUDITAR:
{text}"""

    try:
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        resp = client.chat.completions.create(
            model=ATNE_AUDITOR_MODEL,
            messages=[
                {"role": "system", "content": "Ets un auditor pedagògic que retorna NOMÉS JSON vàlid, sense marcadors markdown ni text addicional."},
                {"role": "user", "content": prompt},
            ],
            max_tokens=1200,
            temperature=0.2,
            timeout=20,
        )
        raw = (resp.choices[0].message.content or "").strip()
    except Exception as e:
        print(f"[LLM Auditor] Error: {type(e).__name__}: {e}")
        result["error"] = f"{type(e).__name__}: {e}"
        return result

    # Netejar markdown fences si existeixen
    if raw.startswith("```"):
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)

    try:
        data = json.loads(raw)
        avisos_raw = data.get("avisos", [])
        if not isinstance(avisos_raw, list):
            avisos_raw = []
        clean_avisos: list[dict] = []
        for a in avisos_raw[:10]:
            if isinstance(a, dict) and a.get("fragment"):
                clean_avisos.append({
                    "tipus": str(a.get("tipus", "")).strip()[:30],
                    "fragment": str(a.get("fragment", "")).strip()[:300],
                    "motiu": str(a.get("motiu", "")).strip()[:300],
                })
        result["avisos"] = clean_avisos
        result["disponible"] = True
    except Exception as e:
        print(f"[LLM Auditor] JSON parse error: {e}. Raw: {raw[:200]}")
        result["error"] = f"JSON invàlid: {e}"

    return result


def post_process_catalan(text: str, target_mecr: str = "", enable_lt: bool = True,
                         enable_auditor: bool = None, etapa: str = "",
                         lang: str = "ca") -> dict:
    """
    Pipeline complet de qualitat català per a un text generat o adaptat.

    Capes:
      1. LanguageTool (determinista): ortografia, gramàtica, paraules desconegudes
      2. Llegibilitat (heurística): paraules/frase + % paraules llargues vs MECR
      3. LLM Auditor (opcional): avisos pedagògics qualitatius (GPT-4o-mini)

    Les capes 1 i 3 s'executen EN PARAL·LEL amb ThreadPoolExecutor.

    Retorna un dict amb:
      - text: text final (amb correccions auto-aplicades si enable_lt)
      - n_correccions, correccions: LanguageTool
      - paraules_sospitoses: paraules no trobades al diccionari
      - avisos_estil: avisos estilístics de LanguageTool (no crítics)
      - llegibilitat: indicadors + estat vs MECR objectiu
      - avisos_auditor: warnings qualitatius del LLM auditor
      - lt_disponible, auditor_disponible, auditor_model
    """
    if enable_auditor is None:
        enable_auditor = _AUDITOR_ENABLED_RUNTIME

    if not text or not text.strip():
        return {
            "text": text,
            "n_correccions": 0,
            "correccions": [],
            "paraules_sospitoses": [],
            "avisos_estil": [],
            "llegibilitat": _readability_score("", target_mecr),
            "avisos_auditor": [],
            "caracters_exotics": [],
            "lt_disponible": False,
            "auditor_disponible": False,
            "auditor_model": ATNE_AUDITOR_MODEL,
        }

    # Capa 0 (determinista, ràpid): detectar caràcters exòtics abans de res.
    #    Capta glitches com '홈olatge' (coreà) o lletres ciríl·liques inserides
    #    per errors de tokenització.
    caracters_exotics = _exotic_char_scan(text)

    # Paral·lelitzar LanguageTool + LLM Auditor per estalviar temps
    lt_result = None
    audit_result = {"avisos": [], "disponible": False, "model": ATNE_AUDITOR_MODEL}

    with concurrent.futures.ThreadPoolExecutor(max_workers=2) as pool:
        lt_future = pool.submit(_languagetool_full_analysis, text, lang) if enable_lt else None
        audit_future = pool.submit(_llm_audit, text, target_mecr, etapa) if enable_auditor else None

        if lt_future is not None:
            try:
                lt_result = lt_future.result(timeout=45)
            except Exception as e:
                print(f"[LT] Timeout/Error: {e}")
                lt_result = None
        if audit_future is not None:
            try:
                audit_result = audit_future.result(timeout=30)
            except Exception as e:
                print(f"[Auditor] Timeout/Error: {e}")
                audit_result = {"avisos": [], "disponible": False, "model": ATNE_AUDITOR_MODEL, "error": str(e)}

    if lt_result:
        text_final = lt_result["text_corregit"]
        correccions = lt_result["correccions"]
        desconegudes = lt_result["paraules_desconegudes"]
        avisos_estil = lt_result["avisos_estil"]
        lt_disponible = lt_result["lt_disponible"]
    else:
        text_final = text
        correccions = []
        desconegudes = []
        avisos_estil = []
        lt_disponible = False

    llegibilitat = _readability_score(text_final, target_mecr)

    return {
        "text": text_final,
        "n_correccions": len(correccions),
        "correccions": correccions,
        "paraules_sospitoses": desconegudes,
        "avisos_estil": avisos_estil,
        "llegibilitat": llegibilitat,
        "avisos_auditor": audit_result.get("avisos", []),
        "caracters_exotics": caracters_exotics,
        "lt_disponible": lt_disponible,
        "auditor_disponible": audit_result.get("disponible", False),
        "auditor_model": audit_result.get("model", ATNE_AUDITOR_MODEL),
    }


@app.post("/api/refine-text")
async def refine_text(payload: dict = Body(...)):
    """
    Refina un text existent segons una instrucció.
    No regenera des de zero — modifica el text actual.

    Payload:
        text: str (required) — el text a refinar
        preset: str (opcional) — clau de REFINE_PRESETS per instruccions ràpides
        instruccio: str (opcional) — instrucció lliure del docent
    """
    text = (payload.get("text") or "").strip()
    if not text:
        return JSONResponse({"error": "Cal proporcionar el text a refinar."}, status_code=400)

    preset = (payload.get("preset") or "").strip().lower()
    instruccio_lliure = (payload.get("instruccio") or "").strip()
    model_override = (payload.get("model") or "").strip()
    mecr = (payload.get("mecr") or "").strip()

    # Preset "catala" → LanguageTool (determinista, no LLM)
    if preset == "catala" and not instruccio_lliure:
        corrected, n_canvis, canvis = _languagetool_correct(text)
        return {
            "text": corrected,
            "paraules": len(corrected.split()),
            "preset_aplicat": "catala",
            "mode": "languagetool",
            "n_canvis": n_canvis,
            "canvis": canvis,
        }

    instruccio_final = ""
    if preset == "enriquir":
        # Instrucció dinàmica: usa el nivell MECR si és conegut, sinó auto-avaluació.
        instruccio_final = corpus_reader.get_enriquir_instruction(mecr or None)
    elif preset and preset in REFINE_PRESETS:
        instruccio_final = REFINE_PRESETS[preset]
    if instruccio_lliure:
        if instruccio_final:
            instruccio_final += "\n\nA més, el docent demana específicament: " + instruccio_lliure
        else:
            instruccio_final = instruccio_lliure

    if not instruccio_final:
        return JSONResponse(
            {"error": "Cal especificar un preset o una instrucció lliure."},
            status_code=400,
        )

    # Per ampliar/escurcar, injectar objectius numèrics explícits basats en
    # el recompte real de paraules per evitar que Gemma 4 faci el contrari.
    paraules_in = len(text.split())
    if preset == "ampliar":
        target_min = int(paraules_in * 1.30)
        target_max = int(paraules_in * 1.50)
        instruccio_final += (
            f"\n\nRECOMPTE OBLIGATORI: l'original té {paraules_in} paraules. "
            f"El text resultant HA DE TENIR entre {target_min} i {target_max} paraules. "
            f"Si retornes menys de {target_min} paraules, és un ERROR. "
            f"Compta les paraules del resultat abans de retornar-lo."
        )
    elif preset == "escurcar":
        target_min = int(paraules_in * 0.60)
        target_max = int(paraules_in * 0.75)
        instruccio_final += (
            f"\n\nRECOMPTE OBLIGATORI: l'original té {paraules_in} paraules. "
            f"El text resultant HA DE TENIR entre {target_min} i {target_max} paraules. "
            f"Si retornes més de {target_max} paraules, és un ERROR. "
            f"Compta les paraules del resultat abans de retornar-lo."
        )

    prompt = f"""# ROL
Ets un expert en lingüística catalana. Has de REFINAR el text que
et passen segons les instruccions del docent. NO el regeneris des de
zero — modifica'l mantenint l'estructura general i el contingut.

# REGLES
1. Tot el text final ha de ser en català estàndard normatiu (IEC).
2. NO afegeixis explicacions meta del tipus "Aquí tens el text refinat:".
   Retorna NOMÉS el text corregit, directament.
3. NO canviïs el contingut substantiu del text. Només aplica les
   instruccions de refinament demanades.
4. Si detectes paraules en altres llengües (francès, castellà, anglès),
   tradueix-les al català correcte sempre.
5. Vigila errors típics: "ser vius" → "èssers vius", "ensemble" → "conjunt",
   apostrofacions, concordances de gènere i nombre.

# INSTRUCCIONS DEL DOCENT
{instruccio_final}

# TEXT A REFINAR
{text}

# RETORNA EL TEXT REFINAT (NOMÉS EL TEXT, RES MÉS)"""

    # Sprint 1B: el model del refine es resol via _MODEL_CONFIG["refine"]
    # (configurable des de /admin) amb override puntual via payload.model.
    refine_model = _model_for("refine", override=model_override)
    # Fallback: si el model principal falla (quota, timeout...), prova gpt-4o-mini.
    fallback_model = "gpt-4o-mini" if refine_model != "gpt-4o-mini" else "gpt-4.1-mini"
    result = ""
    model_used = refine_model
    for attempt_model in (refine_model, fallback_model):
        try:
            result = _call_llm(attempt_model, prompt, "")
            result = clean_gemini_output(result).strip()
            result = _post_process_llm_output(result)
            model_used = attempt_model
            break
        except Exception as e:
            print(f"[ATNE] refine-text fallada amb {attempt_model}: {e}", flush=True)
            if attempt_model == fallback_model:
                return JSONResponse(
                    {"error": _safe_error(e, "Error refinant el text. Torna-ho a provar en uns instants.")},
                    status_code=500,
                )

    if not result:
        return JSONResponse({"error": "L'LLM ha retornat un text buit."}, status_code=500)

    return {
        "text": result,
        "paraules": len(result.split()),
        "preset_aplicat": preset or None,
        "model_used": model_used,
    }


# ── Adaptació (SSE stream) ─────────────────────────────────────────────────

# ── Multinivell: desplaçar MECR ±1 per a mode grup ─────────────────────────

_MECR_SCALE = ["pre-A1", "A1", "A2", "B1", "B2"]

def _shift_mecr(mecr: str, shift: int) -> str:
    """Desplaça el MECR N graons amb límits pre-A1..B2."""
    try:
        idx = _MECR_SCALE.index(mecr)
    except ValueError:
        return mecr
    new_idx = max(0, min(len(_MECR_SCALE) - 1, idx + shift))
    return _MECR_SCALE[new_idx]

# Mapa de nivells per a mode grup: label → desplaçament relatiu al MECR base
LEVEL_SHIFTS = {
    "accessible": -1,
    "estandard": 0,
    "exigent": +1,
}


# ── Resolució canònica de paràmetres (Fase B, 2026-05-15) ─────────────────
#
# Substitueix progressivament les 6 implementacions paral·leles de càlcul
# de MECR/DUA que viuen al frontend i al backend. Vegeu
# adaptation/params_resolver.py per a la lògica.
@app.post("/api/derive-params")
async def derive_params_endpoint(payload: dict = Body(...)):
    """
    Calcula MECR + DUA canònic per a un perfil donat.

    Payload:
      caracteristiques: dict (mateixa estructura que la del backend);
                        també s'accepta 'profile' amb 'caracteristiques' a dins.
      etapa: str ("infantil" | "primaria" | "ESO" | "batxillerat" | "FP" | "")
      curs: str ("I5" | "1r ESO" | ...) — codi canònic; opcional
      override_mecr: str | null — si el docent ha triat un MECR manual

    Retorna:
      {ok, mecr, dua, motiu, trace}
    """
    from adaptation.params_resolver import resolve_params, MECR_ORDER

    chars = payload.get("caracteristiques")
    if chars is None:
        # Format alternatiu: profile pot venir embolcallant les caracteristiques
        profile = payload.get("profile") or {}
        chars = profile.get("caracteristiques") or {}
    if not isinstance(chars, dict):
        return JSONResponse(
            {"ok": False, "error": "'caracteristiques' ha de ser un dict"},
            status_code=400,
        )

    etapa = (payload.get("etapa") or "").strip()
    curs = (payload.get("curs") or "").strip()
    override_mecr = payload.get("override_mecr")
    if override_mecr and override_mecr not in MECR_ORDER:
        override_mecr = None  # ignorem overrides invàlids silenciosament

    try:
        result = resolve_params(chars, etapa=etapa, curs=curs, override_mecr=override_mecr)
    except Exception as e:
        return JSONResponse(
            {"ok": False, "error": f"resolve_params failed: {type(e).__name__}: {e}"},
            status_code=500,
        )

    return {"ok": True, **result}


@app.post("/api/adapt")
async def adapt_stream(request: Request, payload: dict = Body(...)):
    _rate_check(f"adapt:{request.client.host}", 15, 60)
    text = payload.get("text", "")
    profile = payload.get("profile", {})
    context = payload.get("context", {})
    params = payload.get("params", {})
    model = payload.get("model", "")  # mistral | gemma4 | (buit = default ATNE_MODEL)
    docent_id = payload.get("docent_id", "")

    if not text.strip():
        return JSONResponse({"error": "Cal proporcionar un text"}, status_code=400)

    # Nivells a generar. Per defecte: una sola versió.
    # Si arriba 'levels' amb més d'un element, generem cada un en paral·lel
    # amb MECR ajustat (accessible=-1, estandard=0, exigent=+1).
    levels = params.get("levels") or ["single"]
    base_mecr = params.get("mecr_sortida", "B1")

    async def gen():
        events: list[dict] = []
        done_count = {"n": 0}

        def make_cb(level_id: str):
            def _cb(ev):
                # Afegim la identificació del nivell a cada event, perquè el
                # frontend pugui enrutar-lo a la pestanya corresponent.
                ev_tagged = {**ev, "level": level_id}
                if ev.get("type") == "done":
                    done_count["n"] += 1
                    # El 'done' global l'enviem quan tots els nivells han acabat.
                    # Reemetem un 'done_level' individual perquè el frontend sàpiga
                    # que aquest nivell concret ja està llest.
                    ev_tagged["type"] = "done_level"
                events.append(ev_tagged)
            return _cb

        # Rotació a nivell de GRUP (decisió 2026-04-22): si el config està en
        # mode "rotate", triem el model UNA sola vegada i el fem servir a totes
        # les branques del multinivell. Així garantim que les 3 versions del
        # mateix grup són del mateix LLM (no barregem estils a dins d'un grup).
        # Si l'usuari ha indicat un model explícit al payload, té prioritat.
        rotated_model = _model_for("adapt", override=model or "")

        loop = asyncio.get_event_loop()
        with concurrent.futures.ThreadPoolExecutor(max_workers=max(3, len(levels))) as pool:
            tasks = []
            for lvl in levels:
                shift = LEVEL_SHIFTS.get(lvl, 0)
                params_lvl = {**params, "mecr_sortida": _shift_mecr(base_mecr, shift)}
                # level 'single' ⇒ passem id buit (frontend tracta com a mode alumne)
                level_id = lvl if lvl != "single" else ""
                t = loop.run_in_executor(
                    pool,
                    lambda p=params_lvl, l=level_id: run_adaptation(
                        text, profile, context, p, make_cb(l), model_override=rotated_model,
                        docent_id=docent_id
                    ),
                )
                tasks.append(t)

            total = len(tasks)
            all_done = lambda: all(t.done() for t in tasks)
            while not all_done():
                while events:
                    yield f"data: {json.dumps(events.pop(0), ensure_ascii=False)}\n\n"
                yield ": keepalive\n\n"
                await asyncio.sleep(0.5)
            while events:
                yield f"data: {json.dumps(events.pop(0), ensure_ascii=False)}\n\n"

            # Recollim excepcions de tasks que hagin fallat silenciosament
            for t in tasks:
                try:
                    t.result()
                except Exception as task_err:
                    yield f"data: {json.dumps({'type': 'error', 'error': str(task_err)}, ensure_ascii=False)}\n\n"

            # 'done' global quan tots els nivells han acabat
            yield f"data: {json.dumps({'type': 'done', 'total_levels': total}, ensure_ascii=False)}\n\n"

    # Bug 1 (2026-04-19): faltaven headers anti-buffering al SSE. Sense això,
    # nginx/Cloud Run poden bufferar els events i el docent veu la pantalla
    # congelada 60-90s. Mateixos headers que /api/generate-text-stream.
    return StreamingResponse(
        gen(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


# ── Flash (mode ràpid amb prompt MVP) ───────────────────────────────────────

# Mapping (curs, adaptacio) → MECR — idèntic al de mpv/server_mpv.py
_FLASH_CURS_MECR: dict[tuple[str, str], str] = {
    # Base al_nivell = pre-A1 (Infantil P3/P4/P5: lectura emergent compartida amb adult)
    ("infantil",    "molt_simplificat"): "pre-A1",
    ("infantil",    "simplificat"):      "pre-A1",
    ("infantil",    "al_nivell"):        "pre-A1",
    ("infantil",    "enriquiment"):      "A1",
    # Base al_nivell = A1 (Cicle Inicial: comprensió literal, Decret 175/2022)
    ("primaria_12", "molt_simplificat"): "pre-A1",
    ("primaria_12", "simplificat"):      "A1",
    ("primaria_12", "al_nivell"):        "A1",
    ("primaria_12", "enriquiment"):      "A2",
    # Base al_nivell = A1 (Cicle Mitjà: literal + inferencial bàsic)
    ("primaria_34", "molt_simplificat"): "pre-A1",
    ("primaria_34", "simplificat"):      "A1",
    ("primaria_34", "al_nivell"):        "A1",
    ("primaria_34", "enriquiment"):      "A2",
    # Base al_nivell = A2 (Cicle Superior: inferencial + interpretatiu)
    ("primaria_56", "molt_simplificat"): "A1",
    ("primaria_56", "simplificat"):      "A1",
    ("primaria_56", "al_nivell"):        "A2",
    ("primaria_56", "enriquiment"):      "B1",
    # Base al_nivell = B1 (1r-2n ESO: interpretatiu + valoratiu)
    ("eso_12",      "molt_simplificat"): "A1",
    ("eso_12",      "simplificat"):      "A2",
    ("eso_12",      "al_nivell"):        "B1",
    ("eso_12",      "enriquiment"):      "B2",
    # Base al_nivell = B2 (3r-4t ESO: crític, textos argumentatius)
    ("eso_34",      "molt_simplificat"): "B1",
    ("eso_34",      "simplificat"):      "B1",
    ("eso_34",      "al_nivell"):        "B2",
    ("eso_34",      "enriquiment"):      "C1",
    # Base al_nivell = B2 (Batxillerat: crític matisat, textos especialitzats)
    ("batxillerat", "molt_simplificat"): "B1",
    ("batxillerat", "simplificat"):      "B2",
    ("batxillerat", "al_nivell"):        "B2",
    ("batxillerat", "enriquiment"):      "C1",
}

_FLASH_NIVELL_MAP: dict[str, str] = {
    "A1":          "A1 — Lectura Fàcil estricta (AENOR UNE 153101:2018): frases ≤10 paraules, "
                   "vocabulari bàsic, veu activa, una idea per frase, sense subordinades.",
    "A2":          "A2 — Lectura Fàcil adaptada: frases curtes i directes, vocabulari freqüent, "
                   "estructura simple, explica termes amb paraules conegudes.",
    "B1":          "B1 — Llenguatge planer: frases clares, vocabulari estàndard, "
                   "explica termes tècnics entre parèntesis.",
    "B2":          "B2 — Rigor curricular: vocabulari tècnic quan cal, estructura clara, frases fluides.",
    "C1":          "C1 — Text acadèmic estàndard: vocabulari tècnic precís, estructures complexes admeses.",
    "enriquiment": "Enriquiment — Taxonomia de Bloom (anàlisi, síntesi, avaluació): "
                   "aprofundeix conceptes, afegeix connexions interdisciplinàries, invita a la reflexió crítica.",
}

_FLASH_PERFIL_MAP: dict[str, str] = {
    "nouvingut":        "Nouvingut: vocabulari d'alta freqüència, frases curtes, "
                        "explica referents culturals no universals.",
    "tdah":             "TDAH (principis UDL): paràgrafs de 2-3 línies màxim, paraules clau en **negreta**, "
                        "idea principal al principi de cada bloc. "
                        "Defineix els termes tècnics la primera vegada entre parèntesis: terme (definició breu). "
                        "Afegeix una línia de progrés [Secció X de N] entre blocs. "
                        "Inclou una pregunta breu de verificació al final de cada bloc. "
                        "Afegeix un resum curt al final. "
                        "IMPORTANT: no allarguis el text — manté la mateixa extensió que l'original.",
    "dislexia":         "Dislèxia: paraules curtes i freqüents, frases simples (màx. 12 paraules), "
                        "evita sigles i abreviatures, evita encadenar prefixos i sufixos.",
    "tea":              "TEA: llenguatge literal i directe, evita metàfores i ironies, "
                        "estructura previsible i ordenada amb passos numerats, frases afirmatives.",
    "tdl":              "TDL: redueix la densitat lèxica al mínim. "
                        "Estructura SVO estricta (Subjecte-Verb-Object): evita passives i subordinades. "
                        "Cada terme tècnic apareix en 2-3 contextos lleugerament diferents.",
    "di":               "Discapacitat intel·lectual: frases de màxim 8 paraules, una sola idea per frase. "
                        "Cada concepte abstracte amb un exemple tangible i quotidià immediat. "
                        "Vocabulari d'ús quotidià, evita tecnicismes llevat que siguin imprescindibles.",
    "altes_capacitats": "ALERTA — Altes capacitats: PROHIBIT SIMPLIFICAR. "
                        "Mantén o augmenta la complexitat lingüística i conceptual original. "
                        "Afegeix profunditat: excepcions, fronteres del coneixement, debats oberts, "
                        "connexions interdisciplinàries i pensament crític.",
    "vulnerabilitat":   "Vulnerabilitat socioeconòmica: evita supòsits culturals implícits. "
                        "Utilitza exemples i referents accessibles i universals. "
                        "Tona propera i acollidora, sense condescendència.",
    "trastorn_emocional": "Trastorn emocional / ansietat: evita llenguatge que generi pressió o urgència. "
                        "Tona tranquil·la i neutral. Divideix la informació en passos petits.",
    "discalculia":      "Discalcúlia: substitueix nombres abstractes per representacions concretes. "
                        "Desglossa qualsevol seqüència numèrica pas a pas sense saltar-ne cap. "
                        "Per a mesures o estadístiques, afegeix una analogia quotidiana tangible. "
                        "Defineix el vocabulari quantitatiu en context quotidià ('un terç → una de cada tres parts'). "
                        "Si el text exigeix una operació mental implícita, fes-la explícita.",
    "comprensio_lectora": "Comprensió lectora desacoblada: l'alumne descodifica bé però NO construeix significat sol. "
                          "Fes EXPLÍCITES totes les relacions causa-efecte i inferències del text original (no per proximitat: amb connector causal). "
                          "Tradueix nominalitzacions abstractes a estructura SVO concreta ('la dilatació del metall' → 'el metall es fa més gran quan s'escalfa'). "
                          "Substitueix pronoms i referents ambigus pel nom complet. "
                          "A cada bloc, explicita el propòsit ('Llegeix per saber X') i tanca amb una mini-recapitulació ('Idea central: X. Detalls: Y, Z.'). "
                          "Intercala micro-preguntes de metacognició cada 2-3 paràgrafs ('Què acabes de llegir? Resumeix-ho.').",
    "tdc":              "Dispraxia: estructura lineal vertical (evita columnes paral·leles i taules amb >3 columnes). "
                        "Procediments en passos numerats, un pas per línia (no condensis amb 'i després'). "
                        "Si demanes resposta escrita, prioritza formats curts (opció múltiple, V/F, completar amb una paraula) sobre redacció lliure. "
                        "Si el text és llarg, retalla al nucli essencial: la fatiga motora afecta abans que la comprensió.",
}


def _build_flash_system_prompt(
    nivell: str,
    perfils: list[str],
    l1: str = "",
    complements: list[str] | None = None,
    lang: str = "ca",
    caracteristiques: dict | None = None,
) -> str:
    """Prompt simplificat del mode Flash.

    Pot rebre opcionalment `caracteristiques` (estructura canònica del frontend
    nou). Quan està disponible, enriqueix puntualment el glossari amb
    transliteració quan l'alumne nouvingut no usa alfabet llatí. Manté el
    caràcter concís de Flash — no és el lloc per a tot el catàleg.
    """
    from adaptation.lang_config import get_lang_label
    if complements is None:
        complements = ["glossari", "preguntes"]
    lang_label = get_lang_label(lang)

    # Extreu info canònica de nouvingut (si arriba)
    nouv = (caracteristiques or {}).get("nouvingut") or {}
    alfabet_llati = nouv.get("alfabet_llati")
    l1_canonical = nouv.get("l1") or l1
    if isinstance(alfabet_llati, str):
        alfabet_llati = alfabet_llati.lower() not in ("no", "false", "0")
    p = (
        f"Ets un assistent pedagògic especialitzat en adaptació de textos educatius.\n"
        f"Adapta el text que t'enviaré EN {lang_label.upper()}.\n"
        "IMPORTANT: la primera línia ha de ser sempre un títol en format `## Títol`. "
        "Si el text original ja en té, conserva'l. Si no en té, crea'n un de breu i descriptiu del contingut. "
        "Si el text adaptat supera 100 paraules sense apartats, organitza'l en 2-3 subapartats `### Apartat`. "
        "No escriguis cap frase introductòria, cap comentari sobre el que has fet ni cap text de tancament.\n\n"
    )
    p += "NIVELL:\n" + _FLASH_NIVELL_MAP.get(nivell, _FLASH_NIVELL_MAP["B1"]) + "\n"
    if perfils:
        p += "\nPERFILS DE L'ALUMNAT:\n"
        for pf in perfils:
            if pf in _FLASH_PERFIL_MAP:
                p += f"- {_FLASH_PERFIL_MAP[pf]}\n"
    comp_lines = []
    if "glossari" in complements:
        if "nouvingut" in perfils and l1_canonical:
            extra_alfabet = ""
            if alfabet_llati is False:
                extra_alfabet = (
                    f" La L1 ({l1_canonical}) NO usa alfabet llatí: afegeix també "
                    "transliteració fonètica del terme català al costat de la traducció."
                )
            comp_lines.append(
                f"- GLOSSARI: taula markdown 3 columnes | Terme | Traducció ({l1_canonical}, alfabet original) | Explicació (màx. 8 paraules en {lang_label}) |. "
                f"6-10 termes: prioritza termes curriculars del text, paraules ambigues pel nivell i col·locacions clau.{extra_alfabet}"
            )
        else:
            comp_lines.append(
                f"- GLOSSARI: taula markdown 2 columnes | Terme | Explicació (màx. 8 paraules en {lang_label}) |. "
                f"6-10 termes: prioritza termes curriculars del text, paraules ambigues pel nivell i col·locacions clau."
            )
    if "preguntes" in complements:
        comp_lines.append(
            f"- PREGUNTES DE COMPRENSIÓ: 5-6 preguntes en {lang_label} en 3 moments — "
            "1 abans de llegir (hipòtesi o connexió prèvia), 1 durant (inferència o lèxic en context), "
            "3-4 després (1 literal, 1-2 inferencials, 1 crítica). Cada pregunta comença amb «- »."
        )
    if "resum" in complements:
        comp_lines.append(f"- RESUM: 3-5 frases en {lang_label} que resumeixin les idees principals del text adaptat.")
    if comp_lines:
        p += (
            "\nCOMPLEMENTS (afegeix al final del text adaptat, separats amb EXACTAMENT "
            "aquests títols en majúscules: GLOSSARI / PREGUNTES DE COMPRENSIÓ / RESUM):\n"
        )
        p += "\n".join(comp_lines) + "\n"
    return p


def _parse_flash_response(raw: str) -> dict:
    """Separa el text adaptat i els complements (glossari, preguntes, resum).

    Patrons robustos: admeten markdown (##, **), accent variat (COMPRENSIO/COMPRENSIÓ),
    i capçaleres sense línia en blanc prèvia. Retorna dict: adapted, glossari, preguntes, resum.
    """
    import re
    text = raw.strip()

    # Prefix markdown opcional: ## / ** (amb espais addicionals possibles)
    _MD = r"(?:\*{1,2}|#{1,3}[ \t]*)?"
    # Suffix: qualsevol combinació de *, : i espais fins a fi de línia
    # Cobreix: "GLOSSARI", "GLOSSARI:", "**GLOSSARI**", "**GLOSSARI:**", etc.
    _SUF = r"[ \t:*]*"

    patterns = {
        # ca: GLOSSARI | es: GLOSARIO | en: GLOSSARY | fr: GLOSSAIRE | eu: HIZTEGIA | gl: GLOSARIO
        "glossari": re.compile(
            rf"^[ \t]*{_MD}(?:GLOSSARI|GLOSARIO|GLOSSARY|GLOSSAIRE|HIZTEGIA){_SUF}$",
            re.IGNORECASE | re.MULTILINE,
        ),
        # ca: PREGUNTES (DE COMPRENSIÓ) | es: PREGUNTAS | en: QUESTIONS | fr: QUESTIONS | eu: GALDERAK
        "preguntes": re.compile(
            rf"^[ \t]*{_MD}(?:PREGUNTES(?:[ \t]+DE[ \t]+COMPRENS(?:[IÍ][ÓO]|IO))?|PREGUNTAS(?:[ \t]+DE[ \t]+COMPRENSI[OÓ]N)?|QUESTIONS?(?:[ \t]+(?:OF|DE)[ \t]+\w+)?|GALDERAK){_SUF}$",
            re.IGNORECASE | re.MULTILINE,
        ),
        # ca: RESUM | es: RESUMEN | en: SUMMARY | fr: RÉSUMÉ/RESUME | eu: LABURPENA | gl: RESUMO
        "resum": re.compile(
            rf"^[ \t]*{_MD}(?:RESUM(?:EN|O)?|SUMMARY|R[EÉ]SUM[EÉ]?|LABURPENA){_SUF}$",
            re.IGNORECASE | re.MULTILINE,
        ),
    }

    marks = []  # (pos_start, pos_end, key)
    for key, pat in patterns.items():
        m = pat.search(text)
        if m:
            marks.append((m.start(), m.end(), key))
    marks.sort()

    result = {"adapted": text, "glossari": "", "preguntes": "", "resum": ""}
    if not marks:
        return result

    result["adapted"] = text[: marks[0][0]].strip()
    for i, (start, end, key) in enumerate(marks):
        next_start = marks[i + 1][0] if i + 1 < len(marks) else len(text)
        result[key] = text[end:next_start].strip()
    return result


@app.post("/api/adapt-flash")
async def adapt_flash(payload: dict = Body(...)):
    """Mode Flash: adapta amb el prompt MVP (senzill) i SSE.

    Payload:
      text (str):         text original
      curs (str):         primaria_12 | primaria_34 | primaria_56 | eso_12 | eso_34 | batxillerat
      adaptacio (str):    al_nivell | simplificat | molt_simplificat | enriquiment
      perfils (list[str]):clínics: dislexia | tdah | tea | nouvingut | altes_capacitats | di | vulnerabilitat
      l1 (str):           llengua materna (si nouvingut)
      tipus (str):        grup | alumne  (meta-dada)
      docent_id (str):    per desar a history
    """
    text        = (payload.get("text") or "").strip()
    curs        = payload.get("curs", "eso_12")
    adaptacio   = payload.get("adaptacio", "al_nivell")
    perfils     = payload.get("perfils") or []
    l1          = (payload.get("l1") or "").strip()
    tipus       = payload.get("tipus", "grup")
    complements = payload.get("complements") or ["glossari", "preguntes"]
    docent_id   = (payload.get("docent_id") or "").strip()
    lang        = (payload.get("lang") or "ca").strip()
    # Camp opcional: estructura canònica enriquida (subvariables per condició).
    # El frontend nou (profile-canonical.js) la inclou automàticament; els clients
    # antics no la passen i el comportament queda com abans.
    caracteristiques = payload.get("caracteristiques") or None

    if not text:
        return JSONResponse({"error": "El text és buit."}, status_code=400)

    nivell        = _FLASH_CURS_MECR.get((curs, adaptacio), "B1")
    model_id      = _model_for("adapt_flash")
    system_prompt = _build_flash_system_prompt(nivell, perfils, l1, complements, lang=lang, caracteristiques=caracteristiques)
    t0 = __import__("time").time()

    async def gen():
        yield f"data: {json.dumps({'type': 'step', 'msg': 'Generant adaptació…', 'model': model_id}, ensure_ascii=False)}\n\n"

        try:
            loop = asyncio.get_event_loop()
            raw = await loop.run_in_executor(
                None,
                lambda: _call_llm(model_id, system_prompt, text),
            )
        except Exception as exc:
            yield f"data: {json.dumps({'type': 'error', 'msg': str(exc)}, ensure_ascii=False)}\n\n"
            return

        parts = _parse_flash_response(raw)
        adapted   = parts["adapted"]
        glossari  = parts["glossari"]
        preguntes = parts["preguntes"]
        resum     = parts.get("resum", "")
        duration_ms = int((__import__("time").time() - t0) * 1000)

        yield f"data: {json.dumps({'type': 'result', 'adapted': adapted, 'glossari': glossari, 'preguntes': preguntes, 'resum': resum}, ensure_ascii=False)}\n\n"
        yield f"data: {json.dumps({'type': 'done'}, ensure_ascii=False)}\n\n"

        # Desa a history en background (no bloqueja el SSE)
        try:
            n_in  = len(text.split())
            n_out = len(adapted.split())
            row = {
                "profile_name":  ", ".join(perfils) if perfils else "al_nivell",
                "profile_json":  {"perfils": perfils, "curs": curs, "adaptacio": adaptacio, "l1": l1, "tipus": tipus},
                "context_json":  {"curs": curs, "etapa": curs.split("_")[0] if "_" in curs else curs},
                "params_json":   {"nivell": nivell, "adaptacio": adaptacio},
                "original_text": text,
                "adapted_text":  adapted,
                "model_used":    model_id,
                "endpoint":      "/api/adapt-flash",
                "duration_ms":   duration_ms,
                "etapa":         curs.split("_")[0] if "_" in curs else curs,
                "curs":          curs,
                "n_words_in":    n_in,
                "n_words_out":   n_out,
                "via":           "flash",
                "mode":          "flash",
                "prompt_version": ATNE_PROMPT_VERSION,
            }
            if docent_id:
                row["docent_hash"] = docent_id
            requests.post(
                f"{SUPABASE_URL}/rest/v1/history",
                headers={**SUPABASE_HEADERS, "Prefer": "return=minimal"},
                json=row,
                timeout=8,
            )
        except Exception:
            pass  # No bloquejar el SSE per error de desa

    return StreamingResponse(
        gen(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


# ── Il·lustracions (complement beta) ────────────────────────────────────────

@app.post("/api/illustration")
async def resolve_illustration(payload: dict = Body(...)):
    """Resol un marcador [IMATGE: concept_ca] disparant Wikimedia + FLUX en
    paral·lel. Retorna les dues opcions perquè el docent en triï.

    Body:
      concept (str, obligatori): concepte curt en català del marcador.
      context (dict, opcional): {mecr, subject, ...}.
      style (str, opcional): preset FLUX (default aquarela_storybook).
      seed (int, opcional): seed compartida pel document (default 42).
    """
    from adaptation.illustrations import resolve_marker, STYLE_SPINES, DEFAULT_STYLE

    concept = (payload.get("concept") or "").strip()
    if not concept:
        raise HTTPException(status_code=400, detail="falta concept")
    context = payload.get("context") or {}
    style = payload.get("style") or DEFAULT_STYLE
    if style not in STYLE_SPINES:
        style = DEFAULT_STYLE
    seed = int(payload.get("seed") or 42)

    resolution = await asyncio.to_thread(
        resolve_marker, concept, context, style, seed
    )
    return resolution.to_dict()


@app.post("/api/illustrations/batch")
async def resolve_illustrations_batch(payload: dict = Body(...)):
    """Resol tots els marcadors d'un text en paral·lel.

    Body:
      text (str, obligatori): text adaptat amb marcadors [IMATGE: ...].
      context (dict, opcional): {mecr, subject, ...}.
      style (str, opcional): preset FLUX.
      seed (int, opcional): seed compartida.

    Retorna: {markers: [{start, end, concept, resolution}, ...]}
    """
    from adaptation.illustrations import resolve_all_markers, STYLE_SPINES, DEFAULT_STYLE

    text = payload.get("text") or ""
    if not text:
        raise HTTPException(status_code=400, detail="falta text")
    context = payload.get("context") or {}
    style = payload.get("style") or DEFAULT_STYLE
    if style not in STYLE_SPINES:
        style = DEFAULT_STYLE
    seed = int(payload.get("seed") or 42)

    results = await asyncio.to_thread(
        resolve_all_markers, text, context, style, seed
    )
    return {"markers": results, "count": len(results), "style": style}


# ── Exportació ──────────────────────────────────────────────────────────────

@app.post("/api/export")
async def export_doc(payload: dict = Body(...)):
    # TODO: Dockerfile ha d'instal·lar fonts-noto o fonts-amiri perquè àrab es
    # renderitzi a Cloud Run (aquest server.py afegeix els paths i els rangs
    # Unicode, però si la font no està instal·lada al contenidor els glifs
    # continuen sortint com tofu). Veure bug 3 (2026-04-19).
    fmt = payload.get("format", "txt")
    adapted = payload.get("adapted", "")
    original = payload.get("original", "")
    # Per defecte NO incloem el text original: l'usuari vol el text adaptat net.
    # Opt-in amb include_original=true si algun flux futur el necessita.
    include_original = bool(payload.get("include_original", False))
    profile_name = payload.get("profile_name", "adaptacio")
    # Sprint 1C (2026-04-22): complements actius seleccionats pel docent.
    # Format: { backend_key: markdown_string }. Ordre de presentació fix.
    complements = payload.get("complements") or {}

    # Mapa de títols en català + ordre de presentació al document exportat.
    # Si afegeixes un complement al catàleg, recorda afegir-lo aquí.
    _COMP_TITLES: dict[str, str] = {
        "glossari":                 "Glossari",
        "esquema_visual":           "Esquema visual",
        "preguntes_comprensio":     "Preguntes de comprensió",
        "mapa_conceptual":          "Mapa conceptual",
        "mapa_mental":              "Mapa mental",
        "bastides":                 "Bastides (scaffolding)",
        "activitats_aprofundiment": "Activitats d'aprofundiment",
        "pictogrames":              "Pictogrames",
    }
    _COMP_ORDER: list[str] = [
        "glossari", "esquema_visual", "mapa_conceptual", "mapa_mental",
        "bastides", "preguntes_comprensio", "activitats_aprofundiment",
        "pictogrames",
    ]

    def _complements_as_markdown() -> str:
        """Serialitza els complements com a seccions markdown ordenades.
        Cada secció arrenca amb `## Títol` i separa amb línia en blanc.
        Retorna string buit si no hi ha cap complement. Ja es neteja d'una
        possible duplicació del header (alguns prompts retornen el títol
        inclòs dins del markdown del complement)."""
        parts: list[str] = []
        for key in _COMP_ORDER:
            md = (complements.get(key) or "").strip()
            if not md:
                continue
            title = _COMP_TITLES.get(key, key)
            # Treu un header duplicat al principi del markdown si apareix
            # (ex: "## Glossari\n...") perquè no es dupliqui amb el nostre.
            md_clean = re.sub(
                r'^\s*#+\s*' + re.escape(title) + r'\s*\n',
                '', md, count=1, flags=re.IGNORECASE,
            )
            parts.append(f"## {title}\n\n{md_clean}")
        return "\n\n".join(parts)

    # Si el docent ha demanat complements, els afegim al final del markdown
    # abans de processar-lo per a cada format. Així el pipeline de
    # PDF/DOCX/TXT els tracta igual que qualsevol secció del text adaptat.
    complements_md = _complements_as_markdown()
    if complements_md:
        adapted = (adapted.rstrip() + "\n\n---\n\n" + complements_md).strip()

    import tempfile
    safe_name = re.sub(r'[^\w\s-]', '', profile_name).strip().replace(" ", "_")[:30] or "adaptacio"
    timestamp = time.strftime("%Y%m%d_%H%M")
    base_name = f"ATNE_{safe_name}_{timestamp}"

    # Sprint 1C (2026-04-22): sanititzacions pre-PDF per a complements
    # específics que el LLM genera amb caràcters box-drawing (mapa_conceptual,
    # esquema_visual) — es filtrarien a `pdf_safe` i es perdria tota
    # l'estructura. Substituïm per equivalents ASCII abans del filtre.
    _BOX_TO_ASCII = {
        "│": "|", "├": "|-", "└": "`-", "┌": ".-", "┐": "-.", "┘": "-'",
        "┤": "-|", "┬": "-T-", "┴": "-+-", "┼": "-+-",
        "─": "-", "━": "-", "═": "=",
        "╔": "+", "╗": "+", "╚": "+", "╝": "+",
        "║": "|", "╟": "+", "╢": "+", "╤": "+", "╧": "+",
        "▶": ">", "▲": "^", "▼": "v", "◆": "*", "●": "*", "■": "*",
    }

    def _sanitize_for_pdf(text: str) -> str:
        """Substitueix caràcters box-drawing i altres glifs no-renderitzables
        per equivalents ASCII abans de pdf_safe. Evita que mapes conceptuals
        i esquemes visuals es perdin completament."""
        for src, dst in _BOX_TO_ASCII.items():
            text = text.replace(src, dst)
        return text

    def pdf_safe(text):
        """Filtra text per PDF: només caràcters que Arial/Liberation pot renderitzar.
        Whitelist: ASCII, Latin Extended, àrab, puntuació general, fletxes, math."""
        text = _sanitize_for_pdf(text)
        cleaned = []
        for c in text:
            cp = ord(c)
            if c in (' ', '\t', '\n', '\r'):
                cleaned.append(c)
            elif 0x20 <= cp <= 0x7E:       # ASCII imprimible
                cleaned.append(c)
            elif 0x00A0 <= cp <= 0x024F:   # Latin Extended (à é í ò ú ç · ñ ü)
                cleaned.append(c)
            elif 0x0600 <= cp <= 0x06FF:   # Àrab bàsic
                cleaned.append(c)
            elif 0x0750 <= cp <= 0x077F:   # Àrab suplementari
                cleaned.append(c)
            elif 0xFB50 <= cp <= 0xFDFF:   # Formes de presentació A (àrab)
                cleaned.append(c)
            elif 0xFE70 <= cp <= 0xFEFF:   # Formes de presentació B (àrab)
                cleaned.append(c)
            elif 0x2000 <= cp <= 0x206F:   # Puntuació general (— ' " …)
                cleaned.append(c)
            elif 0x2070 <= cp <= 0x209F:   # Superíndexs/subíndexs (₂)
                cleaned.append(c)
            elif 0x20A0 <= cp <= 0x20CF:   # Símbols de moneda (€ £)
                cleaned.append(c)
            elif 0x2190 <= cp <= 0x21FF:   # Fletxes (↓ → ← ↑)
                cleaned.append(c)
            elif 0x2200 <= cp <= 0x22FF:   # Símbols matemàtics (± × ÷)
                cleaned.append(c)
            # Tot el resta (emojis, box-drawing, variation selectors, ZWJ): omès
        return "".join(cleaned)

    def clean_for_plain(text):
        """Neteja markdown per a TXT pla."""
        text = text.replace("**", "").replace("*", "").replace("`", "")
        text = re.sub(r'^#{1,4}\s+', '', text, flags=re.MULTILINE)  # Treure headings #
        text = re.sub(r'^\|[-:\s|]+\|$', '', text, flags=re.MULTILINE)  # Treure separadors taula
        return text

    if fmt == "txt":
        tmp = Path(tempfile.gettempdir()) / f"{base_name}.txt"
        clean_adapted = clean_for_plain(adapted)
        parts = [f"ADAPTACIÓ ATNE — {profile_name}", "=" * 50, "", clean_adapted]
        if include_original and original:
            parts += ["", "=" * 50, "TEXT ORIGINAL:", "", original]
        content = "\n".join(parts)
        # BOM UTF-8 perquè Windows/Notepad el reconegui correctament
        with open(tmp, "w", encoding="utf-8-sig") as f:
            f.write(content)
        return FileResponse(tmp, filename=f"{base_name}.txt", media_type="text/plain; charset=utf-8")

    elif fmt == "docx":
        from docx import Document as DocxDocument
        from docx.shared import Pt, Inches
        import base64
        import io
        doc = DocxDocument()
        doc.styles["Normal"].font.size = Pt(11)
        doc.add_heading(f"Adaptació ATNE — {profile_name}", level=1)

        def add_inline(p, text):
            """Afegeix text a un paràgraf gestionant **negretes** i _cursives_ inline."""
            # Primer partim per negretes, després per cursives dins de cada tros.
            for part in re.split(r'(\*\*.*?\*\*)', text):
                if part.startswith("**") and part.endswith("**") and len(part) > 4:
                    run = p.add_run(part[2:-2])
                    run.bold = True
                else:
                    for sub in re.split(r'(_[^_]+_)', part):
                        if sub.startswith("_") and sub.endswith("_") and len(sub) > 2:
                            run = p.add_run(sub[1:-1])
                            run.italic = True
                        elif sub:
                            p.add_run(sub)

        def try_add_image(caption: str, url: str) -> bool:
            """Descarrega la imatge i l'incrusta. Retorna True si OK, False si falla.
            Suporta URLs http(s) (Pollinations, Wikimedia, Cloud Storage) i data: URIs."""
            try:
                if url.startswith("data:"):
                    # data:image/png;base64,...
                    header, b64 = url.split(",", 1)
                    raw = base64.b64decode(b64)
                elif url.startswith(("http://", "https://")):
                    # Wikimedia Commons rebutja el User-Agent per defecte de
                    # python-requests (403). Posem un UA descriptiu com fa
                    # adaptation/illustrations.py al search.
                    headers = {
                        "User-Agent": "ATNE-FJE-EducationalBot/1.0 (https://atne.fje.edu)",
                        "Accept": "image/*,*/*;q=0.8",
                    }
                    r = requests.get(url, timeout=20, headers=headers)
                    r.raise_for_status()
                    ct = (r.headers.get("Content-Type") or "").lower()
                    if not ct.startswith("image/"):
                        print(f"[export docx] resposta no-imatge ({ct}) per {url[:80]}")
                        return False
                    raw = r.content
                else:
                    return False
                stream = io.BytesIO(raw)
                # Amplada max 6 polzades (~15cm) perquè càpiga bé dins d'un A4.
                doc.add_picture(stream, width=Inches(6))
                if caption:
                    p = doc.add_paragraph()
                    run = p.add_run(caption)
                    run.italic = True
                    run.font.size = Pt(10)
                return True
            except Exception as e:
                print(f"[export docx] no s'ha pogut incrustar imatge {url[:80]}: {e}")
                return False

        img_re = re.compile(r'!\[([^\]]*)\]\(([^)]+)\)')

        for line in adapted.split("\n"):
            stripped = line.strip()
            # Imatge markdown ![caption](url) — provem d'incrustar-la.
            img_match = img_re.match(stripped)
            if img_match:
                caption, url = img_match.group(1), img_match.group(2)
                if not try_add_image(caption, url):
                    # Fallback: text amb caption
                    p = doc.add_paragraph()
                    run = p.add_run(f"[Imatge: {caption}]")
                    run.italic = True
                continue
            # H1 (#) → estil Heading 1; H2 (##) → Heading 2; H3 (###) → Heading 3
            if line.startswith("# ") and not line.startswith("## "):
                doc.add_heading(line[2:].replace("**", ""), level=1)
            elif line.startswith("## "):
                doc.add_heading(line[3:].replace("**", ""), level=2)
            elif line.startswith("### "):
                doc.add_heading(line[4:].replace("**", ""), level=3)
            elif stripped.startswith("|"):
                clean = stripped.replace("|", "  ").strip()
                if clean and not all(c in "-: " for c in clean):
                    p = doc.add_paragraph(clean)
                    p.style.font.size = Pt(11)
            elif stripped.startswith("```"):
                continue
            elif re.match(r'^[\s]*[-*]\s+(?!\*)', line) and not stripped.startswith("**"):
                clean = re.sub(r'^[\s]*[-*]\s+', '', line).strip()
                p = doc.add_paragraph(style="List Bullet")
                add_inline(p, clean)
            elif re.match(r'^[\s]*\d+\.\s+', line):
                clean = re.sub(r'^[\s]*\d+\.\s+', '', line).strip()
                p = doc.add_paragraph(style="List Number")
                add_inline(p, clean)
            elif stripped:
                p = doc.add_paragraph()
                add_inline(p, stripped)

        if include_original and original:
            doc.add_page_break()
            doc.add_heading("Text original", level=2)
            for line in original.split("\n"):
                if line.strip():
                    doc.add_paragraph(line)
        tmp = Path(tempfile.gettempdir()) / f"{base_name}.docx"
        doc.save(str(tmp))
        return FileResponse(tmp, filename=f"{base_name}.docx",
                          media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    elif fmt == "pdf":
        from fpdf import FPDF
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        # Font Unicode del sistema (Arial suporta català, àrab, etc.)
        font_name = "Helvetica"  # fallback
        # Bug 3 (2026-04-19): ordre de prioritat pensat perquè a Cloud Run (Linux)
        # agafi una font amb glifs àrabs abans de caure en Liberation (que NO té
        # àrab i retorna tofu). El Dockerfile ha d'instal·lar fonts-noto o
        # fonts-amiri perquè els primers paths Linux existeixin.
        for ttf_normal, ttf_bold, fname in [
            ("C:/Windows/Fonts/arial.ttf", "C:/Windows/Fonts/arialbd.ttf", "ArialUni"),
            ("C:/Windows/Fonts/segoeui.ttf", "C:/Windows/Fonts/segoeuib.ttf", "SegoeUI"),
            # Linux (Cloud Run / Docker) — primer les que suporten àrab
            ("/usr/share/fonts/truetype/noto/NotoSansArabic-Regular.ttf",
             "/usr/share/fonts/truetype/noto/NotoSansArabic-Bold.ttf", "NotoArabic"),
            ("/usr/share/fonts/truetype/amiri/amiri-regular.ttf",
             "/usr/share/fonts/truetype/amiri/amiri-bold.ttf", "Amiri"),
            ("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
             "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf", "LiberationSans"),
            ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
             "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", "DejaVuSans"),
        ]:
            if Path(ttf_normal).exists():
                try:
                    pdf.add_font(fname, "", ttf_normal)
                    pdf.add_font(fname, "B", ttf_bold)
                    font_name = fname
                    break
                except Exception:
                    continue
        pdf.add_page()
        w = pdf.w - pdf.l_margin - pdf.r_margin  # amplada efectiva

        def pdf_clean(txt):
            """Neteja markdown inline i caràcters no-renderitzables per PDF."""
            txt = pdf_safe(txt)
            return txt.replace("**", "").replace("*", "").replace("`", "")

        def pdf_write_line(line):
            """Escriu una línia al PDF amb detecció de format."""
            stripped = line.strip()
            if not stripped:
                pdf.ln(3)
                return
            # Heading # (H1)
            if line.startswith("# ") and not line.startswith("## "):
                pdf.set_font(font_name, "B", 15)
                pdf.ln(5)
                pdf.multi_cell(w, 9, pdf_clean(line[2:]), align="L",
                               new_x="LMARGIN", new_y="NEXT")
                pdf.set_font(font_name, "", 11)
                pdf.ln(2)
            # Heading ## (H2)
            elif line.startswith("## "):
                pdf.set_font(font_name, "B", 13)
                pdf.ln(4)
                pdf.multi_cell(w, 8, pdf_clean(line[3:]), align="L",
                               new_x="LMARGIN", new_y="NEXT")
                pdf.set_font(font_name, "", 11)
                pdf.ln(2)
            # Heading ### (H3)
            elif line.startswith("### "):
                pdf.set_font(font_name, "B", 11)
                pdf.ln(2)
                pdf.multi_cell(w, 7, pdf_clean(line[4:]), align="L",
                               new_x="LMARGIN", new_y="NEXT")
                pdf.set_font(font_name, "", 11)
            # Taules markdown
            elif stripped.startswith("|"):
                clean = stripped.replace("|", "  ").strip()
                # Saltar línies separadores (|---|---|)
                if clean and not all(c in "-: " for c in clean):
                    pdf.set_font(font_name, "", 9)
                    pdf.multi_cell(w, 5, pdf_clean(clean), align="L",
                                   new_x="LMARGIN", new_y="NEXT")
                    pdf.set_font(font_name, "", 11)
            # Bloc de codi
            elif stripped.startswith("```"):
                pass  # Saltar delimitadors
            # Bullets: "- text" o "* text" (però NO "**negreta**")
            elif re.match(r'^[\s]*[-*]\s+(?!\*)', line) and not stripped.startswith("**"):
                clean = re.sub(r'^[\s]*[-*]\s+', '', line).strip()
                pdf.multi_cell(w, 6, f"  - {pdf_clean(clean)}", align="L",
                               new_x="LMARGIN", new_y="NEXT")
            # Llistes numerades: "1. text"
            elif re.match(r'^[\s]*\d+\.\s+', line):
                clean = re.sub(r'^[\s]*\d+\.\s+', '', line).strip()
                num = re.match(r'[\s]*(\d+)\.', line).group(1)
                pdf.multi_cell(w, 6, f"  {num}. {pdf_clean(clean)}", align="L",
                               new_x="LMARGIN", new_y="NEXT")
            # Línia normal
            else:
                pdf.multi_cell(w, 6, pdf_clean(stripped), align="L",
                               new_x="LMARGIN", new_y="NEXT")

        pdf.set_font(font_name, "B", 16)
        pdf.cell(w, 10, pdf_safe(f"Adaptació ATNE — {profile_name}"), new_x="LMARGIN", new_y="NEXT")
        pdf.set_font(font_name, "", 11)
        pdf.ln(5)

        for line in adapted.split("\n"):
            pdf_write_line(line)

        if include_original and original:
            pdf.add_page()
            pdf.set_font(font_name, "B", 13)
            pdf.cell(w, 10, "Text original", new_x="LMARGIN", new_y="NEXT")
            pdf.set_font(font_name, "", 11)
            pdf.ln(3)
            for orig_line in original.split("\n"):
                if orig_line.strip():
                    pdf.multi_cell(w, 6, pdf_safe(orig_line), align="L",
                                   new_x="LMARGIN", new_y="NEXT")
                else:
                    pdf.ln(3)
        # Footer
        pdf.ln(10)
        pdf.set_font(font_name, "", 8)
        pdf.cell(w, 5, "Generat per ATNE — Jesuïtes Educació",
                 new_x="LMARGIN", new_y="NEXT", align="C")
        tmp = Path(tempfile.gettempdir()) / f"{base_name}.pdf"
        pdf.output(str(tmp))
        return FileResponse(tmp, filename=f"{base_name}.pdf", media_type="application/pdf")

    return JSONResponse({"error": f"Format '{fmt}' no suportat"}, status_code=400)


# ── Admin (configuració runtime del pilot) ─────────────────────────────────

@app.get("/admin", response_class=HTMLResponse)
async def admin_page():
    """Serveix el dashboard /admin. L'auth es gestiona via JS al navegador
    contra /api/admin/whoami (no bloquegem la pàgina aquí perquè el flow
    de login viu dins la mateixa HTML)."""
    html_path = UI_DIR / "admin.html"
    if html_path.exists():
        return HTMLResponse(
            html_path.read_text(encoding="utf-8"),
            headers={"Cache-Control": "no-cache, no-store, must-revalidate"},
        )
    return HTMLResponse("<h1>Admin no disponible</h1>", status_code=404)


@app.get("/admin/pilot", response_class=HTMLResponse)
async def admin_pilot_page():
    """Serveix el dashboard `/admin/pilot` — funnel d'adopció + edit-rate +
    cost per model + feedback agregat per al pilot 2026-04-20→2026-05-08.
    L'auth es gestiona via JS al navegador (mateix patró que /admin)."""
    html_path = UI_DIR / "admin-pilot.html"
    if html_path.exists():
        return HTMLResponse(
            html_path.read_text(encoding="utf-8"),
            headers={"Cache-Control": "no-cache, no-store, must-revalidate"},
        )
    return HTMLResponse("<h1>Pilot dashboard no disponible</h1>", status_code=404)


# ── Cuina (dashboard intern) ───────────────────────────────────────────────

@app.get("/cuina", response_class=HTMLResponse)
async def cuina_page(_: bool = Depends(_require_admin)):
    """Serveix la pàgina de cuina (flux + catàleg d'instruccions). Requereix rol admin."""
    html_path = UI_DIR / "cuina.html"
    if html_path.exists():
        return HTMLResponse(
            html_path.read_text(encoding="utf-8"),
            headers={"Cache-Control": "no-cache, no-store, must-revalidate"},
        )
    return HTMLResponse("<h1>Cuina no disponible</h1>", status_code=404)


@app.get("/pipeline", response_class=HTMLResponse)
async def pipeline_page(_: bool = Depends(_require_admin)):
    """Pàgina viva /pipeline. Requereix rol admin."""
    html_path = UI_DIR / "pipeline.html"
    if html_path.exists():
        return HTMLResponse(
            html_path.read_text(encoding="utf-8"),
            headers={"Cache-Control": "no-cache, no-store, must-revalidate"},
        )
    return HTMLResponse("<h1>Pipeline no disponible</h1>", status_code=404)


@app.get("/saber-ne", response_class=HTMLResponse)
async def saber_ne_page():
    """Serveix la pàgina Saber-ne+ (fonaments pedagògics per a docents)."""
    html_path = UI_DIR / "saber-ne.html"
    if html_path.exists():
        return HTMLResponse(
            html_path.read_text(encoding="utf-8"),
            headers={"Cache-Control": "no-cache, no-store, must-revalidate"},
        )
    return HTMLResponse("<h1>Saber-ne+ no disponible</h1>", status_code=404)


@app.get("/avaluacio", response_class=HTMLResponse)
async def avaluacio_page():
    """Serveix el hub d'avaluació i decisions (Bloc 3 de Saber-ne+)."""
    html_path = UI_DIR / "avaluacio.html"
    if html_path.exists():
        return HTMLResponse(
            html_path.read_text(encoding="utf-8"),
            headers={"Cache-Control": "no-cache, no-store, must-revalidate"},
        )
    return HTMLResponse("<h1>Avaluació no disponible</h1>", status_code=404)


@app.get("/demo", response_class=HTMLResponse)
async def demo_v7_page():
    """Serveix el mockup visual V7 de Stitch (Pas 1 redissenyat). Estàtic."""
    html_path = UI_DIR / "v7_demo.html"
    if html_path.exists():
        return HTMLResponse(
            html_path.read_text(encoding="utf-8"),
            headers={"Cache-Control": "no-cache, no-store, must-revalidate"},
        )
    return HTMLResponse("<h1>Demo V7 no disponible</h1>", status_code=404)


@app.get("/dashboard", response_class=HTMLResponse)
async def dashboard_xat9_page():
    """Serveix el dashboard del Xat 9 (experiment A/B multi-model)."""
    html_path = UI_DIR / "dashboard_xat9.html"
    if html_path.exists():
        return HTMLResponse(
            html_path.read_text(encoding="utf-8"),
            headers={"Cache-Control": "no-cache, no-store, must-revalidate"},
        )
    return HTMLResponse("<h1>Dashboard Xat 9 no disponible</h1>", status_code=404)


@app.get("/dashboard_complements", response_class=HTMLResponse)
async def dashboard_complements_page():
    html_path = UI_DIR / "dashboard_complements.html"
    if html_path.exists():
        return HTMLResponse(
            html_path.read_text(encoding="utf-8"),
            headers={"Cache-Control": "no-cache, no-store, must-revalidate"},
        )
    return HTMLResponse("<h1>Dashboard Complements no disponible</h1>", status_code=404)


@app.get("/dashboard_questions", response_class=HTMLResponse)
async def dashboard_questions_page():
    html_path = UI_DIR / "dashboard_questions.html"
    if html_path.exists():
        return HTMLResponse(
            html_path.read_text(encoding="utf-8"),
            headers={"Cache-Control": "no-cache, no-store, must-revalidate"},
        )
    return HTMLResponse("<h1>Dashboard Questions no disponible</h1>", status_code=404)


@app.get("/informe_fje")
async def informe_fje_md():
    """Descàrrega de l'informe executiu FJE en Markdown."""
    md_path = Path(__file__).parent / "tests" / "experiment_ab" / "informe_executiu_FJE.md"
    if md_path.exists():
        return FileResponse(
            str(md_path),
            filename="informe_executiu_FJE.md",
            media_type="text/markdown; charset=utf-8",
        )
    return HTMLResponse("<h1>Informe FJE no disponible</h1>", status_code=404)


@app.get("/informe_tecnic")
async def informe_tecnic_md():
    """Descàrrega de l'informe tècnic multi-model en Markdown."""
    md_path = Path(__file__).parent / "tests" / "experiment_ab" / "informe_multi_model.md"
    if md_path.exists():
        return FileResponse(
            str(md_path),
            filename="informe_multi_model.md",
            media_type="text/markdown; charset=utf-8",
        )
    return HTMLResponse("<h1>Informe tècnic no disponible</h1>", status_code=404)


@app.get("/validacio", response_class=HTMLResponse)
async def validacio_page():
    """Serveix la pàgina de validació humana."""
    html_path = UI_DIR / "validacio.html"
    if html_path.exists():
        return HTMLResponse(html_path.read_text(encoding="utf-8"),
                            headers={"Cache-Control": "no-cache, no-store, must-revalidate"})
    return HTMLResponse("<h1>Validació no disponible</h1>", status_code=404)


@app.get("/validacio_data.json")
async def validacio_data():
    """Serveix les dades de validació estàtiques (det)."""
    json_path = UI_DIR / "validacio_data.json"
    if json_path.exists():
        return JSONResponse(json.loads(json_path.read_text(encoding="utf-8")))
    return JSONResponse([])


@app.get("/api/validacio/{tanda}")
async def api_validacio_tanda(tanda: str):
    """Retorna textos adaptats + puntuacions per a una tanda (prompt_mode)."""
    import sqlite3
    db_path = Path(__file__).parent / "tests" / "results" / "evaluations.db"
    if not db_path.exists():
        return JSONResponse([])
    conn = sqlite3.connect(str(db_path))
    conn.row_factory = sqlite3.Row
    # Agafar tots els casos d'aquesta tanda amb avaluació GPT
    rows = conn.execute("""
        SELECT g.cas_id, g.generator, g.perfil_nom, g.mecr, g.dua, g.text_adaptat_paraules,
               g.text_original, g.text_adaptat, g.system_prompt, g.perfil_id,
               e.a1_coherencia, e.a2_correccio, e.a3_llegibilitat,
               e.b1_fidelitat, e.b2_adequacio_perfil, e.b3_scaffolding, e.b4_cultura,
               e.c1_potencial, e.puntuacio_global
        FROM multi_llm_generations g
        LEFT JOIN multi_v2_evaluations e ON e.generation_id = g.id AND e.judge = 'gpt4mini'
        WHERE g.run_id = 'multi_v2' AND g.prompt_mode = ?
          AND g.text_adaptat IS NOT NULL AND g.text_adaptat != ''
          AND g.generator IN ('gemini', 'gpt', 'mistral')
        ORDER BY g.cas_id, g.generator
        LIMIT 200
    """, (tanda,)).fetchall()
    data = []
    for r in rows:
        data.append({
            "cas_id": r["cas_id"], "generator": r["generator"],
            "perfil": r["perfil_nom"] or "", "perfil_id": r["perfil_id"] or "",
            "mecr": r["mecr"] or "", "dua": r["dua"] or "",
            "paraules": r["text_adaptat_paraules"] or 0,
            "original": r["text_original"] or "", "adaptat": r["text_adaptat"] or "",
            "system_prompt": r["system_prompt"] or "",
            "gpt": {
                "A1": r["a1_coherencia"] or 0, "A2": r["a2_correccio"] or 0,
                "A3": r["a3_llegibilitat"] or 0, "B1": r["b1_fidelitat"] or 0,
                "B2": r["b2_adequacio_perfil"] or 0, "B3": r["b3_scaffolding"] or 0,
                "B4": r["b4_cultura"] or 0, "C1": r["c1_potencial"] or 0,
                "global": r["puntuacio_global"] or 0,
            }
        })
    conn.close()
    return JSONResponse(data)


@app.get("/api/catalog")
async def api_catalog():
    """Retorna el catàleg complet d'instruccions amb metadades."""
    from instruction_catalog import CATALOG, PROFILE_INSTRUCTION_MAP

    category_labels = {
        "A": "Adaptació Lingüística",
        "B": "Estructura i Organització",
        "C": "Suport Cognitiu",
        "D": "Multimodalitat",
        "E": "Contingut Curricular",
        "F": "Avaluació i Comprensió",
        "G": "Personalització Lingüística",
        "H": "Adaptacions per Perfil",
    }

    items = []
    for iid, instr in CATALOG.items():
        item = {
            "id": iid,
            "text": instr["text"],
            "activation": instr["activation"],
            "category": iid.split("-")[0],
        }
        if "profiles" in instr:
            item["profiles"] = instr["profiles"]
        if "mecr_levels" in instr:
            item["mecr_levels"] = instr["mecr_levels"]
        if "mecr_detail" in instr:
            item["mecr_detail"] = instr["mecr_detail"]
        if "suppress_if" in instr:
            item["suppress_if"] = instr["suppress_if"]
        if "subvar_conditions" in instr:
            item["subvar_conditions"] = instr["subvar_conditions"]
        if "complement" in instr:
            item["complement"] = instr["complement"]
        items.append(item)

    return {
        "items": items,
        "total": len(items),
        "category_labels": category_labels,
        "profile_map": PROFILE_INSTRUCTION_MAP,
    }


# ── API Corpus (visor de documents MD) ────────────────────────────────────

CORPUS_DIR = Path(__file__).parent / "corpus"

@app.get("/api/corpus")
async def api_corpus_list():
    """Retorna la llista de fitxers del corpus amb títol i mòdul."""
    files = []
    for f in sorted(CORPUS_DIR.glob("*.md")):
        name = f.stem
        modul = name.split("_")[0]  # M1, M2, M3
        modul_noms = {
            "M1": "Subjecte (perfils alumnat)",
            "M2": "Mètode (metodologies)",
            "M3": "Llengua",
        }
        files.append({
            "filename": f.name,
            "stem": name,
            "modul": modul,
            "modul_nom": modul_noms.get(modul, modul),
            "titol": name.split("_", 1)[1].replace("-", " ").replace("·", "·").title() if "_" in name else name,
            "size_kb": round(f.stat().st_size / 1024, 1),
        })
    return {"files": files, "total": len(files)}


@app.get("/api/corpus/{filename}")
async def api_corpus_file(filename: str):
    """Retorna el contingut d'un fitxer del corpus."""
    filepath = CORPUS_DIR / filename
    if not filepath.exists() or not filepath.suffix == ".md":
        return {"error": "Fitxer no trobat"}, 404
    content = filepath.read_text(encoding="utf-8")
    # Extreure seccions (## headings)
    sections = []
    current = {"title": "Introducció", "content": ""}
    for line in content.split("\n"):
        if line.startswith("## "):
            if current["content"].strip():
                sections.append(current)
            current = {"title": line[3:].strip(), "content": ""}
        else:
            current["content"] += line + "\n"
    if current["content"].strip():
        sections.append(current)
    return {
        "filename": filename,
        "content": content,
        "sections": sections,
        "length": len(content),
        "words": len(content.split()),
    }


# ── API Prompt Preview ───────────────────────────────────────────────────

@app.post("/api/prompt-preview")
async def api_prompt_preview(request: Request):
    """Genera el prompt complet per a un perfil donat (sense cridar l'LLM)."""
    data = await request.json()
    profile = data.get("profile", {})
    context = data.get("context", {"etapa": "ESO", "curs": "3r"})
    params = data.get("params", {"mecr_sortida": "B2", "dua": "Core"})

    mecr = params.get("mecr_sortida", "B2")

    # 1. Instruccions filtrades
    # Bug fix (2026-04-19): instruction_filter.get_instructions() retorna un
    # dict {macrodirectives, audit, suppressed, stats}, no una llista. Els IDs
    # actius els recorrem des de `macrodirectives` (només incloses, NO
    # suprimides — el camp `audit` també conté entrades amb motiu "suprimit").
    filtered = instruction_filter.get_instructions(profile, params)
    instructions_text = instruction_filter.format_instructions_for_prompt(filtered)
    instruction_ids = []
    for _macro_id, _macro in filtered.get("macrodirectives", {}).items():
        for _instr in _macro.get("instruccions", []):
            instruction_ids.append(_instr["id"])
    suppressed_ids = list(filtered.get("suppressed", []))
    stats = filtered.get("stats", {})

    # 2. Persona-audience
    persona = build_persona_audience(profile, context, mecr)

    # 3. Prompt complet (sense RAG — no fem cerca real)
    prompt = build_system_prompt(profile, context, params, "[Aquí anirien 8-12 fragments del corpus FJE cercats per similitud vectorial]")

    # 4. Capes separades per visualització
    identity = corpus_reader.get_identity()
    dua_block = corpus_reader.get_dua_block(params.get("dua", "Core")) or ""
    genre_block = ""
    genre = params.get("genere_discursiu", "")
    if genre:
        genre_block = corpus_reader.get_genre_block(genre) or ""
    fewshot = corpus_reader.get_fewshot_example(mecr) or ""

    return {
        "prompt_complet": prompt,
        "prompt_length": len(prompt),
        "prompt_words": len(prompt.split()),
        "capes": {
            "identitat": identity,
            "instruccions": instructions_text,
            "dua": dua_block,
            "genere": genre_block,
            "persona_audience": persona,
            "fewshot": fewshot,
        },
        "instruction_ids": instruction_ids,
        "total_instructions": stats.get("total_instruccions", len(instruction_ids)),
        "suppressed_ids": suppressed_ids,
        "stats": stats,
    }


# ── API Stats catàleg (font única de veritat) ───────────────────────────────

@app.get("/api/stats-instruccions")
async def api_stats_instruccions():
    """Recompte viu d'instruccions del catàleg per tipus d'activació i macro.

    Font única de veritat per a qualsevol documentació o UI que necessiti
    el nombre d'instruccions. No hardcodejar en prosa — consultar aquí.
    """
    return instruction_catalog.get_catalog_stats()


# ── API Avaluació (dashboard) ──────────────────────────────────────────────

@app.get("/eval", response_class=HTMLResponse)
async def eval_dashboard():
    """Serveix el dashboard d'avaluació (resultats complets)."""
    html_path = UI_DIR / "eval_dashboard.html"
    if html_path.exists():
        return HTMLResponse(html_path.read_text(encoding="utf-8"))
    return HTMLResponse("<h1>Dashboard no disponible</h1>", status_code=404)


@app.get("/eval/progress", response_class=HTMLResponse)
async def eval_progress_page():
    """Serveix la pàgina de monitoratge en temps real."""
    html_path = UI_DIR / "eval_progress.html"
    if html_path.exists():
        return HTMLResponse(html_path.read_text(encoding="utf-8"))
    return HTMLResponse("<h1>Pàgina de progrés no disponible</h1>", status_code=404)


@app.get("/eval/results", response_class=HTMLResponse)
async def eval_results_page():
    """Serveix la pàgina de resultats amb infografies."""
    html_path = UI_DIR / "eval_results.html"
    if html_path.exists():
        return HTMLResponse(html_path.read_text(encoding="utf-8"))
    return HTMLResponse("<h1>Resultats no disponibles</h1>", status_code=404)


@app.get("/eval/cases", response_class=HTMLResponse)
async def eval_cases_page():
    """Serveix la pàgina de visualització de casos (textos adaptats)."""
    html_path = UI_DIR / "eval_cases.html"
    if html_path.exists():
        return HTMLResponse(html_path.read_text(encoding="utf-8"))
    return HTMLResponse("<h1>Cases no disponibles</h1>", status_code=404)


@app.get("/api/eval/cases/{run_id}")
async def eval_cases_detail(run_id: str, limit: int = 20, offset: int = 0,
                            perfil: str = "", etapa: str = "", genere: str = ""):
    """Retorna els casos amb textos adaptats complets per visualitzar."""
    try:
        import eval_db
        conn = eval_db.init_db()

        # Build optional filters
        filters = ""
        params = [run_id]
        if perfil:
            filters += " AND c1.perfil_id = ?"
            params.append(perfil)
        if etapa:
            filters += " AND c1.etapa = ?"
            params.append(etapa)
        if genere:
            filters += " AND c1.genere = ?"
            params.append(genere)

        rows = conn.execute(f"""
            SELECT c1.cas_id, c1.text_id, c1.perfil_id, c1.etapa, c1.genere, c1.mecr, c1.dua,
                   c1.perfils_actius,
                   c1.text_adaptat as text_hc, c1.puntuacio_forma as forma_hc,
                   c1.system_prompt_length as prompt_hc_len,
                   c2.text_adaptat as text_rag, c2.puntuacio_forma as forma_rag,
                   c2.recall as recall_rag,
                   c2.system_prompt_length as prompt_rag_len,
                   c2.total_instruccions_enviades as instr_rag
            FROM eval_cases c1
            JOIN eval_cases c2 ON c1.cas_id = c2.cas_id AND c1.run_id = c2.run_id
            WHERE c1.run_id = ? AND c1.branca = 'hardcoded' AND c2.branca = 'rag'
            {filters}
            ORDER BY c1.cas_id
            LIMIT ? OFFSET ?
        """, (*params, limit, offset)).fetchall()

        total = conn.execute(f"""
            SELECT COUNT(DISTINCT c1.cas_id) FROM eval_cases c1
            JOIN eval_cases c2 ON c1.cas_id = c2.cas_id AND c1.run_id = c2.run_id
            WHERE c1.run_id = ? AND c1.branca = 'hardcoded' AND c2.branca = 'rag'
            {filters}
        """, params).fetchone()[0]

        # Get comparative judgements if available
        judgements = {}
        jrows = conn.execute(
            "SELECT cas_id, judge, winner, justification FROM comparative_judgements WHERE run_id = ?",
            (run_id,)
        ).fetchall()
        for j in jrows:
            cid = j["cas_id"]
            if cid not in judgements:
                judgements[cid] = []
            judgements[cid].append({"judge": j["judge"], "winner": j["winner"],
                                    "justification": j["justification"]})

        cases = []
        for r in rows:
            cases.append({
                "cas_id": r["cas_id"],
                "text_id": r["text_id"],
                "perfil_id": r["perfil_id"],
                "etapa": r["etapa"],
                "genere": r["genere"],
                "mecr": r["mecr"],
                "dua": r["dua"],
                "perfils": r["perfils_actius"],
                "text_hc": r["text_hc"],
                "text_rag": r["text_rag"],
                "forma_hc": r["forma_hc"],
                "forma_rag": r["forma_rag"],
                "recall_rag": r["recall_rag"],
                "prompt_hc_len": r["prompt_hc_len"],
                "prompt_rag_len": r["prompt_rag_len"],
                "instr_rag": r["instr_rag"],
                "judgements": judgements.get(r["cas_id"], []),
            })

        conn.close()
        return JSONResponse({"cases": cases, "total": total, "offset": offset, "limit": limit})
    except Exception as e:
        return JSONResponse({"error": str(e)})


@app.get("/api/eval/originals")
async def eval_originals():
    """Retorna els textos originals del dataset."""
    import json as json_mod
    data_path = Path(__file__).parent / "tests" / "test_data.json"
    with open(data_path, encoding="utf-8") as f:
        data = json_mod.load(f)
    return JSONResponse({"textos": data["textos"], "perfils": data["perfils"]})


@app.get("/api/eval/comparative")
async def eval_comparative():
    """Retorna totes les dades del judici comparatiu per al dashboard."""
    try:
        import eval_db
        conn = eval_db.init_db()

        # Judicis per jutge i guanyador
        judges = conn.execute("""
            SELECT judge, winner, COUNT(*) as n
            FROM comparative_judgements
            GROUP BY judge, winner ORDER BY judge, winner
        """).fetchall()

        # Per perfil
        by_profile = conn.execute("""
            SELECT judge, perfil_id, winner, COUNT(*) as n
            FROM comparative_judgements
            GROUP BY judge, perfil_id, winner
            ORDER BY perfil_id, judge
        """).fetchall()

        # Per criteri
        criteria = {}
        for c in ["c1_winner", "c2_winner", "c3_winner", "c4_winner", "c5_winner"]:
            rows = conn.execute(f"""
                SELECT judge, {c} as winner, COUNT(*) as n
                FROM comparative_judgements WHERE {c} IS NOT NULL AND {c} != ''
                GROUP BY judge, {c}
            """).fetchall()
            criteria[c.replace("_winner", "")] = [dict(r) for r in rows]

        # Justificacions (últims 20)
        justifications = conn.execute("""
            SELECT cas_id, judge, winner, confidence, justification
            FROM comparative_judgements
            ORDER BY id DESC LIMIT 20
        """).fetchall()

        # Mètriques forma de la BD
        forma = conn.execute("""
            SELECT branca,
                   AVG(puntuacio_forma) as avg_forma,
                   AVG(recall) as avg_recall,
                   AVG(text_adaptat_length) as avg_len,
                   COUNT(*) as n
            FROM eval_cases WHERE run_id = (SELECT run_id FROM eval_runs ORDER BY timestamp DESC LIMIT 1)
            GROUP BY branca
        """).fetchall()

        conn.close()

        return JSONResponse({
            "judges": [dict(r) for r in judges],
            "by_profile": [dict(r) for r in by_profile],
            "criteria": criteria,
            "justifications": [dict(r) for r in justifications],
            "forma": [dict(r) for r in forma],
        })
    except Exception as e:
        return JSONResponse({"error": str(e)})


@app.get("/api/eval/runs")
async def eval_runs():
    """Llista totes les execucions d'avaluació."""
    try:
        import eval_db
        conn = eval_db.init_db()
        runs = eval_db.get_all_runs(conn)
        conn.close()
        return JSONResponse(runs)
    except Exception as e:
        return JSONResponse({"error": str(e), "runs": []})


@app.get("/api/eval/run/{run_id}")
async def eval_run_detail(run_id: str):
    """Retorna totes les dades d'una execució (per al dashboard)."""
    try:
        import eval_db
        conn = eval_db.init_db()
        data = eval_db.export_run_json(conn, run_id)
        conn.close()
        return JSONResponse(data)
    except Exception as e:
        return JSONResponse({"error": str(e)})


@app.get("/api/eval/progress")
async def eval_progress():
    """Monitoratge en temps real del batch en curs — llegeix directament la BD."""
    try:
        import eval_db
        conn = eval_db.init_db()

        # Últim run
        last_run = conn.execute(
            "SELECT run_id, timestamp, total_cases, notes FROM eval_runs ORDER BY timestamp DESC LIMIT 1"
        ).fetchone()
        if not last_run:
            return JSONResponse({"status": "idle", "message": "Cap execució registrada"})

        run_id = last_run["run_id"]
        total_expected = last_run["total_cases"] or 0

        # Casos completats (cada cas genera 2 files: HC + RAG)
        count = conn.execute(
            "SELECT COUNT(*) as n FROM eval_cases WHERE run_id = ?", (run_id,)
        ).fetchone()["n"]
        cases_done = count // 2  # 2 branques per cas

        # Comparacions completades (BLOC 3)
        evals_done = conn.execute(
            "SELECT COUNT(*) as n FROM eval_comparisons WHERE run_id = ?", (run_id,)
        ).fetchone()["n"]

        # Mètriques parcials
        stats = conn.execute("""
            SELECT branca,
                   AVG(puntuacio_forma) as avg_forma,
                   AVG(puntuacio_fons) as avg_fons,
                   AVG(recall) as avg_recall,
                   AVG(text_adaptat_length) as avg_len
            FROM eval_cases WHERE run_id = ?
            GROUP BY branca
        """, (run_id,)).fetchall()

        branca_stats = {}
        for row in stats:
            branca_stats[row["branca"]] = {
                "avg_forma": round(row["avg_forma"], 3) if row["avg_forma"] else None,
                "avg_fons": round(row["avg_fons"], 3) if row["avg_fons"] else None,
                "avg_recall": round(row["avg_recall"], 3) if row["avg_recall"] else None,
                "avg_len": int(row["avg_len"]) if row["avg_len"] else None,
            }

        # Últims 5 casos processats
        recent = conn.execute("""
            SELECT cas_id, branca, puntuacio_forma, text_adaptat_length
            FROM eval_cases WHERE run_id = ?
            ORDER BY id DESC LIMIT 10
        """, (run_id,)).fetchall()
        recent_list = [
            {"cas_id": r["cas_id"], "branca": r["branca"],
             "forma": r["puntuacio_forma"], "len": r["text_adaptat_length"]}
            for r in recent
        ]

        # Veredictes (si n'hi ha)
        veredictes = {}
        if evals_done > 0:
            vrows = conn.execute("""
                SELECT veredicte, COUNT(*) as n
                FROM eval_comparisons WHERE run_id = ?
                GROUP BY veredicte
            """, (run_id,)).fetchall()
            veredictes = {r["veredicte"] or "sense_eval": r["n"] for r in vrows}

        # ── Multi-LLM progress ──
        multi_llm = {}
        try:
            mlg = conn.execute("""
                SELECT generator, prompt_mode, COUNT(*) as n, AVG(puntuacio_forma) as avg_forma,
                       AVG(text_adaptat_paraules) as avg_words
                FROM multi_llm_generations WHERE error IS NULL
                GROUP BY generator, prompt_mode ORDER BY generator
            """).fetchall()
            multi_llm["generations"] = [dict(r) for r in mlg]

            mle = conn.execute("""
                SELECT judge, eval_type, COUNT(*) as n
                FROM multi_llm_evaluations
                GROUP BY judge, eval_type ORDER BY judge
            """).fetchall()
            multi_llm["evaluations"] = [dict(r) for r in mle]

            mlr = conn.execute("""
                SELECT cas_id, generator, prompt_mode, text_adaptat_paraules as words
                FROM multi_llm_generations WHERE error IS NULL
                ORDER BY id DESC LIMIT 10
            """).fetchall()
            multi_llm["recent"] = [dict(r) for r in mlr]

            multi_llm["total_generations"] = conn.execute(
                "SELECT COUNT(*) FROM multi_llm_generations WHERE error IS NULL"
            ).fetchone()[0]
            multi_llm["total_evaluations"] = conn.execute(
                "SELECT COUNT(*) FROM multi_llm_evaluations"
            ).fetchone()[0]
        except Exception:
            pass

        # ── Multi-v2 progress (mateixa BD, taules multi_v2_*) ──
        multi_v2 = {}
        try:
            ind = conn.execute(
                "SELECT judge, COUNT(*) as n FROM multi_v2_evaluations "
                "WHERE run_id='multi_v2' GROUP BY judge"
            ).fetchall()
            multi_v2["individual"] = {r["judge"]: r["n"] for r in ind}

            tri = conn.execute(
                "SELECT judge, COUNT(*) as n FROM multi_v2_trios "
                "WHERE run_id='multi_v2' GROUP BY judge"
            ).fetchall()
            multi_v2["trios"] = {r["judge"]: r["n"] for r in tri}

            cro = conn.execute(
                "SELECT judge, pair, COUNT(*) as n FROM multi_v2_cross "
                "WHERE run_id='multi_v2' GROUP BY judge, pair"
            ).fetchall()
            cross_data = {}
            for r in cro:
                cross_data.setdefault(r["judge"], {})[r["pair"]] = r["n"]
            multi_v2["cross"] = cross_data

            recent_v2 = conn.execute("""
                SELECT e.cas_id, e.judge, g.generator, g.prompt_mode, e.puntuacio_global
                FROM multi_v2_evaluations e
                LEFT JOIN multi_llm_generations g ON e.generation_id = g.id
                WHERE e.run_id = 'multi_v2'
                ORDER BY e.id DESC LIMIT 8
            """).fetchall()
            multi_v2["recent"] = [dict(r) for r in recent_v2]
            multi_v2["_n_recent_raw"] = len(recent_v2)
        except Exception:
            pass

        conn.close()

        is_running = cases_done < total_expected
        pct = round(cases_done / total_expected * 100, 1) if total_expected > 0 else 0

        return JSONResponse({
            "status": "running" if is_running else "completed",
            "run_id": run_id,
            "notes": last_run["notes"],
            "timestamp": last_run["timestamp"],
            "progress": {
                "cases_done": cases_done,
                "total_expected": total_expected,
                "evals_done": evals_done,
                "percentage": pct,
            },
            "partial_stats": branca_stats,
            "veredictes": veredictes,
            "recent_cases": recent_list,
            "multi_llm": multi_llm,
            "multi_v2": multi_v2,
        })
    except Exception as e:
        return JSONResponse({"status": "error", "message": str(e)})


@app.get("/api/eval/v2debug")
async def eval_v2debug():
    """Debug endpoint per testar multi_v2 recent query."""
    import eval_db
    conn = eval_db.init_db()
    try:
        n_evals = conn.execute("SELECT COUNT(*) FROM multi_v2_evaluations").fetchone()[0]
        n_gens = conn.execute("SELECT COUNT(*) FROM multi_llm_generations").fetchone()[0]
        rows = conn.execute("""
            SELECT e.cas_id, e.judge, g.generator, g.prompt_mode, e.puntuacio_global
            FROM multi_v2_evaluations e
            LEFT JOIN multi_llm_generations g ON e.generation_id = g.id
            WHERE e.run_id = 'multi_v2'
            ORDER BY e.id DESC LIMIT 3
        """).fetchall()
        conn.close()
        return JSONResponse({
            "n_evals": n_evals, "n_gens": n_gens,
            "recent": [dict(r) for r in rows]
        })
    except Exception as e:
        conn.close()
        return JSONResponse({"error": str(e)})


# ── Docents i perfils personalitzats ────────────────────────────────────────

def _docent_id_from_login(login: str) -> str:
    """SHA256(login lowercase) → 16 hex chars."""
    return hashlib.sha256(login.lower().encode()).hexdigest()[:16]


# Inicialitza el set de docent_ids admin a partir de ATNE_ADMIN_LOGINS
if _ADMIN_LOGINS_RAW:
    for _login in _ADMIN_LOGINS_RAW.split(","):
        _login = _login.strip()
        if _login:
            _ADMIN_DOCENT_IDS.add(_docent_id_from_login(_login))


def _alias_from_login(login: str) -> str:
    """'miquel.amor' → 'Miquel' | 'mamor' → 'Mamor'."""
    local = login.split("@")[0]   # treu domini si ve amb @
    return local.split(".")[0].capitalize()


@app.post("/api/docent/login")
async def docent_login(payload: dict = Body(...)):
    """Identifica el docent per login lanet. Crea o recupera el registre a atne_docents.

    Retorna {ok, docent_id, alias, is_new}.
    """
    # Accepta 'login' (lanet) o 'email' (compatibilitat)
    login = (payload.get("login") or payload.get("email") or "").strip().lower()
    if not login:
        return JSONResponse({"ok": False, "error": "Login obligatori"}, status_code=400)

    email = login  # guardem el login al camp email per coherència amb la taula

    docent_id = _docent_id_from_login(login)
    alias = _alias_from_login(login)

    # Comprova si ja existeix
    resp = requests.get(
        f"{SUPABASE_URL}/rest/v1/atne_docents?id=eq.{docent_id}&select=id,alias",
        headers=SUPABASE_HEADERS, timeout=5,
    )
    if resp.status_code == 200 and resp.json():
        return {"ok": True, "docent_id": docent_id, "alias": resp.json()[0]["alias"], "is_new": False}

    # Crea el registre
    ins = requests.post(
        f"{SUPABASE_URL}/rest/v1/atne_docents",
        headers={**SUPABASE_HEADERS, "Prefer": "return=minimal"},
        json={"id": docent_id, "email": email, "alias": alias},
        timeout=5,
    )
    if ins.status_code not in (200, 201):
        return JSONResponse({"ok": False, "error": "Error creant docent"}, status_code=500)
    return {"ok": True, "docent_id": docent_id, "alias": alias, "is_new": True}


@app.patch("/api/docent/alias")
async def update_docent_alias(payload: dict = Body(...)):
    """Actualitza el nom visible del docent (alias).

    Retorna {ok, alias}.
    """
    docent_id = (payload.get("docent_id") or "").strip()
    alias = (payload.get("alias") or "").strip()
    if not docent_id or not alias:
        return JSONResponse({"ok": False, "error": "docent_id i alias obligatoris"}, status_code=400)
    if len(alias) > 60:
        return JSONResponse({"ok": False, "error": "Nom massa llarg"}, status_code=400)

    resp = requests.patch(
        f"{SUPABASE_URL}/rest/v1/atne_docents?id=eq.{docent_id}",
        headers={**SUPABASE_HEADERS, "Prefer": "return=minimal"},
        json={"alias": alias},
        timeout=5,
    )
    if resp.status_code not in (200, 204):
        return JSONResponse({"ok": False, "error": "Error actualitzant alias"}, status_code=500)
    return {"ok": True, "alias": alias}


@app.get("/api/docent/profiles")
async def get_docent_profiles(docent_id: str = ""):
    """Retorna els perfils personalitzats del docent."""
    docent_id = docent_id.strip()
    if not docent_id:
        return JSONResponse({"ok": False, "error": "docent_id buit"}, status_code=400)
    resp = requests.get(
        f"{SUPABASE_URL}/rest/v1/atne_custom_profiles"
        f"?docent_id=eq.{docent_id}&order=created_at.asc",
        headers=SUPABASE_HEADERS, timeout=5,
    )
    if resp.status_code != 200:
        return JSONResponse({"ok": False, "error": "Error llegint perfils"}, status_code=500)
    return {"ok": True, "profiles": [
        {"_db_id": r["id"], **r["profile_data"]} for r in resp.json()
    ]}


@app.post("/api/docent/profiles")
async def save_docent_profile(payload: dict = Body(...)):
    """Desa un perfil personalitzat per al docent a atne_custom_profiles."""
    docent_id = (payload.get("docent_id") or "").strip()
    profile = payload.get("profile")
    if not docent_id or not profile:
        return JSONResponse({"ok": False, "error": "docent_id i profile són obligatoris"}, status_code=400)
    # Logging del pilot: si el docent ha fet override manual del MECR, ho
    # deixem rastre per analitzar patrons post-pilot (project_parking_lot #63).
    if isinstance(profile, dict) and profile.get("mecr_is_overridden"):
        print(
            f"[ATNE:mecr-override] docent={docent_id} type={profile.get('type','?')} "
            f"curs={profile.get('curs_id','?')} auto={profile.get('mecr_auto','?')} "
            f"→ override={profile.get('mecr_override','?')}",
            flush=True,
        )
    resp = requests.post(
        f"{SUPABASE_URL}/rest/v1/atne_custom_profiles",
        headers={**SUPABASE_HEADERS, "Prefer": "return=representation"},
        json={"docent_id": docent_id, "profile_data": profile},
        timeout=5,
    )
    if resp.status_code not in (200, 201):
        return JSONResponse({"ok": False, "error": "Error desant perfil"}, status_code=500)
    return {"ok": True, "id": resp.json()[0]["id"]}


@app.patch("/api/docent/profiles/{profile_id}")
async def update_docent_profile(profile_id: str, payload: dict = Body(...)):
    """Actualitza un perfil personalitzat (només el propietari)."""
    docent_id = (payload.get("docent_id") or "").strip()
    profile = payload.get("profile")
    if not docent_id or not profile:
        return JSONResponse({"ok": False, "error": "docent_id i profile obligatoris"}, status_code=400)
    resp = requests.patch(
        f"{SUPABASE_URL}/rest/v1/atne_custom_profiles"
        f"?id=eq.{profile_id}&docent_id=eq.{docent_id}",
        headers={**SUPABASE_HEADERS, "Prefer": "return=minimal"},
        json={"profile_data": profile},
        timeout=5,
    )
    return {"ok": resp.status_code in (200, 204)}


@app.delete("/api/docent/profiles/{profile_id}")
async def delete_docent_profile(profile_id: str, docent_id: str = ""):
    """Elimina un perfil personalitzat (només el propietari pot eliminar-lo)."""
    docent_id = docent_id.strip()
    if not docent_id:
        return JSONResponse({"ok": False, "error": "docent_id buit"}, status_code=400)
    resp = requests.delete(
        f"{SUPABASE_URL}/rest/v1/atne_custom_profiles"
        f"?id=eq.{profile_id}&docent_id=eq.{docent_id}",
        headers=SUPABASE_HEADERS, timeout=5,
    )
    return {"ok": resp.status_code in (200, 204)}


@app.get("/api/docent/is-admin")
async def check_is_admin(docent_id: str = ""):
    """Comprova si el docent té rol admin (ATNE_ADMIN_LOGINS env o is_admin a atne_docents)."""
    docent_id = docent_id.strip()
    if not docent_id:
        return {"ok": True, "is_admin": False}
    if docent_id in _ADMIN_DOCENT_IDS:
        return {"ok": True, "is_admin": True}
    if not SUPABASE_URL:
        return {"ok": True, "is_admin": False}
    resp = requests.get(
        f"{SUPABASE_URL}/rest/v1/atne_docents?id=eq.{docent_id}&select=is_admin",
        headers=SUPABASE_HEADERS, timeout=5,
    )
    if resp.status_code == 200 and resp.json():
        return {"ok": True, "is_admin": bool(resp.json()[0].get("is_admin", False))}
    return {"ok": True, "is_admin": False}


@app.get("/api/docent/list-admins")
async def list_admins(_: bool = Depends(_require_admin)):
    """Retorna la llista d'admins: env var (permanents) + Supabase (is_admin=true)."""
    result: list[dict] = []
    env_ids: set[str] = set()

    if _ADMIN_LOGINS_RAW:
        for _lg in _ADMIN_LOGINS_RAW.split(","):
            _lg = _lg.strip()
            if _lg:
                did = _docent_id_from_login(_lg)
                env_ids.add(did)
                result.append({"docent_id": did, "alias": _alias_from_login(_lg),
                                "email": _lg, "permanent": True})

    if SUPABASE_URL:
        resp = requests.get(
            f"{SUPABASE_URL}/rest/v1/atne_docents?is_admin=eq.true&select=id,alias,email",
            headers=SUPABASE_HEADERS, timeout=5,
        )
        if resp.status_code == 200:
            for row in resp.json():
                if row["id"] not in env_ids:
                    result.append({"docent_id": row["id"],
                                   "alias": row.get("alias") or row["id"],
                                   "email": row.get("email") or "",
                                   "permanent": False})

    return {"ok": True, "admins": result}


@app.post("/api/docent/set-admin")
async def set_admin(payload: dict = Body(...), _: bool = Depends(_require_admin)):
    """Concedeix o revoca rol admin a un docent (requereix sessió admin).
    Accepta {login} (lanet) o {docent_id}. Fa upsert per si el docent no existeix encara."""
    login = (payload.get("login") or "").strip().lower()
    docent_id = (payload.get("docent_id") or "").strip()
    is_admin = bool(payload.get("is_admin", True))

    if not docent_id and login:
        docent_id = _docent_id_from_login(login)
    if not docent_id:
        return JSONResponse({"ok": False, "error": "login o docent_id obligatori"}, status_code=400)

    upsert: dict = {"id": docent_id, "is_admin": is_admin}
    if login:
        upsert["email"] = login
        upsert["alias"] = _alias_from_login(login)

    resp = requests.post(
        f"{SUPABASE_URL}/rest/v1/atne_docents",
        headers={**SUPABASE_HEADERS, "Prefer": "resolution=merge-duplicates,return=minimal"},
        json=upsert,
        timeout=5,
    )
    return {"ok": resp.status_code in (200, 201, 204),
            "docent_id": docent_id,
            "alias": upsert.get("alias", docent_id)}


# ── Main ────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8000))
    print("=" * 50)
    print("  ATNE — Adaptador de Textos")
    print(f"  http://localhost:{port}")
    try:
        _stats = instruction_catalog.get_catalog_stats()
        _parts = ", ".join(f"{k}={v}" for k, v in sorted(_stats["per_activation"].items()))
        print(f"  Catàleg: {_stats['total']} instruccions ({_parts})")
    except Exception:
        pass
    print("=" * 50)
    uvicorn.run(app, host="0.0.0.0", port=port)
