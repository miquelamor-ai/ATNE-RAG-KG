"""
Adaptació de documents PDF i PPTX amb preservació de format.

PDF:  PyMuPDF span-level redact+reinsert
PPTX: python-pptx run-consolidation

Pipeline:
  1. extract_*_text_map(bytes) → [{id, text, meta}]
  2. batch_adapt_text_map(text_map, model, system_prompt) → {id: adapted_text}
  3. inject_*_adapted(bytes, dict{id→adapted_text}) → bytes
"""

import io
import json
import re

# Molts PDFs codifiquen l'apòstrof amb un glifo que PyMuPDF mapeja a U+00B7 (·).
# Normalitzem: · entre lletres → ' (excepte l·l que és la L geminada catalana).
_MID_DOT_RE = re.compile(r'([A-Za-zÀ-ÖØ-öø-ÿ])(·)([A-Za-zÀ-ÖØ-öø-ÿ])')

def _fix_apostrophes(text: str) -> str:
    def _repl(m):
        l, r = m.group(1), m.group(3)
        if l.lower() == 'l' and r.lower() == 'l':
            return m.group(0)
        return l + "'" + r
    return _MID_DOT_RE.sub(_repl, text)


# ── PDF ─────────────────────────────────────────────────────────────────────

def is_scanned_pdf(pdf_bytes: bytes) -> bool:
    """Retorna True si el PDF no té text seleccionable."""
    import fitz
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    total_chars = sum(len(page.get_text()) for page in doc)
    doc.close()
    return total_chars < 50


def extract_pdf_text_map(pdf_bytes: bytes) -> list:
    """Extreu blocs de text nadiu d'un PDF. Retorna [{id, text, page, bbox, fontsize, color}]."""
    import fitz
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    items = []
    for page_num, page in enumerate(doc):
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
        for block_num, block in enumerate(blocks):
            if block["type"] != 0:
                continue
            lines_text = []
            for line in block["lines"]:
                line_text = "".join(span["text"] for span in line["spans"])
                if line_text.strip():
                    lines_text.append(line_text.strip())
            block_text = _fix_apostrophes(" ".join(lines_text))
            if len(block_text.strip()) < 3:
                continue
            first_span = None
            for line in block["lines"]:
                if line["spans"]:
                    first_span = line["spans"][0]
                    break
            origin = list(first_span["origin"]) if first_span and "origin" in first_span else [block["bbox"][0], block["bbox"][3]]
            items.append({
                "id": f"p{page_num}_b{block_num}",
                "text": block_text,
                "page": page_num,
                "bbox": list(block["bbox"]),
                "origin": origin,
                "fontsize": round(first_span["size"], 1) if first_span else 11.0,
                "font": (first_span["font"] if first_span else "Helvetica"),
                "color": (first_span["color"] if first_span else 0),
            })
    doc.close()
    return items


def inject_pdf_adapted(pdf_bytes: bytes, adapted: dict) -> bytes:
    """Reinjecta text adaptat al PDF original. adapted: {id → text_adaptat}."""
    import fitz
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    for page_num, page in enumerate(doc):
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
        to_inject = []
        for block_num, block in enumerate(blocks):
            if block["type"] != 0:
                continue
            bid = f"p{page_num}_b{block_num}"
            if bid not in adapted:
                continue
            new_text = adapted[bid]
            bbox = fitz.Rect(block["bbox"])
            first_span = None
            for line in block["lines"]:
                if line["spans"]:
                    first_span = line["spans"][0]
                    break
            fontsize = round(first_span["size"], 1) if first_span else 11.0
            origin = first_span["origin"] if first_span and "origin" in first_span else (block["bbox"][0], block["bbox"][3])
            color_int = first_span["color"] if first_span else 0
            r = ((color_int >> 16) & 0xFF) / 255
            g = ((color_int >> 8) & 0xFF) / 255
            b = (color_int & 0xFF) / 255
            to_inject.append((bbox, origin, new_text, fontsize, (r, g, b)))
        # Primer redactar tot el bloc, després reinsertar
        for bbox, _, _, _, _ in to_inject:
            page.add_redact_annot(bbox, fill=(1, 1, 1))
        if to_inject:
            page.apply_redactions()
        for bbox, origin, new_text, fontsize, color in to_inject:
            # insert_text usa l'origen (baseline), més fiable que insert_textbox per text curt
            page.insert_text(origin, new_text, fontsize=fontsize, color=color)
    out = io.BytesIO()
    doc.save(out)
    doc.close()
    return out.getvalue()


# ── PPTX ────────────────────────────────────────────────────────────────────

def extract_pptx_text_map(pptx_bytes: bytes) -> list:
    """Extreu paràgrafs de text d'un PPTX. Retorna [{id, text, slide, shape, para}]."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    items = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                para_text = para.text.strip()
                if len(para_text) < 3:
                    continue
                items.append({
                    "id": f"s{slide_idx}_sh{shape_idx}_p{para_idx}",
                    "text": para_text,
                    "slide": slide_idx,
                    "shape": shape_idx,
                    "para": para_idx,
                    "shape_name": shape.name,
                })
    return items


def inject_pptx_adapted(pptx_bytes: bytes, adapted: dict) -> bytes:
    """
    Reinjecta text adaptat al PPTX original preservant format visual.
    Consolida els runs de cada paràgraf en un (hereta format del primer run).
    adapted: {id → text_adaptat}
    """
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                bid = f"s{slide_idx}_sh{shape_idx}_p{para_idx}"
                if bid not in adapted or not para.runs:
                    continue
                new_text = adapted[bid]
                first_run = para.runs[0]
                for run in para.runs[1:]:
                    para._p.remove(run._r)
                first_run.text = new_text
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ── Batch LLM ────────────────────────────────────────────────────────────────

def batch_adapt_text_map(text_map: list, model_id: str, system_prompt: str) -> dict:
    """
    Envia tots els segments al LLM en lots de 30 i retorna {id: text_adaptat}.
    En cas d'error de parsing del JSON, conserva el text original.
    """
    from adaptation.llm_clients import _call_llm
    result = {}
    BATCH = 30
    for i in range(0, len(text_map), BATCH):
        chunk = text_map[i:i + BATCH]
        segments_json = json.dumps(
            [{"id": x["id"], "text": x["text"]} for x in chunk],
            ensure_ascii=False,
            indent=2,
        )
        user_msg = (
            "Adapta cadascun dels segments de text del document "
            "aplicant les instruccions del sistema.\n\n"
            "IMPORTANT:\n"
            "- Retorna ÚNICAMENT un objecte JSON vàlid: {\"id\": \"text adaptat\", ...}\n"
            "- Mantén els mateixos id exactes, no n'omitis cap\n"
            "- Preserva majúscules inicials i puntuació final de cada segment\n"
            "- No afegeixis complements ni seccions extra\n\n"
            f"SEGMENTS:\n{segments_json}"
        )
        raw = _call_llm(model_id, system_prompt, user_msg)
        m = re.search(r'\{.*\}', raw, re.DOTALL)
        if m:
            try:
                result.update(json.loads(m.group(0)))
                continue
            except json.JSONDecodeError:
                pass
        # Fallback: text original
        for item in chunk:
            result[item["id"]] = item["text"]
    return result
